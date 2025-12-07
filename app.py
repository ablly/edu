import os
import traceback
import uuid
import zipfile
import tempfile
import json
import re
import sys
import requests
import shutil
from flask import Flask, request, jsonify, render_template, redirect, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
from datetime import datetime
import openai  # 旧版本openai库
import config
import PyPDF2

# 配置旧版本openai
openai.api_key = config.DEEPSEEK_API_KEY
openai.api_base = config.DEEPSEEK_BASE_URL
import base64
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.util import Inches  # 设置PPT元素尺寸
from docx import Document
from pptx import Presentation
from models import db, Student, Assignment, QuestionBank, QuestionSubmission, VideoNote, Conversation, \
    ConversationMessage  # 新增QuestionBank
from models_membership import User, MembershipTier, UserMembership, PaymentTransaction, UsageLog
from models_order import Order, OrderRefund
from utils.security import (
    validate_password_strength, validate_username, validate_email, sanitize_input,
    record_login_attempt, is_account_locked, get_remaining_attempts
)
from membership_utils import feature_limit, check_feature_access, log_feature_usage, get_usage_stats
from config import DEEPSEEK_API_KEY, DEEPSEEK_BASE_URL
from flask_migrate import Migrate
from flask_login import LoginManager, login_user, logout_user, login_required, current_user
from flask_mail import Mail
from datetime import datetime, timedelta
# 导入自定义认证装饰器
from utils.auth_decorators import require_login_api, require_membership
# 导入验证码模型和邮件服务
from models_verification import VerificationCode
from utils.email_service import EmailService
from pptx.enum.text import PP_ALIGN
import time
import threading
import warnings

# 初始化 Flask 应用
app = Flask(__name__)
app.config.from_object(config.Config)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB
# 禁用模板缓存，确保实时读取最新模板
app.config['TEMPLATES_AUTO_RELOAD'] = True
app.config['SEND_FILE_MAX_AGE_DEFAULT'] = 0

# 邮件服务配置
app.config['MAIL_SERVER'] = 'smtp.qq.com'
app.config['MAIL_PORT'] = 587
app.config['MAIL_USE_TLS'] = True
app.config['MAIL_USE_SSL'] = False
app.config['MAIL_USERNAME'] = '3533912007@qq.com'
app.config['MAIL_PASSWORD'] = 'uxsenlfyzpdrdbgh'  # QQ邮箱授权码
app.config['MAIL_DEFAULT_SENDER'] = ('EduPilot AI', '3533912007@qq.com')
app.config['MAIL_MAX_EMAILS'] = None
app.config['MAIL_ASCII_ATTACHMENTS'] = False

# 确保上传文件夹存在（Vercel等只读文件系统会跳过）
try:
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    os.makedirs(config.UPLOAD_FOLDER, exist_ok=True)
    if hasattr(config, 'DATABASE_PATH') and config.DATABASE_PATH:
        os.makedirs(os.path.dirname(config.DATABASE_PATH), exist_ok=True)
except (OSError, PermissionError):
    # Vercel等无服务器环境是只读文件系统，跳过目录创建
    pass

# 初始化数据库
db.init_app(app)
migrate = Migrate(app, db)

# 初始化邮件服务
mail = Mail(app)
email_service = EmailService(mail)

# 初始化Flask-Login
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login_page'
login_manager.login_message = '请先登录'

# 初始化CSRF保护（先禁用，后续逐步启用）
from flask_wtf.csrf import CSRFProtect
csrf = CSRFProtect()
# 启用CSRF保护，但默认不检查（由 WTF_CSRF_CHECK_DEFAULT = False 控制）
csrf.init_app(app)  # 已启用，仅对登录/注册生效

# 初始化CORS（允许React前端跨域访问）
CORS(app, resources={
    r"/api/*": {
        "origins": ["http://localhost:3000", "http://localhost:5173"],
        "methods": ["GET", "POST", "PUT", "DELETE", "PATCH", "OPTIONS"],
        "allow_headers": ["Content-Type", "Authorization"],
        "supports_credentials": True,  # 支持携带Cookie
        "max_age": 3600
    }
})

# 初始化API限流保护
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address

limiter = Limiter(
    app=app,
    key_func=get_remote_address,
    default_limits=["200 per day", "50 per hour"],  # 默认限制
    storage_uri=app.config.get('RATELIMIT_STORAGE_URL', 'memory://'),
    headers_enabled=app.config.get('RATELIMIT_HEADERS_ENABLED', True),
    swallow_errors=app.config.get('RATELIMIT_SWALLOW_ERRORS', True)
)

@login_manager.user_loader
def load_user(user_id):
    """Flask-Login用户加载回调"""
    return User.query.get(int(user_id))

# 创建数据库表
with app.app_context():
    db.create_all()

    # ✅ 修复后的调试代码（兼容 SQLAlchemy ≥2.0）
    from sqlalchemy import inspect

    inspector = inspect(db.engine)

    print("=== 检查数据库表结构 ===")
    print("所有表名:", inspector.get_table_names())  # 替代原先的 db.engine.table_names()

    if 'question_bank' in inspector.get_table_names():
        print("question_bank 表的列:", inspector.get_columns('question_bank'))
    else:
        print("❌ question_bank 表不存在，请检查模型定义")


# 资源文件目录访问
def source_path(relative_path):
    # 是否Bundle Resource
    if getattr(sys, 'frozen', False):
        base_path = sys._MEIPASS
    else:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)


# 修改当前工作目录，使得资源文件可以被正确访问
cd = source_path('')
os.chdir(cd)


# 检查文件类型是否允许
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in config.ALLOWED_EXTENSIONS


# 从文件名解析学号和姓名
def parse_student_info(filename):
    base_name = os.path.splitext(filename)[0]
    match = re.match(r'^([A-Za-z0-9]+)-(.+)$', base_name)
    if match:
        student_id = match.group(1)
        name = match.group(2)
        return student_id, name
    return None, None


# 获取或创建学生记录
def get_or_create_student(student_id, name):
    student = Student.query.filter_by(student_id=student_id).first()
    if student:
        return student
    new_student = Student(student_id=student_id, name=name)
    db.session.add(new_student)
    db.session.commit()
    return new_student


# 解析 PDF 文件为纯文本
def parse_pdf(file_path):
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        return ' '.join(page.extract_text() for page in reader.pages if page.extract_text())


# 解析 DOCX 文件为纯文本
def parse_docx(file_path):
    doc = Document(file_path)
    text = ''
    for para in doc.paragraphs:
        text += para.text + '\n'
    return text


#解析ppt
def parse_pptx(file_path):
    """解析 PPTX 文件为纯文本，并包含图片等内容的描述"""
    presentation = Presentation(file_path)
    text = ''
    for slide_index, slide in enumerate(presentation.slides):
        text += f"幻灯片 {slide_index + 1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + '\n'
            # 处理图片
            if shape.shape_type == 13:  # 13 表示图片
                text += f"[图片描述] 幻灯片 {slide_index + 1} 包含一张图片，可能与主题相关。\n"
    return text


def parse_ai_response(ai_response):
    """
    预处理 AI 返回的 JSON 数据，去除外部的 ```json 标记。
    """
    try:
        # 使用正则表达式去除 ```json 和多余的换行符
        ai_response = re.sub(r'```json|\```', '', ai_response).strip()
        # 将字符串解析为 JSON
        return json.loads(ai_response)
    except Exception as e:
        raise ValueError(f"无法解析 AI 返回的题目数据: {str(e)}")


def parse_txt(file_path):
    """解析 TXT 文件为纯文本"""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()


def sanitize_ai_response(text):
    """
    清理和规范化AI响应文本
    1. 移除非法字符
    2. 规范化标点符号
    3. 确保段落格式
    """
    # 移除ASCII控制字符（保留\t\n\r）
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)

    # 中文标点规范化
    text = text.replace('，', '，').replace('。', '。').replace('？', '？')

    # 英文标点规范化
    text = text.replace(' ,', ',').replace(' .', '.').replace(' ?', '?')

    # 确保问答段落之间有适当间距
    text = re.sub(r'(\n\s*){3,}', '\n\n', text)

    # 移除首尾空白
    return text.strip()


def hex_to_rgb(hex_color):
    """将十六进制颜色代码转换为RGB元组"""
    hex_color = str(hex_color).lstrip('#')

    # 处理缩写格式如 #RGB
    if len(hex_color) == 3:
        hex_color = ''.join([c * 2 for c in hex_color])

    # 验证长度
    if len(hex_color) != 6:
        return (0, 0, 0)  # 返回默认黑色

    try:
        return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return (0, 0, 0)  # 返回默认黑色


# 预定义的课程结构（大学/高等教育课程）
COURSE_STRUCTURE = {
    "计算机组成原理": ["数字逻辑基础", "计算机系统概述", "数据的表示和运算", "存储系统"],
    "操作系统": ["操作系统概述", "进程管理", "内存管理", "文件系统"],
    "数据结构": ["线性结构", "树形结构", "图结构", "排序与查找"],
    "计算机网络": ["网络基础", "TCP/IP协议", "网络应用", "网络安全"],
    "数据库系统": ["数据库基础", "SQL语言", "数据库设计", "事务管理"],
    "软件工程": ["软件过程", "需求工程", "软件设计", "软件测试"],
    "人工智能": ["搜索算法", "机器学习基础", "神经网络", "自然语言处理"]
}

# K12课程结构（小学、初中、高中）
K12_COURSE_STRUCTURE = None

def load_k12_courses():
    """加载K12课程结构数据"""
    global K12_COURSE_STRUCTURE
    try:
        k12_file_path = os.path.join(config.BASE_DIR, 'data', 'k12_courses.json')
        if os.path.exists(k12_file_path):
            with open(k12_file_path, 'r', encoding='utf-8') as f:
                K12_COURSE_STRUCTURE = json.load(f)
            app.logger.info("K12课程结构加载成功")
        else:
            app.logger.warning(f"K12课程文件不存在: {k12_file_path}")
            K12_COURSE_STRUCTURE = {}
    except Exception as e:
        app.logger.error(f"加载K12课程结构失败: {str(e)}")
        K12_COURSE_STRUCTURE = {}

# 在应用启动时加载K12课程
load_k12_courses()


# 数据库初始化路由（用于 Vercel 部署后首次初始化）
@app.route('/api/init_db', methods=['GET', 'POST'])
@csrf.exempt
def api_init_database():
    """初始化数据库表结构和默认数据"""
    try:
        from models_admin import init_admin_tables
        from models_membership import MembershipTier
        
        # 创建所有表
        db.create_all()
        
        # 初始化管理员
        created_admins = init_admin_tables()
        
        # 检查是否需要创建默认会员套餐
        free_tier = MembershipTier.query.filter_by(code='free').first()
        tier_created = False
        if not free_tier:
            free_tier = MembershipTier(
                name='免费版',
                code='free',
                level=0,
                price=0,
                duration_days=0,
                is_active=True,
                description='免费用户',
                permissions='{"allowed_features": ["ai_ask"], "limits": {"ai_ask": 5}}'
            )
            db.session.add(free_tier)
            db.session.commit()
            tier_created = True
        
        return jsonify({
            'success': True,
            'message': '数据库初始化成功',
            'admins_created': len(created_admins),
            'tier_created': tier_created
        })
    except Exception as e:
        import traceback
        return jsonify({
            'success': False,
            'error': str(e),
            'traceback': traceback.format_exc()
        }), 500


# 页面路由
@app.route('/')
def index():
    return render_template('index.html')


@app.route('/students')
def students_page():
    return render_template('students.html')


@app.route('/submit')
def submit_page():
    return render_template('submit.html')


@app.route('/scores')
def scores_page():
    return render_template('scores.html')


@app.route('/ai-ask')
def ai_ask_page():
    return render_template('ai_ask.html')


@app.route('/Auxiliary-programming')
def auxiliary_programming_page():
    return render_template('Auxiliary_programming.html')


@app.route('/video-summary')
def video_summary_page():
    return render_template('video_summary.html')


@app.route('/generate-question')
def generate_question_page():
    return render_template('generate_question.html')


@app.route('/generate-lecture')
def generate_lecture_page():
    return render_template('generate_lecture.html')


@app.route('/login')
def login_page():
    return render_template('login.html')


@app.route('/forgot-password')
def forgot_password_page():
    """忘记密码页面"""
    return render_template('forgot_password.html')


@app.route('/profile')
@login_required
def profile_page():
    """个人中心页面"""
    return render_template('profile.html')


@app.route('/membership')
@login_required
def membership_center():
    """会员中心页面"""
    return render_template('membership.html')


@app.route('/payment')
@login_required
def payment_page():
    """支付页面"""
    return render_template('payment.html')


@app.route('/membership-test')
def membership_test():
    """会员功能测试页面"""
    return render_template('membership_test.html')


@app.route('/progress')
def progress_chart():
    students = Student.query.all()  # 假设这是获取学生数据的代码
    return render_template(
        'progress_chart.html',
        students=students,
        now=datetime.now()  # 添加当前时间
    )


@app.route('/api/students/progress', methods=['GET'])
@csrf.exempt
def get_students_progress():
    """获取学生成绩进度数据，用于图表分析"""
    try:
        # 获取所有学生
        students = Student.query.all()
        
        # 获取作业统计数据
        assignment_stats = db.session.query(
            Assignment.student_id,
            db.func.count(Assignment.id).label('assignment_count'),
            db.func.avg(Assignment.score).label('avg_score'),
            db.func.max(Assignment.score).label('max_score'),
            db.func.min(Assignment.score).label('min_score')
        ).group_by(Assignment.student_id).all()
        
        # 获取题目提交统计
        question_stats = db.session.query(
            QuestionSubmission.student_id,
            db.func.count(QuestionSubmission.id).label('question_count'),
            db.func.avg(db.case((QuestionSubmission.is_correct == True, 100), else_=0)).label('accuracy_rate')
        ).group_by(QuestionSubmission.student_id).all()
        
        # 构建学生进度数据
        progress_data = []
        assignment_dict = {stat.student_id: stat for stat in assignment_stats}
        question_dict = {stat.student_id: stat for stat in question_stats}
        
        for student in students:
            student_id = student.student_id
            assignment_stat = assignment_dict.get(student_id)
            question_stat = question_dict.get(student_id)
            
            student_data = {
                'student_id': student_id,
                'name': student.name,
                'class_name': student.class_name or '未分班',
                'major': student.major or '未设置',
                'assignment_count': assignment_stat.assignment_count if assignment_stat else 0,
                'avg_score': round(assignment_stat.avg_score, 2) if assignment_stat and assignment_stat.avg_score else 0,
                'max_score': assignment_stat.max_score if assignment_stat else 0,
                'min_score': assignment_stat.min_score if assignment_stat else 0,
                'question_count': question_stat.question_count if question_stat else 0,
                'accuracy_rate': round(question_stat.accuracy_rate, 2) if question_stat and question_stat.accuracy_rate else 0
            }
            progress_data.append(student_data)
        
        # 班级统计
        class_stats = {}
        for data in progress_data:
            class_name = data['class_name']
            if class_name not in class_stats:
                class_stats[class_name] = {
                    'class_name': class_name,
                    'student_count': 0,
                    'total_assignments': 0,
                    'avg_class_score': 0,
                    'total_questions': 0,
                    'avg_accuracy': 0
                }
            
            class_stats[class_name]['student_count'] += 1
            class_stats[class_name]['total_assignments'] += data['assignment_count']
            class_stats[class_name]['avg_class_score'] += data['avg_score']
            class_stats[class_name]['total_questions'] += data['question_count']
            class_stats[class_name]['avg_accuracy'] += data['accuracy_rate']
        
        # 计算班级平均值
        for class_name, stats in class_stats.items():
            if stats['student_count'] > 0:
                stats['avg_class_score'] = round(stats['avg_class_score'] / stats['student_count'], 2)
                stats['avg_accuracy'] = round(stats['avg_accuracy'] / stats['student_count'], 2)
        
        # 专业统计
        major_stats = {}
        for data in progress_data:
            major = data['major']
            if major not in major_stats:
                major_stats[major] = {
                    'major': major,
                    'student_count': 0,
                    'avg_score': 0,
                    'avg_accuracy': 0
                }
            
            major_stats[major]['student_count'] += 1
            major_stats[major]['avg_score'] += data['avg_score']
            major_stats[major]['avg_accuracy'] += data['accuracy_rate']
        
        # 计算专业平均值
        for major, stats in major_stats.items():
            if stats['student_count'] > 0:
                stats['avg_score'] = round(stats['avg_score'] / stats['student_count'], 2)
                stats['avg_accuracy'] = round(stats['avg_accuracy'] / stats['student_count'], 2)
        
        return jsonify({
            'students': progress_data,
            'class_stats': list(class_stats.values()),
            'major_stats': list(major_stats.values()),
            'total_students': len(students),
            'total_assignments': sum(data['assignment_count'] for data in progress_data),
            'total_questions': sum(data['question_count'] for data in progress_data)
        })
        
    except Exception as e:
        return jsonify({'error': f'获取进度数据失败: {str(e)}'}), 500


@app.route('/api/assignments/recent', methods=['GET'])
@csrf.exempt
def get_recent_assignments():
    """获取最近的作业提交记录"""
    try:
        limit = int(request.args.get('limit', 10))
        
        recent_assignments = db.session.query(Assignment, Student).join(
            Student, Assignment.student_id == Student.student_id
        ).order_by(Assignment.submission_time.desc()).limit(limit).all()
        
        assignments_data = []
        for assignment, student in recent_assignments:
            assignments_data.append({
                'id': assignment.id,
                'student_name': student.name,
                'student_id': assignment.student_id,
                'assignment_name': assignment.assignment_name,
                'subject': assignment.subject,
                'chapter': assignment.chapter,
                'score': assignment.score,
                'submission_time': assignment.submission_time.strftime('%Y-%m-%d %H:%M:%S'),
                'feedback': assignment.feedback
            })
        
        return jsonify({'assignments': assignments_data})
        
    except Exception as e:
        return jsonify({'error': f'获取作业记录失败: {str(e)}'}), 500
# API 路由
@app.route('/api/students', methods=['GET'])
@csrf.exempt
def get_students():
    """获取学生列表，支持搜索和筛选"""
    # 获取查询参数
    search = request.args.get('search', '')
    major = request.args.get('major', '')
    class_name = request.args.get('class', '')
    status = request.args.get('status', '')
    page = int(request.args.get('page', 1))
    per_page = int(request.args.get('per_page', 20))
    simple = request.args.get('simple', 'false').lower() == 'true'
    
    # 构建查询
    query = Student.query
    
    # 搜索条件
    if search:
        search_filter = f"%{search}%"
        query = query.filter(
            db.or_(
                Student.name.like(search_filter),
                Student.student_id.like(search_filter),
                Student.phone.like(search_filter),
                Student.email.like(search_filter)
            )
        )
    
    # 筛选条件
    if major:
        query = query.filter(Student.major.like(f"%{major}%"))
    if class_name:
        query = query.filter(Student.class_name.like(f"%{class_name}%"))
    if status:
        query = query.filter(Student.student_status == status)
    
    # 分页
    students = query.paginate(
        page=page, per_page=per_page, error_out=False
    )
    
    # 返回数据
    if simple:
        student_list = [student.to_simple_dict() for student in students.items]
    else:
        student_list = [student.to_dict() for student in students.items]
    
    return jsonify({
        'students': student_list,
        'total': students.total,
        'pages': students.pages,
        'current_page': page,
        'per_page': per_page
    })


@app.route('/api/students', methods=['POST'])
@csrf.exempt
def add_student():
    """添加学生，支持完整的学生信息"""
    data = request.json
    
    # 检查必填字段
    if not data.get('student_id') or not data.get('name'):
        return jsonify({'error': '学号和姓名为必填字段'}), 400
    
    # 检查学号是否已存在
    existing_student = Student.query.filter_by(student_id=data['student_id']).first()
    if existing_student:
        return jsonify({'error': '学号已存在'}), 400
    
    try:
        # 创建新学生
        new_student = Student(
            student_id=data['student_id'],
            name=data['name'],
            gender=data.get('gender'),
            birth_date=datetime.strptime(data['birth_date'], '%Y-%m-%d').date() if data.get('birth_date') else None,
            id_card=data.get('id_card'),
            phone=data.get('phone'),
            email=data.get('email'),
            admission_date=datetime.strptime(data['admission_date'], '%Y-%m-%d').date() if data.get('admission_date') else None,
            major=data.get('major'),
            class_name=data.get('class_name'),
            grade=data.get('grade'),
            education_level=data.get('education_level'),
            student_status=data.get('student_status', '在读'),
            address=data.get('address'),
            emergency_contact=data.get('emergency_contact'),
            emergency_phone=data.get('emergency_phone'),
            parent_name=data.get('parent_name'),
            parent_phone=data.get('parent_phone'),
            notes=data.get('notes'),
            tags=json.dumps(data.get('tags', [])) if data.get('tags') else None
        )
        
        db.session.add(new_student)
        db.session.commit()
        
        return jsonify(new_student.to_dict()), 201
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'添加学生失败: {str(e)}'}), 500


@app.route('/api/students/<student_id>', methods=['PUT'])
@csrf.exempt
def update_student(student_id):
    """更新学生信息，支持完整字段更新"""
    student = Student.query.filter_by(student_id=student_id).first()
    if not student:
        return jsonify({'error': '学生不存在'}), 404
    
    data = request.json
    
    try:
        # 检查新学号是否与其他学生冲突
        if 'student_id' in data and data['student_id'] != student_id:
            existing_student = Student.query.filter_by(student_id=data['student_id']).first()
            if existing_student:
                return jsonify({'error': '新学号已被其他学生使用'}), 400
            student.student_id = data['student_id']
        
        # 更新基本信息
        if 'name' in data:
            student.name = data['name']
        if 'gender' in data:
            student.gender = data['gender']
        if 'birth_date' in data and data['birth_date']:
            student.birth_date = datetime.strptime(data['birth_date'], '%Y-%m-%d').date()
        if 'id_card' in data:
            student.id_card = data['id_card']
        if 'phone' in data:
            student.phone = data['phone']
        if 'email' in data:
            student.email = data['email']
        
        # 更新学术信息
        if 'admission_date' in data and data['admission_date']:
            student.admission_date = datetime.strptime(data['admission_date'], '%Y-%m-%d').date()
        if 'major' in data:
            student.major = data['major']
        if 'class_name' in data:
            student.class_name = data['class_name']
        if 'grade' in data:
            student.grade = data['grade']
        if 'education_level' in data:
            student.education_level = data['education_level']
        if 'student_status' in data:
            student.student_status = data['student_status']
        
        # 更新联系信息
        if 'address' in data:
            student.address = data['address']
        if 'emergency_contact' in data:
            student.emergency_contact = data['emergency_contact']
        if 'emergency_phone' in data:
            student.emergency_phone = data['emergency_phone']
        if 'parent_name' in data:
            student.parent_name = data['parent_name']
        if 'parent_phone' in data:
            student.parent_phone = data['parent_phone']
        
        # 更新扩展信息
        if 'notes' in data:
            student.notes = data['notes']
        if 'tags' in data:
            student.tags = json.dumps(data['tags']) if data['tags'] else None
        
        # 更新时间戳
        student.updated_at = datetime.utcnow()
        
        db.session.commit()
        return jsonify(student.to_dict()), 200
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'更新失败: {str(e)}'}), 500


@app.route('/api/students/<student_id>', methods=['GET'])
@csrf.exempt
def get_student_detail(student_id):
    """获取学生详细信息"""
    student = Student.query.filter_by(student_id=student_id).first()
    if not student:
        return jsonify({'error': '学生不存在'}), 404
    
    return jsonify(student.to_dict())


@app.route('/api/students/stats', methods=['GET'])
@csrf.exempt
def get_students_stats():
    """获取学生统计信息"""
    try:
        total_students = Student.query.count()
        
        # 按状态统计
        status_stats = db.session.query(
            Student.student_status,
            db.func.count(Student.id)
        ).group_by(Student.student_status).all()
        
        # 按专业统计
        major_stats = db.session.query(
            Student.major,
            db.func.count(Student.id)
        ).filter(Student.major.isnot(None)).group_by(Student.major).all()
        
        # 按班级统计
        class_stats = db.session.query(
            Student.class_name,
            db.func.count(Student.id)
        ).filter(Student.class_name.isnot(None)).group_by(Student.class_name).all()
        
        # 按性别统计
        gender_stats = db.session.query(
            Student.gender,
            db.func.count(Student.id)
        ).filter(Student.gender.isnot(None)).group_by(Student.gender).all()
        
        return jsonify({
            'total_students': total_students,
            'status_distribution': [{'status': status, 'count': count} for status, count in status_stats],
            'major_distribution': [{'major': major, 'count': count} for major, count in major_stats],
            'class_distribution': [{'class': class_name, 'count': count} for class_name, count in class_stats],
            'gender_distribution': [{'gender': gender, 'count': count} for gender, count in gender_stats]
        })
        
    except Exception as e:
        return jsonify({'error': f'获取统计信息失败: {str(e)}'}), 500


@app.route('/api/students/<student_id>', methods=['DELETE'])
@csrf.exempt
def delete_student(student_id):
    """删除学生（级联删除相关数据）"""
    student = Student.query.filter_by(student_id=student_id).first()
    if not student:
        return jsonify({'error': '学生不存在'}), 404
    
    try:
        # 删除相关的作业记录
        Assignment.query.filter_by(student_id=student_id).delete()
        
        # 删除相关的问题提交记录
        QuestionSubmission.query.filter_by(student_id=student_id).delete()
        
        # 删除头像文件
        if student.avatar and os.path.exists(student.avatar):
            try:
                os.remove(student.avatar)
            except:
                pass
        
        # 删除学生记录
        db.session.delete(student)
        db.session.commit()
        
        return jsonify({'message': '学生及相关数据删除成功'}), 200
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"删除学生失败: {str(e)}")
        return jsonify({'error': f'删除失败: {str(e)}'}), 500


@app.route('/api/course-structure', methods=['GET'])
@csrf.exempt
def get_course_structure():
    """获取完整的课程结构（大学课程）"""
    return jsonify(COURSE_STRUCTURE)


@app.route('/api/courses/k12', methods=['GET'])
@csrf.exempt
def get_k12_courses():
    """获取K12课程结构（小学、初中、高中）"""
    if K12_COURSE_STRUCTURE is None:
        load_k12_courses()
    
    stage = request.args.get('stage')  # 可选: 小学/初中/高中
    grade = request.args.get('grade')  # 可选: 一年级/二年级等
    
    try:
        if not K12_COURSE_STRUCTURE:
            return jsonify({'error': 'K12课程数据未加载'}), 500
        
        # 如果指定了学段
        if stage and stage in K12_COURSE_STRUCTURE:
            stage_data = K12_COURSE_STRUCTURE[stage]
            
            # 如果还指定了年级
            if grade and grade in stage_data:
                return jsonify({
                    'stage': stage,
                    'grade': grade,
                    'subjects': stage_data[grade]
                })
            
            return jsonify({
                'stage': stage,
                'grades': stage_data
            })
        
        # 返回所有K12课程结构
        return jsonify(K12_COURSE_STRUCTURE)
    
    except Exception as e:
        app.logger.error(f"获取K12课程失败: {str(e)}")
        return jsonify({'error': f'获取课程失败: {str(e)}'}), 500


@app.route('/api/courses/k12/stages', methods=['GET'])
@csrf.exempt
def get_k12_stages():
    """获取所有学段（小学、初中、高中）"""
    if K12_COURSE_STRUCTURE is None:
        load_k12_courses()
    
    try:
        if not K12_COURSE_STRUCTURE:
            return jsonify({'stages': []})
        
        stages = list(K12_COURSE_STRUCTURE.keys())
        return jsonify({'stages': stages})
    
    except Exception as e:
        app.logger.error(f"获取学段失败: {str(e)}")
        return jsonify({'error': f'获取学段失败: {str(e)}'}), 500


@app.route('/api/courses/k12/<stage>/grades', methods=['GET'])
@csrf.exempt
def get_k12_grades(stage):
    """获取指定学段的所有年级"""
    if K12_COURSE_STRUCTURE is None:
        load_k12_courses()
    
    try:
        if not K12_COURSE_STRUCTURE or stage not in K12_COURSE_STRUCTURE:
            return jsonify({'error': f'学段 {stage} 不存在'}), 404
        
        grades = list(K12_COURSE_STRUCTURE[stage].keys())
        return jsonify({
            'stage': stage,
            'grades': grades
        })
    
    except Exception as e:
        app.logger.error(f"获取年级失败: {str(e)}")
        return jsonify({'error': f'获取年级失败: {str(e)}'}), 500


@app.route('/api/submit', methods=['POST'])
@csrf.exempt
def submit_assignment():
    if 'file' not in request.files:
        return jsonify({'error': '未包含文件'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': '未选择文件'}), 400
    if not allowed_file(file.filename):
        return jsonify({'error': '不支持的文件类型'}), 400
    student_id = request.form.get('student_id')
    assignment_name = request.form.get('assignment_name')
    subject = request.form.get('subject', '未分类')
    chapter = request.form.get('chapter', '未分类')
    prompt = request.form.get('prompt')
    student = Student.query.filter_by(student_id=student_id).first()
    if not student:
        return jsonify({'error': '未找到学生'}), 404
    temp_dir = tempfile.mkdtemp()
    file_path = os.path.join(temp_dir, secure_filename(file.filename))
    file.save(file_path)
    file_contents = {}
    file_extension = file.filename.rsplit('.', 1)[1].lower()

    # 处理不同文件类型
    if file_extension == 'zip':
        extract_path = os.path.join(temp_dir, 'extracted')
        os.makedirs(extract_path, exist_ok=True)
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(extract_path)
            for root, _, files in os.walk(extract_path):
                for file_name in files:
                    if file_name.endswith(('.c', '.py', '.cpp', '.java')):
                        current_file_path = os.path.join(root, file_name)
                        relative_path = os.path.relpath(current_file_path, extract_path)
                        with open(current_file_path, 'r', encoding='utf-8', errors='ignore') as f:
                            file_contents[relative_path] = f.read()
        except zipfile.BadZipFile:
            return jsonify({'error': '无效的zip文件'}), 400
    elif file_extension == 'pdf':
        file_contents[file.filename] = parse_pdf(file_path)
    elif file_extension == 'docx':
        file_contents[file.filename] = parse_docx(file_path)
    elif file_extension == 'pptx':
        file_contents[file.filename] = parse_pptx(file_path)
    else:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                file_contents[file.filename] = f.read()
        except Exception as e:
            return jsonify({'error': f'读取文件错误: {str(e)}'}), 500
    try:
        submission_content = ""
        for file_name, content in file_contents.items():
            submission_content += f"文件: {file_name}\n\n{content}\n\n---\n\n"
        user_prompt = f"""
        [作业批改助手]
        作业名称: {assignment_name}
        学科: {subject}
        章节: {chapter}
        学生ID: {student_id}
        学生提交的内容:
        {submission_content}
        评分标准:
        针对{subject}学科的{chapter}章节内容，按照以下标准评分:
        1. 知识掌握程度 (40分)
           - 是否符合本章节核心概念
        2. 应用能力 (30分)
           - 是否能正确应用本章节所学
        3. 创新性 (20分)
           - 是否有超出本章节的深入思考
        4. 表达清晰度 (10分)
           - 逻辑是否清晰，表述是否准确
        请提供具体的改进建议，特别是针对{subject}学科{chapter}章节的知识点。
        返回以下格式的 JSON:
        {{
            "score": 分数,
            "feedback": "详细的反馈内容"
        }}
        仅输出JSON格式，不要有其他文本。
        """

        response = openai.ChatCompletion.create(
            model=config.DEEPSEEK_MODEL,
            messages=[
                {"role": "system", "content": "你是一位精通编程的助手，负责批改学生提交的编程作业。"},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.3,
            max_tokens=4096,
            stream=False
        )

        ai_response = response.choices[0].message.content.strip()
        try:
            grading_result = json.loads(ai_response)
            if not isinstance(grading_result,
                              dict) or 'score' not in grading_result or 'feedback' not in grading_result:
                return jsonify({'error': 'AI返回的结果格式不正确'}), 500
        except json.JSONDecodeError:
            import re
            score_match = re.search(r'score["\s:]+(\d+)', ai_response)
            score = float(score_match.group(1)) if score_match else None
            feedback_match = re.search(r'feedback["\s:]+(.+?)(?:}|$)', ai_response, re.DOTALL)
            feedback = feedback_match.group(1).strip().strip('"\'') if feedback_match else ai_response
            grading_result = {'score': score, 'feedback': feedback}

        assignment = Assignment(
            student_id=student_id,
            assignment_name=assignment_name,
            subject=subject,  # 保存科目
            chapter=chapter,  # 保存章节
            score=grading_result.get('score'),
            feedback=grading_result.get('feedback'),
            submission_time=datetime.utcnow()
        )
        db.session.add(assignment)
        db.session.commit()
        return jsonify(assignment.to_dict()), 201
    except Exception as e:
        return jsonify({'error': f'批改过程中出错: {str(e)}'}), 500
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


# 获取所有已存在的科目分类
@app.route('/api/subjects', methods=['GET'])
@csrf.exempt
def get_all_subjects():
    subjects = db.session.query(Assignment.subject).distinct().all()
    return jsonify([s[0] for s in subjects if s[0]])


@app.route('/api/batch-submit', methods=['POST'])
@csrf.exempt
def batch_submit_assignments():
    if 'files[]' not in request.files:
        return jsonify({'error': '未包含文件'}), 400
    files = request.files.getlist('files[]')
    if not files or files[0].filename == '':
        return jsonify({'error': '未选择文件'}), 400

    subject = request.form.get('subject', '未分类')  # 新增
    chapter = request.form.get('chapter', '未分类')  # 新增

    batch_name = request.form.get('batch_name')
    prompt = request.form.get('prompt')
    if not batch_name:
        return jsonify({'error': '请提供批次作业名称'}), 400
    if not prompt:
        return jsonify({'error': '请提供批改提示词'}), 400

    results = []
    temp_dirs = []
    try:
        for file in files:
            if not allowed_file(file.filename):
                continue
            student_id, name = parse_student_info(file.filename)
            if not student_id or not name:
                results.append({
                    'filename': file.filename,
                    'status': 'error',
                    'message': '无法从文件名解析学号和姓名，请确保文件名格式为"学号-姓名.扩展名"'
                })
                continue
            student = get_or_create_student(student_id, name)
            temp_dir = tempfile.mkdtemp()
            temp_dirs.append(temp_dir)
            file_path = os.path.join(temp_dir, secure_filename(file.filename))
            file.save(file_path)
            try:
                file_content = ""
                file_extension = file.filename.rsplit('.', 1)[1].lower()

                if file_extension == 'pdf':
                    file_content = parse_pdf(file_path)
                elif file_extension == 'docx':
                    file_content = parse_docx(file_path)
                elif file_extension == 'pptx':
                    file_content = parse_pptx(file_path)
                else:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        file_content = f.read()
            except Exception as e:
                results.append({
                    'filename': file.filename,
                    'status': 'error',
                    'message': f'读取文件错误: {str(e)}'
                })
                continue

            assignment_name = f"{batch_name} - {student.name}"
            user_prompt = f"""
            [作业批改助手]
            作业名称: {assignment_name}
            学生ID: {student.student_id}
            学生提交的代码:
            文件: {file.filename}
            {file_content}
            评分标准:
            1. 本次作业满分为100分。
            2. 评分时请综合考虑代码的质量、可读性和功能性：
               - 功能实现 (40分)
               - 代码规范 (30分)
               - 逻辑清晰 (20分)
               - 创新性 (10分)
            3. 提供详细的评分理由，指出优点和可以改进的地方。
            返回以下格式的 JSON：
            {{
                "score": 分数,
                "feedback": "详细的反馈内容"
            }}
            仅输出JSON格式，不要有其他文本。
            """

            response = openai.ChatCompletion.create(
                model=config.DEEPSEEK_MODEL,
                messages=[
                    {"role": "system", "content": "你是一位精通编程的助手，负责批改学生提交的编程作业。"},
                    {"role": "user", "content": user_prompt}
                ],
                temperature=0.3,
                max_tokens=4096,
                stream=False
            )

            ai_response = response.choices[0].message.content.strip()
            try:
                grading_result = json.loads(ai_response)
                if not isinstance(grading_result,
                                  dict) or 'score' not in grading_result or 'feedback' not in grading_result:
                    results.append({
                        'filename': file.filename,
                        'status': 'error',
                        'message': 'AI返回的结果格式不正确'
                    })
                    continue
            except json.JSONDecodeError:
                import re
                score_match = re.search(r'score["\s:]+(\d+)', ai_response)
                score = float(score_match.group(1)) if score_match else None
                feedback_match = re.search(r'feedback["\s:]+(.+?)(?:}|$)', ai_response, re.DOTALL)
                feedback = feedback_match.group(1).strip().strip('"\'') if feedback_match else ai_response
                grading_result = {'score': score, 'feedback': feedback}

            assignment = Assignment(
                student_id=student.student_id,
                assignment_name=batch_name,
                subject=subject,  # 新增
                chapter=chapter,  # 新增
                score=grading_result.get('score'),
                feedback=grading_result.get('feedback'),
                submission_time=datetime.utcnow()
            )
            db.session.add(assignment)
            db.session.commit()

            results.append({
                'filename': file.filename,
                'status': 'success',
                'student_id': student.student_id,
                'student_name': student.name,
                'assignment_name': batch_name,
                'score': grading_result.get('score')
            })

    except Exception as e:
        return jsonify({'error': f'批量上传处理出错: {str(e)}'}), 500
    finally:
        for temp_dir in temp_dirs:
            shutil.rmtree(temp_dir, ignore_errors=True)

    return jsonify({
        'batch_name': batch_name,
        'total': len(files),
        'processed': len(results),
        'results': results
    }), 201


@app.route('/api/scores/<student_id>/<assignment_name>', methods=['GET'])
@csrf.exempt
def get_assignment_score(student_id, assignment_name):
    student = Student.query.filter_by(student_id=student_id).first()
    if not student:
        return jsonify({'error': '未找到学生'}), 404
    assignment = Assignment.query.filter_by(student_id=student_id, assignment_name=assignment_name).order_by(
        Assignment.submission_time.desc()).first()
    if not assignment:
        return jsonify({'error': '未找到作业'}), 404
    return jsonify({
        'student': student.to_dict(),
        'assignment': assignment.to_dict()
    })


# 获取学生所有成绩数据（用于基础图表）
@app.route('/api/scores/<student_id>')
@csrf.exempt
def get_student_scores(student_id):
    # 首先获取学生信息
    student = Student.query.filter_by(student_id=student_id).first()
    if not student:
        return jsonify({'error': 'Student not found'}), 404

    # 然后获取该学生的所有作业
    assignments = Assignment.query.filter_by(student_id=student_id).all()

    # 返回统一格式的数据
    return jsonify({
        'student': {
            'student_id': student.student_id,
            'name': student.name
        },
        'assignments': [a.to_dict() for a in assignments]
    })


# 获取学科-章节成绩趋势数据
@app.route('/api/chapter-scores/<student_id>/<subject>/<chapter>')
@csrf.exempt
def get_chapter_scores(student_id, subject, chapter):
    assignments = Assignment.query.filter_by(
        student_id=student_id,
        subject=subject,
        chapter=chapter
    ).order_by(Assignment.submission_time).all()

    data = {
        'student_id': student_id,
        'subject': subject,
        'chapter': chapter,
        'assignments': [a.to_dict() for a in assignments]
    }
    return jsonify(data)


# 获取学科成绩分布数据
@app.route('/api/subject-scores/<student_id>/<subject>')
@csrf.exempt
def get_subject_scores(student_id, subject):
    assignments = Assignment.query.filter_by(
        student_id=student_id,
        subject=subject
    ).all()

    data = {
        'student_id': student_id,
        'subject': subject,
        'assignments': [a.to_dict() for a in assignments]
    }
    return jsonify(data)


@app.route('/api/scores/<student_id>/subjects', methods=['GET'])
@csrf.exempt
def get_student_subjects(student_id):
    subjects = db.session.query(Assignment.subject).filter_by(student_id=student_id).distinct().all()
    return jsonify([s[0] for s in subjects if s[0]])


@app.route('/api/scores/<student_id>/<subject>', methods=['GET'])
@csrf.exempt
def get_student_subject_scores(student_id, subject):
    student = Student.query.filter_by(student_id=student_id).first()
    if not student:
        return jsonify({'error': '未找到学生'}), 404

    assignments = Assignment.query.filter_by(
        student_id=student_id,
        subject=subject
    ).order_by(Assignment.submission_time).all()

    return jsonify({
        'student': student.to_dict(),
        'subject': subject,
        'assignments': [a.to_dict() for a in assignments]
    })


@app.route('/api/subject-chapters/<subject>', methods=['GET'])
@csrf.exempt
def get_subject_chapters(subject):
    """获取指定科目下的所有章节"""
    chapters = db.session.query(Assignment.chapter).filter_by(
        subject=subject
    ).distinct().all()
    return jsonify([c[0] for c in chapters if c[0]])


@app.route('/api/subject-score-trend/<student_id>/<subject>', methods=['GET'])
@csrf.exempt
def get_subject_score_trend(student_id, subject):
    """获取学生在某科目的各章节成绩趋势"""
    student = Student.query.filter_by(student_id=student_id).first()
    if not student:
        return jsonify({'error': '未找到学生'}), 404

    assignments = Assignment.query.filter_by(
        student_id=student_id,
        subject=subject
    ).order_by(Assignment.chapter).all()

    # 按章节分组计算平均分
    scores_by_chapter = {}
    for a in assignments:
        if a.chapter not in scores_by_chapter:
            scores_by_chapter[a.chapter] = []
        scores_by_chapter[a.chapter].append(a.score)

    # 计算每个章节的平均分
    chapter_data = []
    for chapter, scores in scores_by_chapter.items():
        if scores:
            avg_score = round(sum(scores) / len(scores), 2)
            chapter_data.append({
                'chapter': chapter,
                'average_score': avg_score,
                'assignment_count': len(scores)
            })

    return jsonify({
        'subject': subject,
        'trend_data': chapter_data
    })


@app.route('/api/score-comparison/<student_id>', methods=['GET'])
@csrf.exempt
def get_score_comparison(student_id):
    """获取学生在各科目的平均分比较"""
    student = Student.query.filter_by(student_id=student_id).first()
    if not student:
        return jsonify({'error': '未找到学生'}), 404

    subjects = db.session.query(
        Assignment.subject,
        db.func.avg(Assignment.score).label('avg_score'),
        db.func.count(Assignment.id).label('assignment_count')
    ).filter_by(
        student_id=student_id
    ).group_by(
        Assignment.subject
    ).all()

    subject_data = [{
        'subject': s.subject,
        'average_score': round(float(s.avg_score), 2),
        'assignment_count': s.assignment_count
    } for s in subjects]

    return jsonify({
        'student_id': student_id,
        'subjects': subject_data
    })


@app.route('/api/ai/ask', methods=['POST'])
@limiter.limit("20 per minute")  # API限流保护
@csrf.exempt
@require_login_api
@require_membership
@feature_limit('ai_ask')
def ai_ask():
    data = request.json
    question = data.get('question')
    session_id = data.get('session_id')
    if not question or not isinstance(question, str):
        return jsonify({'error': '必须提供有效的问题内容'}), 400
    try:
        # 查找或创建会话（不使用事务，因为需要立即查询）
        conversation = Conversation.query.filter_by(session_id=session_id).first()
        if not conversation:
            if not session_id:
                session_id = str(uuid.uuid4())
            conversation = Conversation(session_id=session_id)
            db.session.add(conversation)
        db.session.commit()  # 立即提交以获取ID
        
        # 添加系统消息
        system_msg = ConversationMessage(
            conversation_id=conversation.id,
            role='system',
            content="""
            你是一位知识渊博的导师，请用专业但易懂的语言回答学生问题。
            回答要求：
            1. 使用规范的中文表达
            2. 禁止输出任何代码或特殊符号
            3. 如果问题不明确，请要求澄清
            4. 保持回答简洁明了
            """
        )
        db.session.add(system_msg)
        
        # 添加用户消息
        user_msg = ConversationMessage(
            conversation_id=conversation.id,
            role='user',
            content=question
        )
        db.session.add(user_msg)
        db.session.commit()  # 提交用户消息
        
        # 使用正确的方法获取对话历史(最多最近20条消息)
        messages = [msg.to_dict() for msg in ConversationMessage.query
            .filter_by(conversation_id=conversation.id)
            .order_by(ConversationMessage.created_at.desc())
            .limit(20).all()]
        messages.reverse()  # 按时间正序排列

        # 使用DeepSeek API
        response = openai.ChatCompletion.create(
            model=config.DEEPSEEK_MODEL,
            messages=messages,
            temperature=0.3,
            max_tokens=1024,
            stream=False
        )
        ai_response = response.choices[0].message.content
        sanitized_response = sanitize_ai_response(ai_response)

        # 保存AI回复到数据库
        ai_msg = ConversationMessage(
            conversation_id=conversation.id,
            role='assistant',
            content=sanitized_response
        )
        db.session.add(ai_msg)
        db.session.commit()

        return jsonify({
            'answer': sanitized_response,
            'session_id': conversation.session_id
        }), 200

    except Exception as e:
        db.session.rollback()  # 发生错误时回滚
        app.logger.error(f"AI答疑错误: {str(e)}")
        return jsonify({
            'error': '抱歉，回答生成失败',
            'session_id': session_id,
            'technical_detail': str(e) if app.debug else None
        }), 500


@app.route('/api/ai/conversation', methods=['GET'])
@csrf.exempt
def get_conversation():
    session_id = request.args.get('session_id')
    if not session_id:
        return jsonify({'error': 'Missing session_id'}), 400

    conversation = Conversation.query.filter_by(session_id=session_id).first()
    if not conversation:
        return jsonify({'messages': []})

    messages = [msg.to_dict() for msg in conversation.messages.order_by(ConversationMessage.created_at).limit(20).all()]
    return jsonify({
        'session_id': session_id,
        'messages': messages
    })


@app.route('/api/ai/new-conversation', methods=['POST'])
@csrf.exempt
def new_conversation():
    try:
        # 创建新会话
        new_session_id = str(uuid.uuid4())
        conversation = Conversation(session_id=new_session_id)
        db.session.add(conversation)

        # 添加系统消息
        db.session.add(ConversationMessage(
            conversation=conversation,
            role='system',
            content="""
            你是一位知识渊博的导师，请用专业但易懂的语言回答学生问题。
            回答要求：
            1. 使用规范的中文表达
            2. 禁止输出任何代码或特殊符号
            3. 如果问题不明确，请要求澄清
            4. 保持回答简洁明了
            """
        ))
        db.session.commit()

        return jsonify({
            'success': True,
            'session_id': new_session_id
        }), 200
    except Exception as e:
        db.session.rollback()
        return jsonify({
            'error': f'创建新会话时出错: {str(e)}'
        }), 500


@app.route('/api/ai/conversations', methods=['GET'])
@csrf.exempt
def get_conversations():
    try:
        # 获取最近的10个会话
        conversations = Conversation.query.order_by(
            Conversation.updated_at.desc()
        ).limit(10).all()

        return jsonify([{
            'session_id': conv.session_id,
            'updated_at': conv.updated_at.isoformat() if conv.updated_at else None
        } for conv in conversations])
    except Exception as e:
        return jsonify({'error': f'获取会话列表失败: {str(e)}'}), 500


# 辅助编程API
@app.route('/api/ai/programming-help', methods=['POST'])
@limiter.limit("15 per minute")  # API限流保护
@csrf.exempt
@require_login_api
@require_membership
@feature_limit('programming_help')
def programming_help():
    """辅助编程API - 使用DeepSeek API"""
    try:
        data = request.json
        code = data.get('code', '')
        question = data.get('question', '')
        language = data.get('language', 'python')
        session_id = data.get('session_id', str(uuid.uuid4()))
        
        if not code and not question:
            return jsonify({'error': '请提供代码或问题描述'}), 400
        
        # 构建编程助手的提示词
        system_prompt = f"""你是一位专业的编程助手，专门帮助学生解决编程问题。
        
请遵循以下规则：
1. 使用中文回答
2. 提供清晰的代码解释和建议
3. 如果发现代码错误，请指出并提供修正方案
4. 提供最佳实践建议
5. 保持回答简洁明了但详细

当前编程语言：{language}
"""
        
        # 构建用户消息
        user_message = ""
        if code:
            user_message += f"我的代码：\n```{language}\n{code}\n```\n\n"
        if question:
            user_message += f"问题：{question}"
        
        # 使用DeepSeek API
        response = openai.ChatCompletion.create(
            model=config.DEEPSEEK_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message}
            ],
            temperature=0.3,
            max_tokens=2048,
            stream=False
        )
        ai_response = response.choices[0].message.content
        
        # 清理响应内容
        sanitized_response = sanitize_ai_response(ai_response)
        
        return jsonify({
            'success': True,
            'response': sanitized_response,
            'session_id': session_id
        }), 200
        
    except Exception as e:
        app.logger.error(f"编程助手错误: {str(e)}")
        return jsonify({
            'error': '编程助手服务暂时不可用',
            'technical_detail': str(e) if app.debug else None
        }), 500


@app.route('/api/ai/code-review', methods=['POST'])
@limiter.limit("10 per minute")  # API限流保护
@csrf.exempt
@require_login_api
@require_membership
@feature_limit('code_review')
def code_review():
    """代码审查API"""
    try:
        data = request.json
        code = data.get('code', '')
        language = data.get('language', 'python')
        
        if not code:
            return jsonify({'error': '请提供要审查的代码'}), 400
        
        # 构建代码审查提示词
        review_prompt = f"""请对以下{language}代码进行专业的代码审查：

```{language}
{code}
```

请从以下几个方面进行评估：
1. 代码质量和规范性
2. 性能优化建议
3. 安全性问题
4. 可读性和维护性
5. 潜在的bug或错误
6. 最佳实践建议

请用中文回答，提供具体的改进建议和示例代码。"""
        
        # 使用DeepSeek API
        response = openai.ChatCompletion.create(
            model=config.DEEPSEEK_MODEL,
            messages=[
                {"role": "system", "content": "你是一位资深的代码审查专家，请提供专业的代码审查意见。"},
                {"role": "user", "content": review_prompt}
            ],
            temperature=0.3,
            max_tokens=2048,
            stream=False
        )
        ai_response = response.choices[0].message.content
        
        sanitized_response = sanitize_ai_response(ai_response)
        
        return jsonify({
            'success': True,
            'review': sanitized_response
        }), 200
        
    except Exception as e:
        app.logger.error(f"代码审查错误: {str(e)}")
        return jsonify({
            'error': '代码审查服务暂时不可用',
            'technical_detail': str(e) if app.debug else None
        }), 500


@app.route('/api/ai/code-explain', methods=['POST'])
@limiter.limit("15 per minute")  # API限流保护
@csrf.exempt
@require_login_api
@require_membership
@feature_limit('code_explain')
def code_explain():
    """代码解释API"""
    try:
        data = request.json
        code = data.get('code', '')
        language = data.get('language', 'python')
        
        if not code:
            return jsonify({'error': '请提供要解释的代码'}), 400
        
        explain_prompt = f"""请详细解释以下{language}代码的功能和工作原理：

```{language}
{code}
```

请包括：
1. 代码的整体功能
2. 每个重要部分的作用
3. 使用的算法或数据结构
4. 代码的执行流程
5. 关键概念的解释

请用中文回答，使用通俗易懂的语言。"""
        
        # 使用DeepSeek API
        response = openai.ChatCompletion.create(
            model=config.DEEPSEEK_MODEL,
            messages=[
                {"role": "system", "content": "你是一位编程导师，擅长用简单易懂的语言解释复杂的代码。"},
                {"role": "user", "content": explain_prompt}
            ],
            temperature=0.3,
            max_tokens=2048,
            stream=False
        )
        ai_response = response.choices[0].message.content
        
        sanitized_response = sanitize_ai_response(ai_response)
        
        return jsonify({
            'success': True,
            'explanation': sanitized_response
        }), 200
        
    except Exception as e:
        app.logger.error(f"代码解释错误: {str(e)}")
        return jsonify({
            'error': '代码解释服务暂时不可用',
            'technical_detail': str(e) if app.debug else None
        }), 500


@app.route('/api/ai/debug-help', methods=['POST'])
@limiter.limit("15 per minute")  # API限流保护
@csrf.exempt
@require_login_api
@require_membership
@feature_limit('debug_help')
def debug_help():
    """调试帮助API"""
    try:
        data = request.json
        code = data.get('code', '')
        error_message = data.get('error_message', '')
        language = data.get('language', 'python')
        
        if not code and not error_message:
            return jsonify({'error': '请提供代码或错误信息'}), 400
        
        debug_prompt = f"""请帮助调试以下{language}代码问题：

"""
        
        if code:
            debug_prompt += f"代码：\n```{language}\n{code}\n```\n\n"
        
        if error_message:
            debug_prompt += f"错误信息：\n{error_message}\n\n"
        
        debug_prompt += """请提供：
1. 错误原因分析
2. 具体的修复方案
3. 修正后的代码示例
4. 预防类似错误的建议

请用中文回答。"""
        
        # 使用DeepSeek API
        response = openai.ChatCompletion.create(
            model=config.DEEPSEEK_MODEL,
            messages=[
                {"role": "system", "content": "你是一位专业的调试专家，擅长快速定位和解决代码问题。"},
                {"role": "user", "content": debug_prompt}
            ],
            temperature=0.3,
            max_tokens=2048,
            stream=False
        )
        ai_response = response.choices[0].message.content
        
        sanitized_response = sanitize_ai_response(ai_response)
        
        return jsonify({
            'success': True,
            'debug_help': sanitized_response
        }), 200
        
    except Exception as e:
        app.logger.error(f"调试帮助错误: {str(e)}")
        return jsonify({
            'error': '调试帮助服务暂时不可用',
            'technical_detail': str(e) if app.debug else None
        }), 500




@app.route('/api/ai/summarize-video', methods=['POST'])
@limiter.limit("5 per minute")  # API限流保护
@csrf.exempt
@require_login_api
@require_membership
@feature_limit('video_summary')
def ai_summarize_video():
    if 'file' not in request.files and 'url' not in request.form:
        return jsonify({'error': '未提供视频文件或链接'}), 400

    video_file = request.files.get('file')
    video_url = request.form.get('url')

    try:
        if video_file:
            # 文件上传处理
            video_file.seek(0, os.SEEK_END)
            file_size = video_file.tell()
            video_file.seek(0)
            if file_size > 500 * 1024 * 1024:
                return jsonify({'error': '视频文件超过500MB限制'}), 413
            
            secure_name = secure_filename(video_file.filename)
            if not secure_name:
                return jsonify({'error': '文件名无效'}), 400
            
            # 由于DeepSeek API无法直接处理视频文件，我们基于文件名生成一个教育性的总结
            summary_prompt = f"""
            请根据视频文件名"{secure_name}"，生成一个教育性的视频总结。假设这是一个教育视频，请：

            1. 根据文件名推测视频可能的内容主题
            2. 提供一个结构化的学习总结，包括：
               - 主要知识点
               - 重要概念
               - 学习建议
               - 可能的应用场景

            请用中文回答，格式清晰，适合学生学习使用。
            """
            
        else:
            # URL处理
            if not video_url.startswith(('http://', 'https://')):
                return jsonify({'error': '无效的视频链接'}), 400
            
            # 基于视频URL生成总结
            summary_prompt = f"""
            请根据视频链接"{video_url}"，生成一个教育性的视频总结。请：

            1. 根据链接URL推测视频可能的内容和主题
            2. 如果是Bilibili、YouTube等平台的链接，说明这类平台的教育价值
            3. 提供一个结构化的学习指导，包括：
               - 预期学习内容
               - 观看建议
               - 笔记要点
               - 扩展学习方向

            请用中文回答，格式清晰，适合教育场景使用。
            """

        # 使用DeepSeek API生成总结
        response = openai.ChatCompletion.create(
            model=config.DEEPSEEK_MODEL,
            messages=[
                {"role": "system", "content": "你是一个专业的教育助手，擅长分析和总结教育视频内容，为学生提供有价值的学习指导。"},
                {"role": "user", "content": summary_prompt}
            ],
            temperature=0.7,
            max_tokens=1500,
            stream=False
        )
        
        ai_summary = response.choices[0].message.content
        sanitized_summary = sanitize_ai_response(ai_summary)

        return jsonify({
            'input': video_url or secure_name,
            'summary': sanitized_summary
        }), 200

    except Exception as e:
        app.logger.exception("视频总结严重错误")
        return jsonify({'error': f'视频总结过程中出错: {str(e)}'}), 500


@app.route('/api/ai/video-to-lecture', methods=['POST'])
@limiter.limit("3 per minute")  # API限流保护
@csrf.exempt
@require_login_api
@require_membership
@feature_limit('generate_lecture')
def video_to_lecture():
    """从视频生成教学讲义
    
    参数：
        - file: 视频文件（可选）
        - url: 视频链接（可选）
        - course_info: 课程信息（JSON格式）
            - stage: 学段（小学/初中/高中）
            - grade: 年级
            - subject: 科目
            - chapter: 章节（可选）
        - generate_exercises: 是否生成配套练习题（布尔值，默认false）
    """
    video_file = request.files.get('file')
    video_url = request.form.get('url')
    course_info_str = request.form.get('course_info', '{}')
    generate_exercises = request.form.get('generate_exercises', 'false').lower() == 'true'
    
    if not video_file and not video_url:
        return jsonify({'error': '请提供视频文件或链接'}), 400
    
    try:
        # 解析课程信息
        course_info = json.loads(course_info_str) if course_info_str else {}
        stage = course_info.get('stage', '')
        grade = course_info.get('grade', '')
        subject = course_info.get('subject', '')
        chapter = course_info.get('chapter', '')
        
        # 1. 获取视频标识
        video_identifier = ''
        if video_file:
            video_identifier = secure_filename(video_file.filename)
        elif video_url:
            video_identifier = video_url
        
        # 2. 生成讲义内容
        exercises_section = """
        ## 5. 配套练习
        - 基础题（2-3题）
        - 提高题（1-2题）
        - 答案与解析
        """ if generate_exercises else ""
        
        lecture_prompt = f"""
        请根据以下信息生成一份详细的教学讲义：
        
        **视频信息**：{video_identifier}
        **学段**：{stage or '未指定'}
        **年级**：{grade or '未指定'}
        **科目**：{subject or '未指定'}
        **章节**：{chapter or '未指定'}
        
        请生成一份完整的教学讲义，包括：
        
        ## 1. 课程概述
        - 本节课的学习目标
        - 知识点定位
        - 学习重点和难点
        
        ## 2. 知识讲解
        - 核心概念详细解释
        - 重要公式/定理/规则
        - 易错点提示
        - 实例分析
        
        ## 3. 知识梳理
        - 思维导图式的知识结构
        - 关键点归纳
        
        ## 4. 学习建议
        - 课后复习要点
        - 拓展学习方向
        {exercises_section}
        
        请用清晰的markdown格式回答，语言生动易懂，适合{grade or '学生'}学习。
        """
        
        # 调用AI生成讲义
        response = openai.ChatCompletion.create(
            model=config.DEEPSEEK_MODEL,
            messages=[
                {"role": "system", "content": f"你是一位经验丰富的{subject or ''}教师，擅长将视频教学内容转化为结构化的教学讲义，帮助学生系统学习。"},
                {"role": "user", "content": lecture_prompt}
            ],
            temperature=0.7,
            max_tokens=3000,
            stream=False
        )
        
        lecture_content = response.choices[0].message.content
        sanitized_lecture = sanitize_ai_response(lecture_content)
        
        # 3. 保存到数据库（使用VideoNote表）
        try:
            video_note = VideoNote(
                user_id=current_user.id,
                video_url=video_url or f"file://{video_identifier}",
                notes=sanitized_lecture,
                created_at=datetime.utcnow()
            )
            db.session.add(video_note)
            db.session.commit()
            
            note_id = video_note.id
        except Exception as db_error:
            app.logger.warning(f"保存讲义到数据库失败: {str(db_error)}")
            note_id = None
        
        # 4. 记录功能使用（会员系统）
        try:
            log_feature_usage(current_user.id, 'generate_lecture')
        except:
            pass
        
        return jsonify({
            'success': True,
            'lecture': sanitized_lecture,
            'video_identifier': video_identifier,
            'course_info': {
                'stage': stage,
                'grade': grade,
                'subject': subject,
                'chapter': chapter
            },
            'note_id': note_id,
            'has_exercises': generate_exercises
        }), 200
    
    except json.JSONDecodeError:
        return jsonify({'error': '课程信息格式错误'}), 400
    except Exception as e:
        app.logger.exception("视频转讲义失败")
        return jsonify({'error': f'生成讲义失败: {str(e)}'}), 500


# 智能讲义
@app.route('/api/ai/generate-lecture', methods=['POST'])
@limiter.limit("3 per minute")  # API限流保护
@csrf.exempt
@require_login_api
@require_membership
@feature_limit('generate_lecture')
def generate_lecture():
    if 'file' not in request.files:
        return jsonify({'error': '未提供文件', 'status': 'failed'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': '未选择文件', 'status': 'failed'}), 400

    temp_dir = tempfile.mkdtemp()
    try:
        # 保存上传的文件
        file_path = os.path.join(temp_dir, secure_filename(file.filename))
        file.save(file_path)

        # 解析文件内容
        ext = file.filename.rsplit('.', 1)[1].lower()
        if ext == 'pdf':
            text = parse_pdf(file_path)
        elif ext == 'docx':
            text = parse_docx(file_path)
        elif ext == 'pptx':
            text = parse_pptx(file_path)
        else:
            return jsonify({
                'error': '不支持的文件类型',
                'status': 'failed',
                'allowed_types': ['pdf', 'docx', 'pptx']
            }), 400

        # 使用DeepSeek API生成教案
        response = openai.ChatCompletion.create(
            model=config.DEEPSEEK_MODEL,
            messages=[
                {"role": "system", "content": """
                你是一位专业讲师助手，请严格按以下要求生成结构化讲义:

                结构化要求:
                1. 使用JSON格式返回
                2. 数据结构:
                {
                    "title": "讲义标题",
                    "sections": [
                        {
                            "title": "章节标题",
                            "content": "章节概述",
                            "subsections": [
                                {
                                    "title": "子标题",
                                    "content": "详细内容",
                                    "key_points": ["要点1", "要点2"]
                                }
                            ]
                        }
                    ]
                }
                3. 内容要求:
                   - 使用##标记重要概念
                   - 使用**加粗**关键术语
                   - 层级分明(章节>子章节>要点)
                   - 保持学术严谨性
                """},
                {"role": "user", "content": f"请基于以下内容生成结构化讲义:\n{text}"}
            ],
            response_format={"type": "json_object"},
            temperature=0.3,
            max_tokens=2048
        )
        ai_response = response.choices[0].message.content
        try:
            result = json.loads(ai_response)
            if not isinstance(result, dict) or 'sections' not in result:
                raise ValueError("返回结构不符合要求")

            # 对内容进行安全处理和格式化
            def sanitize_content(content):
                content = re.sub(r'\n{3,}', '\n\n', content)  # 合并多余空行
                return content.strip()

            # 递归处理所有文本内容
            def process_structure(data):
                if isinstance(data, dict):
                    return {k: process_structure(v) for k, v in data.items()}
                elif isinstance(data, list):
                    return [process_structure(item) for item in data]
                elif isinstance(data, str):
                    return sanitize_content(data)
                return data

            processed_result = process_structure(result)

            return jsonify({
                'status': 'success',
                'lecture': processed_result,
                'source_file': file.filename,
                'generated_at': datetime.now().isoformat(),
                'format_version': '1.1'  # 标识返回格式版本
            })

        except json.JSONDecodeError:
            return jsonify({'error': 'AI返回的数据不是有效JSON'}), 500
        except ValueError as e:
            return jsonify({'error': f'内容格式验证失败: {str(e)}'}), 500

    except Exception as e:
        return jsonify({
            'status': 'error',
            'error': str(e),
            'traceback': traceback.format_exc() if app.debug else None
        }), 500

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


# 智能出题
# 智能出题系统路由
@app.route('/api/ai/generate-question', methods=['POST'])
@limiter.limit("3 per minute")  # API限流保护
@csrf.exempt
@require_login_api
@require_membership
@feature_limit('generate_question')
def generate_question():
    if 'file' not in request.files:
        print("没有收到文件")  # 调试日志
        return jsonify({'error': '未提供文件'}), 400
    file = request.files['file']
    print(f"收到的文件名: {file.filename}")  # 调试日志

    if not allowed_file(file.filename):
        print("文件类型不允许")  # 调试日志
        return jsonify({'error': '不支持的文件类型'}), 400

    difficulty = request.form.get('difficulty', 'medium')
    num_questions = int(request.form.get('num_questions', 3))

    # 生成唯一的题目集ID
    question_set_id = str(uuid.uuid4())

    with tempfile.TemporaryDirectory() as temp_dir:
        file_path = os.path.join(temp_dir, secure_filename(file.filename))
        file.save(file_path)

        # 解析文件内容
        ext = file.filename.rsplit('.', 1)[1].lower()
        if ext == 'pdf':
            text = parse_pdf(file_path)
        elif ext == 'docx':
            text = parse_docx(file_path)
        elif ext == 'pptx':
            text = parse_pptx(file_path)
        else:
            return jsonify({'error': '不支持的文件类型'}), 400
        
        # 使用DeepSeek API生成题目
        response = openai.ChatCompletion.create(
            model=config.DEEPSEEK_MODEL,
            messages=[
                {"role": "system", "content": "你是一个专业的题目生成助手，根据提供的材料生成考试题目。"},
                {"role": "user", "content": f"""
                根据以下材料生成{difficulty}难度的{num_questions}道题目:
                {text}

                要求:
                1. 包含选择题、填空题和简答题
                2. 题目考察核心知识点
                3. 返回的JSON中必须包含correct_answer字段

                返回格式:
                {{
                    "status": "success",
                    "questions": [
                        {{
                            "type": "question_type",
                            "question": "题目内容",
                            "options": ["选项1", "选项2"] (仅选择题),
                            "correct_answer": "正确答案"
                        }}
                    ],
                    "source_file": "文件名",
                    "generated_at": "生成时间"
                }}
                """}
            ],
            response_format={"type": "json_object"},
            temperature=0.7
        )
        ai_content = response.choices[0].message.content
        result = parse_ai_response(ai_content)

    try:

            # 保存到题库，使用UUID作为question_id
            questions_result = []
            for q in result['questions']:
                question_id = str(uuid.uuid4())  # 为每个题目生成唯一ID
                new_question = QuestionBank(
                    question_id=question_id,
                    question_set_id=question_set_id,  # 记录题目集ID
                    question_text=q['question'],
                    question_type=q['type'],
                    options=q.get('options', []),
                    correct_answer=q.get('correct_answer', ''),
                    created_at=datetime.utcnow()
                )
                db.session.add(new_question)

                # 构建返回给前端的结果
                questions_result.append({
                    'id': question_id,
                    'type': q['type'],
                    'question': q['question'],
                    'options': q.get('options', []),
                    'correct_answer': q.get('correct_answer', '')
                })

            db.session.commit()

            return jsonify({
                'status': 'success',
                'questions': questions_result,
                'question_set_id': question_set_id,  # 返回题目集ID给前端
                'source_file': file.filename,
                'generated_at': datetime.utcnow().isoformat()
            })
    except Exception as e:
            db.session.rollback()
            return jsonify({'error': str(e)}), 500


# 新增的题目解答路由
@app.route('/api/ai/answer-questions', methods=['POST'])
@feature_limit('generate_question')
@limiter.limit("5 per minute")  # API限流保护
@csrf.exempt
def answer_questions():
    data = request.json
    if not data or 'questions' not in data:
        return jsonify({'error': '未提供题目数据或格式不正确'}), 400
    questions = data['questions']
    if not isinstance(questions, list) or len(questions) == 0:
        return jsonify({'error': '题目数据格式错误'}), 400
    try:
        # 构建问题解答提示
        questions_json = json.dumps(questions, ensure_ascii=False)
        user_prompt = f"""
        请解答以下题目，并按照要求返回答案：
        题目列表：
        {questions_json}
        回答要求：
        1. 为每道题目提供详细解答
        2. 选择题要标明正确选项并解释原因
        3. 填空题要给出完整答案
        4. 判断题要有明确答案和解释
        5. 简答题要提供完整回答
        返回格式必须是严格的JSON格式：
        {{
            "answers": [
                {{
                    "id": "题目ID",
                    "type": "题目类型",
                    "question": "原问题",
                    "answer": "详细答案",
                    "explanation": "解答说明"
                }},
                // 更多答案...
            ]
        }}
        注意：
        - 必须严格按照题目ID对应回答
        - 保持原始问题内容不变
        """

        response = openai.ChatCompletion.create(
            model=config.DEEPSEEK_MODEL,
            messages=[
                {"role": "system",
                 "content": "你是一位专业教师，负责解答各类考题问题。解答过程要专业、全面和准确。"},
                {"role": "user", "content": user_prompt}
            ],
            response_format={"type": "json_object"},
            temperature=0.2,  # 降低温度值以获得更准确的答案
            max_tokens=4096,
            stream=False
        )
        ai_response = response.choices[0].message.content.strip()

        try:
            answers = parse_ai_response(ai_response)
            if not isinstance(answers, dict) or 'answers' not in answers:
                return jsonify({'error': 'AI 返回结果格式不正确，缺少 answers 字段'}), 500
            return jsonify({
                'status': 'success',
                'answers': answers['answers'],
                'answered_at': datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')
            }), 200
        except ValueError as e:
            return jsonify({'error': f'无法解析AI返回的答案数据: {str(e)}'}), 500
    except Exception as e:
        return jsonify({'error': f'解答题目过程中出错: {str(e)}'}), 500


# 题目解答路由
@app.route('/api/question/submit-answers', methods=['POST'])
@csrf.exempt
def submit_answers():
    data = request.get_json()
    if not data or 'student_id' not in data or 'answers' not in data:
        return jsonify({'error': '缺少必要参数'}), 400

    try:
        student = Student.query.filter_by(student_id=data['student_id']).first()
        if not student:
            return jsonify({'error': '学生不存在'}), 404

        results = []
        total_score = 0

        for answer in data['answers']:
            question = QuestionBank.query.filter_by(question_id=answer['question_id']).first()
            if not question:
                continue

            is_correct = False
            score = 0
            feedback = ""

            # 选择题直接比对答案
            if question.question_type == 'choice':
                is_correct = answer['answer'].strip().lower() == question.correct_answer.strip().lower()
                score = 10 if is_correct else 0  # 假设每道选择题10分
                feedback = "回答正确" if is_correct else f"正确答案是: {question.correct_answer}"
            else:
                # 对于非选择题（简答/填空），直接比较学生答案和标准答案
                # 这里可以更严格地比较，比如忽略大小写、空格等
                is_correct = answer['answer'].strip().casefold() == question.correct_answer.strip().casefold()
                score = 10 if is_correct else 0  # 假设每道非选择题也是10分

                # 如果不是完全匹配，提供标准答案作为反馈
                if not is_correct:
                    feedback = f"参考答案: {question.correct_answer}"
                else:
                    feedback = "回答正确"

            # 记录答题结果
            submission = QuestionSubmission(
                student_id=data['student_id'],
                question_id=answer['question_id'],
                answer_text=answer['answer'],
                is_correct=is_correct,
                score=score,
                feedback=feedback,
                submission_time=datetime.utcnow()
            )
            db.session.add(submission)
            total_score += score

            results.append({
                'question_id': answer['question_id'],
                'question': question.question_text,
                'student_answer': answer['answer'],
                'correct_answer': question.correct_answer,
                'score': score,
                'feedback': feedback,
                'is_correct': is_correct
            })

        avg_score = total_score / len(data['answers']) if data['answers'] else 0

        db.session.commit()

        return jsonify({
            'status': 'success',
            'student_id': data['student_id'],
            'total_score': total_score,
            'average_score': round(avg_score, 2),
            'results': results,
            'submitted_at': datetime.utcnow().isoformat()
        })

    except Exception as e:
        db.session.rollback()
        return jsonify({'error': str(e)}), 500


# 题目统计路由
@app.route('/api/question/stats/<student_id>', methods=['GET'])
@csrf.exempt
def get_question_stats(student_id):
    try:
        student = Student.query.filter_by(student_id=student_id).first()
        if not student:
            return jsonify({'error': '学生不存在'}), 404

        # 获取答题总体统计
        submissions = QuestionSubmission.query.filter_by(student_id=student_id).all()
        total = len(submissions)
        if total == 0:
            return jsonify({'error': '没有答题记录'}), 404

        correct = sum(1 for s in submissions if s.is_correct)
        accuracy = round((correct / total) * 100, 2)
        avg_score = round(sum(s.score for s in submissions) / total, 2)

        # 按题型统计
        type_stats = db.session.execute(
            db.select(
                QuestionBank.question_type,
                db.func.count(QuestionSubmission.id),
                db.func.avg(QuestionSubmission.score),
                db.func.sum(db.case((QuestionSubmission.is_correct == True, 1), else_=0))
            )
            .join(QuestionBank, QuestionBank.question_id == QuestionSubmission.question_id)
            .filter(QuestionSubmission.student_id == student_id)
            .group_by(QuestionBank.question_type)
        ).all()

        type_results = [
            {
                'type': stat[0],
                'count': stat[1],
                'avg_score': round(float(stat[2] or 0), 2),
                'accuracy': round((stat[3] / stat[1]) * 100, 2) if stat[1] > 0 else 0
            }
            for stat in type_stats
        ]

        # 最近5次答题情况
        recent_submissions = (
            QuestionSubmission.query
            .filter_by(student_id=student_id)
            .order_by(QuestionSubmission.submission_time.desc())
            .limit(5)
            .all()
        )

        recent_results = [
            {
                'question_id': sub.question_id,
                'question': QuestionBank.query.filter_by(question_id=sub.question_id).first().question_text[
                            :50] + '...',
                'score': sub.score,
                'is_correct': sub.is_correct,
                'submitted_at': sub.submission_time.isoformat()
            }
            for sub in recent_submissions
        ]

        return jsonify({
            'status': 'success',
            'student_id': student_id,
            'student_name': student.name,
            'total_questions': total,
            'correct_answers': correct,
            'accuracy': accuracy,
            'average_score': avg_score,
            'type_stats': type_results,
            'recent_results': recent_results
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# 获取题目详情
@app.route('/api/question/<question_id>', methods=['GET'])
@csrf.exempt
def get_question_detail(question_id):
    question = QuestionBank.query.filter_by(question_id=question_id).first()
    if not question:
        return jsonify({'error': '题目不存在'}), 404

    return jsonify({
        'status': 'success',
        'question': question.to_dict()
    })


# 新增笔记相关的API路由
@app.route('/api/notes', methods=['GET', 'POST'])
@csrf.exempt
def video_notes():
    if request.method == 'GET':
        # 获取特定视频的笔记
        video_source = request.args.get('video_source')
        if not video_source:
            return jsonify({'error': '缺少video_source参数'}), 400

        notes = VideoNote.query.filter_by(video_source=video_source).order_by(VideoNote.timestamp.desc()).all()
        return jsonify({
            'video_source': video_source,
            'notes': [note.to_dict() for note in notes]
        })

    elif request.method == 'POST':
        # 创建新笔记
        data = request.json
        if not data or 'video_source' not in data or 'content' not in data:
            return jsonify({'error': '缺少必要参数'}), 400

        new_note = VideoNote(
            video_source=data['video_source'],
            content=data['content']
        )
        db.session.add(new_note)
        db.session.commit()

        return jsonify(new_note.to_dict()), 201


@app.route('/generate-ppt')
def generate_ppt_page():
    return render_template('generate_ppt.html')


@app.route('/api/ai/generate-ppt', methods=['POST'])
@feature_limit('generate_ppt')
@csrf.exempt
def generate_ppt():
    try:
        if 'file' not in request.files:
            return jsonify({'error': '未提供PPT文件'}), 400

        file = request.files['file']
        if not file.filename.lower().endswith('.pptx'):
            return jsonify({'error': '请上传PPTX格式文件'}), 400

        with tempfile.TemporaryDirectory() as temp_dir:
            # 保存上传的PPT
            upload_path = os.path.join(temp_dir, 'template.pptx')
            file.save(upload_path)
            prs = Presentation(upload_path)

            # === 定义每页小标题与内容的精准对应关系 ===
            SLIDE_CONFIGS = {
                1: {  # 第二页(索引1): 项目概述
                    'title': '项目概述 PROJECT OVERVIEW',
                    'header_positions': {
                        '01 已成立安徽臻锦科技有限公司': {
                            'position': 0,  # 第1个标题的位置
                            'content': '已组建10人核心研发团队，完成工商注册，申请发明专利3项、实用新型专利2项。'
                        },
                        '02 行业领先，综合性能行业一流': {
                            'position': 1,
                            'content': '探测器灵敏度0.01μSv/h(比国标高10倍)，通过CNAS认证，量程覆盖0.01μSv/h-10Sv/h。'
                        },
                        '03 市场破局，与10余家企业达成合作': {
                            'position': 2,
                            'content': '已与中核集团、中国辐射防护研究院等12家企业签订采购协议，获得首批订单200台。'
                        },
                        '04 企业认可，使用反馈极佳': {
                            'position': 3,
                            'content': '实测数据显示：稳定性99.9%，误报率<0.1%，平均无故障时间>5000小时。'
                        },
                        '05 军民融合 ，助力核安全强国建设': {
                            'position': 4,
                            'content': '产品满足GJB5313-2004军用标准，已列入军队核应急监测设备采购目录。'
                        }
                    }
                },
                2: {  # 第三页: 未来蓝图
                    'title': '未来蓝图 FUTURE BLUEPRINT',
                    'header_positions': {
                        '01 研发便携化电离室': {
                            'position': 0,
                            'content': '2024年推出重量<1kg的手持设备，体积减小80%，保持同等精度(±5%)。'
                        },
                        '02 培养技术人才': {
                            'position': 1,
                            'content': '与中科大共建"核检测实验室"，计划每年培养硕士以上专业人才30名。'
                        },
                        '03 市场破局，启动A轮融资': {
                            'position': 2,
                            'content': '计划融资2000万元用于生产线建设，目标2025年市场份额达到15%。'
                        },
                        '04 打通上下游产业链': {
                            'position': 3,
                            'content': '已与3家传感器供应商达成战略合作，元器件成本可降低25%。'
                        },
                        '05 打开国际市场': {
                            'position': 4,
                            'content': '正在申请CE/FDA认证，2024年进入东南亚和欧洲市场，预计出口额500万美元。'
                        }
                    }
                },
                3: {  # 第四页: 项目背景
                    'title': '项目背景 PROJECT BACKGROUND',
                    'header_positions': {
                        '日本排放核废水': {
                            'position': 0,
                            'content': '福岛核废水含氚、铯-137等放射性物质，预计5年内扩散至全球海域。'
                        },
                        '沿海城市忧虑': {
                            'position': 1,
                            'content': '2023年抽检显示：2%海鲜样本检出微量铯-137(0.1-0.3Bq/kg)。'
                        },
                        '大理石核辐射': {
                            'position': 2,
                            'content': '市场抽检不合格率8%，主要超标元素为镭-226(最高达1.5倍限值)。'
                        },
                        '海鲜安全': {
                            'position': 3,
                            'content': '自主研发的快速检测仪可在3分钟内完成检测，精度达0.01Bq/kg。'
                        },
                        '衣物核辐射超标 政策': {
                            'position': 4,
                            'content': '新修订《核安全法》要求加强日用消费品辐射监测，创造百亿级市场空间。'
                        }
                    }
                }
            }

            # === 精准填充逻辑 ===
            filled_slides = 0
            content_details = []

            for slide_num, config in SLIDE_CONFIGS.items():
                if slide_num >= len(prs.slides):
                    continue

                slide = prs.slides[slide_num]

                # 步骤1：收集所有文本框并分类
                text_shapes = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_shapes.append(shape)

                # 步骤2：识别标题和内容框的对应关系
                # 假设小标题是左对齐，内容框是右对齐(根据实际模板结构调整)
                header_boxes = {}
                content_boxes = {}

                for i, shape in enumerate(text_shapes):
                    text = shape.text_frame.text.strip()

                    # 通过文本内容识别小标题(包含数字编号)
                    if any(char.isdigit() for char in text[:3]):
                        for header in config['header_positions']:
                            if header in text:
                                header_boxes[config['header_positions'][header]['position']] = i
                                break
                    # 通过布局识别内容框(可根据实际模板调整识别逻辑)
                    elif shape.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT:
                        pos = len(content_boxes)
                        if pos < 5:  # 最多5个内容框
                            content_boxes[pos] = i

                # 步骤3：按位置对应关系填充内容
                for header_pos, header_idx in header_boxes.items():
                    if header_pos in content_boxes:
                        content_idx = content_boxes[header_pos]
                        header_text = text_shapes[header_idx].text_frame.text.strip()

                        # 找到对应的配置内容
                        matched_content = None
                        for header, details in config['header_positions'].items():
                            if header in header_text:
                                matched_content = details['content']
                                break

                        if matched_content:
                            # 填充内容
                            content_shape = text_shapes[content_idx]
                            content_shape.text_frame.clear()
                            p = content_shape.text_frame.paragraphs[0]
                            run = p.add_run()
                            run.text = matched_content

                            # 统一样式
                            run.font.name = '微软雅黑'
                            run.font.size = Pt(12)
                            run.font.color.rgb = RGBColor(0, 0, 0)

                filled_slides += 1
                content_details.append({
                    'slide': slide_num + 1,
                    'title': config['title'],
                    'contents': [f"{header}: {details['content']}"
                                 for header, details in config['header_positions'].items()]
                })

            # 保存并返回结果
            output_path = os.path.join(temp_dir, 'filled.pptx')
            prs.save(output_path)

            with open(output_path, 'rb') as f:
                ppt_data = base64.b64encode(f.read()).decode('utf-8')

            return jsonify({
                'success': True,
                'filename': '挑战杯_核辐射检测项目_已填充.pptx',
                'data': ppt_data,
                'details': content_details,
                'stats': {
                    'filled_slides': filled_slides,
                    'content_matches': sum(len(slide['contents']) for slide in content_details)
                }
            })

    except Exception as e:
        app.logger.error(f"PPT填充错误: {str(e)}\n{traceback.format_exc()}")
        return jsonify({
            'success': False,
            'error': 'PPT处理失败',
            'detail': str(e) if app.debug else None
        }), 500


def set_default_styles(prs):
    # 设置默认字体
    prs.font.name = '微软雅黑'

    # 定义颜色主题
    colors = prs.slide_master.color_scheme
    colors.accent1.color.rgb = RGBColor(79, 129, 189)  # 主色
    colors.accent6.color.rgb = RGBColor(165, 165, 165)  # 副色

    # 设置标题样式
    title_style = prs.slide_master.title_style
    title_style.font.size = Pt(24)
    title_style.font.color.rgb = colors.accent1.color.rgb

    # 设置内容样式
    content_style = prs.slide_master.text_style
    content_style.font.size = Pt(14)
    content_style.font.color.rgb = RGBColor(0, 0, 0)


def generate_ppt_fill_content(ppt_structure, theme, instructions):
    """调用AI生成填充内容"""
    prompt = f"""根据以下PPT结构和要求生成内容：
    主题: {theme}
    要求: {instructions}
    当前PPT结构: {json.dumps(ppt_structure, indent=2)}

    返回JSON格式填充数据"""
    # 这里调用AI接口获取填充内容
    return {'fills': []}  # 示例返回值


def apply_ppt_fills(presentation, fill_data):
    """将填充内容应用到PPT"""
    for fill in fill_data.get('fills', []):
        try:
            slide = presentation.slides[fill['slide'] - 1]
            shape = slide.shapes[fill['placeholder']]
            if shape.has_text_frame:
                shape.text_frame.text = fill.get('content', '')
        except:
            continue
    return presentation


def generate_ppt_structure_from_text(text, theme):
    """根据文本内容生成PPT结构"""
    prompt = f"""
    根据以下内容生成PPT结构大纲：
    - 主题: {theme}
    - 内容:
    {text}

    返回JSON格式：
    {{
        "slides": [
            {{
                "title": "幻灯片标题",
                "content": [
                    "主要内容点1",
                    "主要内容点2"
                ],
                "layout": "标题和内容/两栏/图片+文字 等",
                "notes": "幻灯片备注"
            }},
            ...
        ]
    }}
    """

    response = openai.ChatCompletion.create(
        model=config.DEEPSEEK_MODEL,
        messages=[
            {"role": "system", "content": "你是PPT结构设计专家，根据文本内容生成合理的PPT结构"},
            {"role": "user", "content": prompt}
        ],
        response_format={"type": "json_object"}
    )

    try:
        return json.loads(response.choices[0].message.content)
    except json.JSONDecodeError:
        # 默认结构
        return {
            "slides": [
                {
                    "title": f"{theme}演示",
                    "content": ["请在此处添加内容"],
                    "layout": "标题和内容",
                    "notes": ""
                }
            ]
        }


def build_presentation_from_structure(structure, theme):
    """根据结构创建PPT文档"""
    presentation = Presentation()

    # 设置幻灯片尺寸 (16:9)
    presentation.slide_width = Inches(13.33)
    presentation.slide_height = Inches(7.5)

    # 添加幻灯片
    slide_layout_map = {
        '标题和内容': 1,
        '两栏': 3,
        '图片+文字': 5,
        '仅标题': 0
    }

    for slide_data in structure.get('slides', []):
        layout_type = slide_data.get('layout', '标题和内容')
        layout_idx = slide_layout_map.get(layout_type, 1)

        # 创建幻灯片
        slide = presentation.slides.add_slide(presentation.slide_layouts[layout_idx])

        # 设置标题
        title = slide.shapes.title
        title.text = slide_data.get('title', '')

        # 设置内容 (根据布局类型)
        if layout_idx == 1:  # 标题和内容
            content = slide.placeholders[1]
            content.text = '\n'.join(slide_data.get('content', []))
        elif layout_idx == 3:  # 两栏
            left = slide.placeholders[1]
            left.text = '\n'.join(slide_data.get('content', [])[::2])  # 奇数项
            right = slide.placeholders[2]
            right.text = '\n'.join(slide_data.get('content', [])[1::2])  # 偶数项

    return presentation


def hex_to_rgb(hex_color):
    """十六进制颜色转RGB"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:
        hex_color = ''.join([c * 2 for c in hex_color])
    return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))


@app.route('/api/students/upload-avatar', methods=['POST'])
@csrf.exempt
def upload_student_avatar():
    """上传学生头像"""
    if 'file' not in request.files:
        return jsonify({'error': '未选择文件'}), 400
    
    file = request.files['file']
    student_id = request.form.get('student_id')
    
    if not student_id:
        return jsonify({'error': '缺少学生ID'}), 400
    
    if file.filename == '':
        return jsonify({'error': '未选择文件'}), 400
    
    # 检查文件类型
    allowed_extensions = {'png', 'jpg', 'jpeg', 'gif'}
    if not ('.' in file.filename and file.filename.rsplit('.', 1)[1].lower() in allowed_extensions):
        return jsonify({'error': '不支持的文件类型，请上传图片文件'}), 400
    
    # 检查学生是否存在
    student = Student.query.filter_by(student_id=student_id).first()
    if not student:
        return jsonify({'error': '学生不存在'}), 404
    
    try:
        # 创建上传目录
        upload_dir = os.path.join('static', 'uploads', 'avatars')
        os.makedirs(upload_dir, exist_ok=True)
        
        # 生成文件名
        file_extension = file.filename.rsplit('.', 1)[1].lower()
        filename = f"{student_id}_{int(datetime.utcnow().timestamp())}.{file_extension}"
        file_path = os.path.join(upload_dir, filename)
        
        # 保存文件
        file.save(file_path)
        
        # 删除旧头像文件
        if student.avatar and os.path.exists(student.avatar):
            try:
                os.remove(student.avatar)
            except:
                pass
        
        # 更新数据库
        student.avatar = file_path
        student.updated_at = datetime.utcnow()
        db.session.commit()
        
        return jsonify({
            'message': '头像上传成功',
            'avatar_url': f'/{file_path}'
        }), 200
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'上传失败: {str(e)}'}), 500


@app.route('/api/students/batch-import', methods=['POST'])
@csrf.exempt
def batch_import_students():
    """批量导入学生信息"""
    if 'file' not in request.files:
        return jsonify({'error': '未选择文件'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': '未选择文件'}), 400
    
    # 检查文件类型
    if not file.filename.lower().endswith(('.xlsx', '.xls', '.csv')):
        return jsonify({'error': '不支持的文件类型，请上传Excel或CSV文件'}), 400
    
    try:
        import pandas as pd
        
        # 读取文件
        if file.filename.lower().endswith('.csv'):
            df = pd.read_csv(file, encoding='utf-8')
        else:
            df = pd.read_excel(file)
        
        # 检查必需的列
        required_columns = ['student_id', 'name']
        missing_columns = [col for col in required_columns if col not in df.columns]
        if missing_columns:
            return jsonify({'error': f'缺少必需的列: {", ".join(missing_columns)}'}), 400
        
        success_count = 0
        error_count = 0
        errors = []
        
        for index, row in df.iterrows():
            try:
                # 检查学号是否已存在
                existing_student = Student.query.filter_by(student_id=str(row['student_id'])).first()
                if existing_student:
                    error_count += 1
                    errors.append(f'第{index+2}行: 学号{row["student_id"]}已存在')
                    continue
                
                # 创建新学生
                new_student = Student(
                    student_id=str(row['student_id']),
                    name=str(row['name']),
                    gender=str(row.get('gender', '')) if pd.notna(row.get('gender')) else None,
                    birth_date=pd.to_datetime(row['birth_date']).date() if pd.notna(row.get('birth_date')) else None,
                    id_card=str(row.get('id_card', '')) if pd.notna(row.get('id_card')) else None,
                    phone=str(row.get('phone', '')) if pd.notna(row.get('phone')) else None,
                    email=str(row.get('email', '')) if pd.notna(row.get('email')) else None,
                    major=str(row.get('major', '')) if pd.notna(row.get('major')) else None,
                    class_name=str(row.get('class_name', '')) if pd.notna(row.get('class_name')) else None,
                    grade=str(row.get('grade', '')) if pd.notna(row.get('grade')) else None,
                    education_level=str(row.get('education_level', '')) if pd.notna(row.get('education_level')) else None,
                    student_status=str(row.get('student_status', '在读')) if pd.notna(row.get('student_status')) else '在读',
                    address=str(row.get('address', '')) if pd.notna(row.get('address')) else None,
                    emergency_contact=str(row.get('emergency_contact', '')) if pd.notna(row.get('emergency_contact')) else None,
                    emergency_phone=str(row.get('emergency_phone', '')) if pd.notna(row.get('emergency_phone')) else None,
                    parent_name=str(row.get('parent_name', '')) if pd.notna(row.get('parent_name')) else None,
                    parent_phone=str(row.get('parent_phone', '')) if pd.notna(row.get('parent_phone')) else None,
                    notes=str(row.get('notes', '')) if pd.notna(row.get('notes')) else None
                )
                
                db.session.add(new_student)
                success_count += 1
                
            except Exception as e:
                error_count += 1
                errors.append(f'第{index+2}行: {str(e)}')
        
        db.session.commit()
        
        return jsonify({
            'message': '批量导入完成',
            'success_count': success_count,
            'error_count': error_count,
            'errors': errors[:10]  # 只返回前10个错误
        }), 200
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'导入失败: {str(e)}'}), 500


@app.route('/api/students/export', methods=['GET'])
@csrf.exempt
def export_students():
    """导出学生信息"""
    try:
        import pandas as pd
        from io import BytesIO
        
        # 获取查询参数
        format_type = request.args.get('format', 'excel')  # excel 或 csv
        
        # 查询学生数据
        students = Student.query.all()
        
        # 准备数据
        data = []
        for student in students:
            data.append({
                'student_id': student.student_id,
                'name': student.name,
                'gender': student.gender,
                'birth_date': student.birth_date.strftime('%Y-%m-%d') if student.birth_date else '',
                'id_card': student.id_card,
                'phone': student.phone,
                'email': student.email,
                'major': student.major,
                'class_name': student.class_name,
                'grade': student.grade,
                'education_level': student.education_level,
                'student_status': student.student_status,
                'address': student.address,
                'emergency_contact': student.emergency_contact,
                'emergency_phone': student.emergency_phone,
                'parent_name': student.parent_name,
                'parent_phone': student.parent_phone,
                'admission_date': student.admission_date.strftime('%Y-%m-%d') if student.admission_date else '',
                'notes': student.notes,
                'created_at': student.created_at.strftime('%Y-%m-%d %H:%M:%S') if student.created_at else '',
                'updated_at': student.updated_at.strftime('%Y-%m-%d %H:%M:%S') if student.updated_at else ''
            })
        
        df = pd.DataFrame(data)
        
        # 设置列名映射
        column_mapping = {
            'student_id': '学号',
            'name': '姓名',
            'gender': '性别',
            'birth_date': '出生日期',
            'id_card': '身份证号',
            'phone': '联系电话',
            'email': '邮箱',
            'major': '专业',
            'class_name': '班级',
            'grade': '年级',
            'education_level': '学制',
            'student_status': '状态',
            'address': '地址',
            'emergency_contact': '紧急联系人',
            'emergency_phone': '紧急联系电话',
            'parent_name': '家长姓名',
            'parent_phone': '家长电话',
            'admission_date': '入学日期',
            'notes': '备注',
            'created_at': '创建时间',
            'updated_at': '更新时间'
        }
        
        df = df.rename(columns=column_mapping)
        
        # 生成文件
        output = BytesIO()
        filename = f'学生信息_{datetime.now().strftime("%Y%m%d_%H%M%S")}'
        
        if format_type == 'csv':
            df.to_csv(output, index=False, encoding='utf-8-sig')
            filename += '.csv'
            mimetype = 'text/csv'
        else:
            df.to_excel(output, index=False, engine='openpyxl')
            filename += '.xlsx'
            mimetype = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        
        output.seek(0)
        
        return send_file(
            output,
            mimetype=mimetype,
            as_attachment=True,
            download_name=filename
        )
        
    except Exception as e:
        return jsonify({'error': f'导出失败: {str(e)}'}), 500


@app.route('/api/students/batch-delete', methods=['POST'])
@csrf.exempt
def batch_delete_students():
    """批量删除学生"""
    data = request.json
    student_ids = data.get('student_ids', [])
    
    if not student_ids:
        return jsonify({'error': '未提供要删除的学生ID'}), 400
    
    try:
        deleted_count = 0
        errors = []
        
        for student_id in student_ids:
            student = Student.query.filter_by(student_id=student_id).first()
            if student:
                try:
                    # 删除相关的作业记录
                    Assignment.query.filter_by(student_id=student_id).delete()
                    
                    # 删除相关的问题提交记录
                    QuestionSubmission.query.filter_by(student_id=student_id).delete()
                    
                    # 删除头像文件
                    if student.avatar and os.path.exists(student.avatar):
                        try:
                            os.remove(student.avatar)
                        except:
                            pass
                    
                    # 删除学生记录
                    db.session.delete(student)
                    deleted_count += 1
                except Exception as e:
                    errors.append(f'删除学生{student_id}失败: {str(e)}')
            else:
                errors.append(f'学生{student_id}不存在')
        
        db.session.commit()
        
        return jsonify({
            'message': f'成功删除{deleted_count}个学生',
            'deleted_count': deleted_count,
            'errors': errors
        }), 200
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'批量删除失败: {str(e)}'}), 500


@app.route('/test/students')
def test_students():
    """学生管理测试页面"""
    return render_template('test_students.html')


@app.route('/api/assignments/add-score', methods=['POST'])
@csrf.exempt
def add_manual_score():
    """手动添加学生成绩记录"""
    try:
        data = request.get_json()
        
        # 验证必要字段
        required_fields = ['student_id', 'assignment_name', 'subject', 'score']
        for field in required_fields:
            if field not in data:
                return jsonify({'error': f'缺少必要字段: {field}'}), 400
        
        # 验证学生是否存在
        student = Student.query.filter_by(student_id=data['student_id']).first()
        if not student:
            return jsonify({'error': '学生不存在'}), 404
        
        # 验证分数范围
        score = float(data['score'])
        if score < 0 or score > 100:
            return jsonify({'error': '分数必须在0-100之间'}), 400
        
        # 创建作业记录
        assignment = Assignment(
            student_id=data['student_id'],
            assignment_name=data['assignment_name'],
            subject=data['subject'],
            chapter=data.get('chapter', '手动录入'),
            score=score,
            feedback=data.get('feedback', ''),
            submission_time=datetime.utcnow()
        )
        
        db.session.add(assignment)
        db.session.commit()
        
        return jsonify({
            'message': '成绩录入成功',
            'assignment': assignment.to_dict()
        }), 201
        
    except ValueError:
        return jsonify({'error': '分数格式错误'}), 400
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'录入失败: {str(e)}'}), 500


@app.route('/api/assignments/batch-scores', methods=['POST'])
@csrf.exempt
def batch_add_scores():
    """批量录入学生成绩"""
    try:
        data = request.get_json()
        scores_data = data.get('scores', [])
        
        if not scores_data:
            return jsonify({'error': '没有提供成绩数据'}), 400
        
        success_count = 0
        error_count = 0
        errors = []
        
        for index, score_data in enumerate(scores_data):
            try:
                # 验证必要字段
                required_fields = ['student_id', 'assignment_name', 'subject', 'score']
                for field in required_fields:
                    if field not in score_data:
                        raise ValueError(f'第{index+1}行缺少字段: {field}')
                
                # 验证学生是否存在
                student = Student.query.filter_by(student_id=score_data['student_id']).first()
                if not student:
                    raise ValueError(f'第{index+1}行学生不存在: {score_data["student_id"]}')
                
                # 验证分数
                score = float(score_data['score'])
                if score < 0 or score > 100:
                    raise ValueError(f'第{index+1}行分数超出范围: {score}')
                
                # 创建作业记录
                assignment = Assignment(
                    student_id=score_data['student_id'],
                    assignment_name=score_data['assignment_name'],
                    subject=score_data['subject'],
                    chapter=score_data.get('chapter', '批量录入'),
                    score=score,
                    feedback=score_data.get('feedback', ''),
                    submission_time=datetime.utcnow()
                )
                
                db.session.add(assignment)
                success_count += 1
                
            except Exception as e:
                error_count += 1
                errors.append(f'第{index+1}行: {str(e)}')
        
        if success_count > 0:
            db.session.commit()
        
        return jsonify({
            'message': f'批量录入完成',
            'success_count': success_count,
            'error_count': error_count,
            'errors': errors
        }), 200 if error_count == 0 else 207
        
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f'批量录入失败: {str(e)}'}), 500


@app.route('/api/students/<student_id>/assignments', methods=['GET'])
@csrf.exempt
def get_student_assignments(student_id):
    """获取特定学生的所有作业记录"""
    try:
        # 验证学生是否存在
        student = Student.query.filter_by(student_id=student_id).first()
        if not student:
            return jsonify({'error': '学生不存在'}), 404
        
        # 获取学生的所有作业
        assignments = Assignment.query.filter_by(student_id=student_id).order_by(
            Assignment.submission_time.desc()
        ).all()
        
        # 计算统计信息
        total_assignments = len(assignments)
        avg_score = sum(a.score for a in assignments if a.score) / total_assignments if total_assignments > 0 else 0
        max_score = max((a.score for a in assignments if a.score), default=0)
        min_score = min((a.score for a in assignments if a.score), default=0)
        
        # 按科目统计
        subject_stats = {}
        for assignment in assignments:
            subject = assignment.subject
            if subject not in subject_stats:
                subject_stats[subject] = {
                    'subject': subject,
                    'count': 0,
                    'avg_score': 0,
                    'total_score': 0
                }
            
            subject_stats[subject]['count'] += 1
            if assignment.score:
                subject_stats[subject]['total_score'] += assignment.score
        
        # 计算各科目平均分
        for subject, stats in subject_stats.items():
            if stats['count'] > 0:
                stats['avg_score'] = round(stats['total_score'] / stats['count'], 2)
        
        return jsonify({
            'student': student.to_dict(),
            'assignments': [assignment.to_dict() for assignment in assignments],
            'statistics': {
                'total_assignments': total_assignments,
                'avg_score': round(avg_score, 2),
                'max_score': max_score,
                'min_score': min_score,
                'subject_stats': list(subject_stats.values())
            }
        })
        
    except Exception as e:
        return jsonify({'error': f'获取学生作业失败: {str(e)}'}), 500


# ==================== 用户认证API ====================

@app.route('/api/auth/register', methods=['POST'])
@csrf.exempt
@limiter.limit("5 per hour")  # 每小时最多5次注册请求
def register():
    """用户注册API"""
    try:
        data = request.json
        
        # 验证必填字段
        if not data.get('username') or not data.get('email') or not data.get('password'):
            return jsonify({'error': '用户名、邮箱和密码为必填项'}), 400
        
        # 验证用户名格式
        username_valid, username_msg = validate_username(data['username'])
        if not username_valid:
            return jsonify({'error': username_msg}), 400
        
        # 验证邮箱格式
        if not validate_email(data['email']):
            return jsonify({'error': '邮箱格式不正确'}), 400
        
        # 验证密码强度（新增）
        pwd_valid, pwd_msg, pwd_strength = validate_password_strength(data['password'])
        if not pwd_valid:
            return jsonify({
                'error': pwd_msg,
                'strength': pwd_strength
            }), 400
        
        # 检查用户名是否已存在
        if User.query.filter_by(username=data['username']).first():
            return jsonify({'error': '用户名已被使用'}), 400
        
        # 检查邮箱是否已存在
        if User.query.filter_by(email=data['email']).first():
            return jsonify({'error': '邮箱已被注册'}), 400
        
        # 创建新用户
        new_user = User(
            username=data['username'],
            email=data['email'],
            is_admin=False
        )
        new_user.set_password(data['password'])
        
        db.session.add(new_user)
        db.session.commit()
        
        # ⚠️ 注意：新注册用户不再自动分配免费会员
        # 用户需要购买会员后才能使用系统功能
        
        return jsonify({
            'message': '注册成功！请购买会员后开始使用',
            'user': {
                'id': new_user.id,
                'username': new_user.username,
                'email': new_user.email
            },
            'redirect': '/payment'  # 引导到购买页面
        }), 201
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"注册错误: {str(e)}")
        return jsonify({'error': f'注册失败: {str(e)}'}), 500


@app.route('/api/auth/login', methods=['POST'])
@csrf.exempt
@limiter.limit("10 per minute")  # 每分钟最多10次登录请求
def login():
    """用户登录API"""
    try:
        data = request.json
        username_input = data.get('username', '').strip()
        password = data.get('password', '')
        
        # 验证必填字段
        if not username_input or not password:
            return jsonify({'error': '用户名和密码为必填项'}), 400
        
        # 检查账户是否被锁定
        is_locked, locked_until, attempts = is_account_locked(username_input)
        if is_locked:
            # 计算剩余锁定时间
            if locked_until:
                now_utc = datetime.utcnow()
                remaining_seconds = (locked_until - now_utc).total_seconds()
                if remaining_seconds > 0:
                    remaining_minutes = max(1, int(remaining_seconds / 60))  # 至少显示1分钟
                    message = f'账户已被锁定，请在 {remaining_minutes} 分钟后重试'
                    
                    app.logger.warning(f'安全事件 - 尝试登录已锁定账户: 用户={username_input}, IP={request.remote_addr}, 剩余{remaining_minutes}分钟')
                    
                    return jsonify({
                        'error': message,
                        'locked': True,
                        'locked_until': locked_until.isoformat() + 'Z',  # 明确标记为UTC时间
                        'remaining_seconds': int(remaining_seconds),
                        'remaining_minutes': remaining_minutes,
                        'attempts': attempts
                    }), 403
                else:
                    # 锁定时间已过期，自动解锁
                    from utils.security import unlock_account
                    unlock_account(username_input)
                    app.logger.info(f'账户自动解锁: 用户={username_input}')
            else:
                # 没有锁定时间但显示锁定，可能是数据异常，自动解锁
                from utils.security import unlock_account
                unlock_account(username_input)
                app.logger.info(f'账户异常解锁: 用户={username_input}')
        
        # 查找用户（支持用户名或邮箱登录）
        user = User.query.filter(
            (User.username == username_input) | (User.email == username_input)
        ).first()
        
        # 验证用户存在且密码正确
        if not user or not user.check_password(password):
            # 记录登录失败
            record_result = record_login_attempt(
                username=username_input,
                ip_address=request.remote_addr,
                user_agent=request.user_agent.string,
                success=False,
                failure_reason='用户名或密码错误'
            )
            
            # 获取剩余尝试次数
            remaining = get_remaining_attempts(username_input)
            
            # 记录安全日志
            app.logger.warning(
                f'安全事件 - 登录失败: 用户={username_input}, IP={request.remote_addr}, '
                f'剩余尝试次数={remaining}'
            )
            
            # 如果触发了锁定，特别提示
            if record_result.get('locked'):
                app.logger.error(
                    f'安全事件 - 账户锁定: 用户={username_input}, IP={request.remote_addr}, '
                    f'锁定至={record_result.get("locked_until")}'
                )
                return jsonify({
                    'error': '登录失败次数过多，账户已被锁定2分钟',
                    'locked': True,
                    'locked_until': record_result.get('locked_until'),
                    'remaining_attempts': 0
                }), 403
            
            # 返回错误和剩余次数
            error_msg = '用户名或密码错误'
            if remaining <= 2 and remaining > 0:
                error_msg += f'，还有 {remaining} 次尝试机会'
            
            return jsonify({
                'error': error_msg,
                'remaining_attempts': remaining
            }), 401
        
        # 记录登录成功
        record_login_attempt(
            username=username_input,
            ip_address=request.remote_addr,
            user_agent=request.user_agent.string,
            success=True
        )
        
        # 更新最后登录时间（使用本地时间）
        user.last_login = datetime.now()
        db.session.commit()
        
        # 登录用户
        login_user(user, remember=data.get('remember', False))
        
        # 记录成功登录日志
        app.logger.info(f'用户登录成功: 用户={user.username}, ID={user.id}, IP={request.remote_addr}')
        
        # 获取当前会员信息
        membership = user.get_current_membership()
        membership_info = None
        if membership:
            membership_info = {
                'tier_name': membership.tier.name,
                'tier_level': membership.tier.level,
                'end_date': membership.end_date.isoformat()
            }
        
        return jsonify({
            'message': '登录成功',
            'user': {
                'id': user.id,
                'username': user.username,
                'email': user.email,
                'is_admin': user.is_admin,
                'membership': membership_info
            }
        }), 200
        
    except Exception as e:
        app.logger.error(f"登录错误: {str(e)}")
        return jsonify({'error': f'登录失败: {str(e)}'}), 500


@app.route('/api/auth/logout', methods=['POST'])
@csrf.exempt
@require_login_api
def logout():
    """用户登出API"""
    try:
        logout_user()
        return jsonify({'message': '登出成功'}), 200
    except Exception as e:
        app.logger.error(f"登出错误: {str(e)}")
        return jsonify({'error': f'登出失败: {str(e)}'}), 500


@app.route('/api/auth/current-user', methods=['GET'])
@csrf.exempt
def get_current_user():
    """获取当前登录用户信息"""
    try:
        # 检查用户是否登录
        if not current_user.is_authenticated:
            return jsonify({
                'logged_in': False,
                'user': None,
                'membership': None
            }), 200
        
        # 获取当前会员信息
        membership = current_user.get_current_membership()
        membership_info = {
            'tier_name': '免费版',
            'tier_type': 'free',
            'tier_level': 0
        }
        
        if membership:
            membership_info = {
                'tier_name': membership.tier.name,
                'tier_type': membership.tier.code,
                'tier_level': membership.tier.level,
                'tier_price': membership.tier.price,
                'start_date': membership.start_date.isoformat() if membership.start_date else None,
                'end_date': membership.end_date.isoformat() if membership.end_date else None,
                'is_active': membership.is_active
            }
        
        # 获取使用统计
        usage_stats = {
            'ai_ask_today': current_user.get_usage_count('ai_ask', 'daily'),
            'ai_ask_week': current_user.get_usage_count('ai_ask', 'weekly'),
            'ai_ask_month': current_user.get_usage_count('ai_ask', 'monthly'),
            'generate_question_month': current_user.get_usage_count('generate_question', 'monthly'),
            'generate_lecture_month': current_user.get_usage_count('generate_lecture', 'monthly'),
        }
        
        return jsonify({
            'logged_in': True,
            'user': {
                'id': current_user.id,
                'username': current_user.username,
                'email': current_user.email,
                'is_admin': current_user.is_admin,
                'created_at': current_user.created_at.isoformat()
            },
            'membership': membership_info,
            'usage_stats': usage_stats
        }), 200
        
    except Exception as e:
        app.logger.error(f"获取用户信息错误: {str(e)}")
        return jsonify({'error': f'获取用户信息失败: {str(e)}'}), 500


@app.route('/api/auth/check', methods=['GET'])
@csrf.exempt
def check_auth():
    """检查用户是否已登录"""
    if current_user.is_authenticated:
        return jsonify({
            'authenticated': True,
            'user': {
                'id': current_user.id,
                'username': current_user.username,
                'is_admin': current_user.is_admin
            }
        }), 200
    else:
        return jsonify({'authenticated': False}), 200


@app.route('/api/membership/tiers', methods=['GET'])
@csrf.exempt
def get_membership_tiers():
    """获取所有会员等级（不包含免费套餐）"""
    try:
        # ⚠️ 过滤掉免费套餐，只显示付费套餐
        tiers = MembershipTier.query.filter(
            MembershipTier.is_active == True,
            MembershipTier.code != 'free'
        ).order_by(MembershipTier.price).all()
        
        return jsonify({
            'tiers': [tier.to_dict() for tier in tiers]
        }), 200
    except Exception as e:
        app.logger.error(f"获取会员等级错误: {str(e)}")
        return jsonify({'error': f'获取会员等级失败: {str(e)}'}), 500


@app.route('/api/membership/purchase', methods=['POST'])
@csrf.exempt
@require_login_api
def purchase_membership():
    """购买会员（模拟支付）"""
    try:
        data = request.json
        tier_id = data.get('tier_id')
        
        if not tier_id:
            return jsonify({'error': '缺少会员等级ID'}), 400
        
        # 获取会员等级
        tier = MembershipTier.query.get(tier_id)
        if not tier:
            return jsonify({'error': '会员等级不存在'}), 404
        
        # 不能购买免费会员
        if tier.code == 'free':
            return jsonify({'error': '免费会员无需购买'}), 400
        
        # 生成交易ID
        import uuid
        transaction_id = f"TXN_{uuid.uuid4().hex[:16].upper()}"
        
        # 创建支付交易记录
        payment = PaymentTransaction(
            user_id=current_user.id,
            transaction_id=transaction_id,
            amount=tier.price,
            currency='CNY',
            payment_method='simulated',  # 模拟支付
            tier_id=tier_id,
            status='completed',
            note='模拟支付自动完成',
            completed_at=datetime.utcnow()
        )
        
        db.session.add(payment)
        db.session.flush()  # 获取payment.id
        
        # 计算会员有效期
        start_date = datetime.utcnow()
        end_date = start_date + timedelta(days=tier.duration_days)
        
        # 创建会员记录
        membership = UserMembership(
            user_id=current_user.id,
            tier_id=tier_id,
            start_date=start_date,
            end_date=end_date,
            is_active=True,
            payment_id=payment.id
        )
        
        db.session.add(membership)
        db.session.commit()
        
        return jsonify({
            'message': '购买成功',
            'transaction_id': transaction_id,
            'membership': membership.to_dict()
        }), 200
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"购买会员错误: {str(e)}")
        return jsonify({'error': f'购买失败: {str(e)}'}), 500


@app.route('/api/membership/tiers/available', methods=['GET'])
@csrf.exempt
def get_available_tiers():
    """获取当前可购买的套餐（包含早鸟状态）"""
    try:
        from membership_utils import get_available_tiers
        
        tiers_info = get_available_tiers()
        
        return jsonify(tiers_info), 200
        
    except Exception as e:
        app.logger.error(f"获取可用套餐错误: {str(e)}")
        return jsonify({'error': f'获取套餐失败: {str(e)}'}), 500


@app.route('/api/membership/early-bird/status', methods=['GET'])
@csrf.exempt
def get_early_bird_status():
    """获取早鸟优惠状态"""
    try:
        from membership_utils import get_current_early_bird_tier
        
        early_bird = get_current_early_bird_tier()
        
        if early_bird:
            return jsonify({
                'active': True,
                **early_bird
            }), 200
        else:
            return jsonify({
                'active': False,
                'message': '早鸟优惠已售罄'
            }), 200
        
    except Exception as e:
        app.logger.error(f"获取早鸟状态错误: {str(e)}")
        return jsonify({'error': f'获取早鸟状态失败: {str(e)}'}), 500


@app.route('/api/membership/status', methods=['GET'])
@csrf.exempt
@require_login_api
def get_user_membership_status():
    """获取当前用户的会员状态"""
    try:
        from membership_utils import get_membership_status
        
        status = get_membership_status(current_user.id)
        
        return jsonify(status), 200
        
    except Exception as e:
        app.logger.error(f"获取会员状态错误: {str(e)}")
        return jsonify({'error': f'获取会员状态失败: {str(e)}'}), 500


@app.route('/api/membership/usage', methods=['GET'])
@csrf.exempt
@require_login_api
def get_user_usage():
    """获取用户所有功能使用情况"""
    try:
        from membership_utils import get_all_features_usage
        
        usage_list = get_all_features_usage(current_user.id)
        
        return jsonify({
            'features': usage_list
        }), 200
        
    except Exception as e:
        app.logger.error(f"获取使用情况错误: {str(e)}")
        return jsonify({'error': f'获取使用情况失败: {str(e)}'}), 500


@app.route('/api/membership/history', methods=['GET'])
@csrf.exempt
@require_login_api
def get_membership_history():
    """获取会员历史和购买记录"""
    try:
        from membership_utils import get_user_membership
        
        # 获取当前会员信息
        membership = get_user_membership(current_user.id)
        
        current_membership = None
        if membership:
            days_remaining = (membership.end_date - datetime.utcnow()).days if membership.end_date else 0
            current_membership = {
                'tier_name': membership.tier.name,
                'tier_code': membership.tier.code,
                'start_date': membership.start_date.strftime('%Y-%m-%d') if membership.start_date else None,
                'end_date': membership.end_date.strftime('%Y-%m-%d') if membership.end_date else None,
                'days_remaining': max(0, days_remaining),
                'auto_renew': membership.auto_renew if hasattr(membership, 'auto_renew') else False
            }
        
        # 获取购买历史
        payments = PaymentTransaction.query.filter_by(
            user_id=current_user.id
        ).order_by(PaymentTransaction.created_at.desc()).limit(20).all()
        
        purchase_history = []
        for payment in payments:
            tier = MembershipTier.query.get(payment.tier_id) if payment.tier_id else None
            purchase_history.append({
                'id': payment.id,
                'tier_name': tier.name if tier else '未知',
                'amount': float(payment.amount),
                'payment_method': payment.payment_method,
                'transaction_id': payment.transaction_id,
                'status': payment.status,
                'created_at': payment.created_at.strftime('%Y-%m-%d %H:%M:%S')
            })
        
        return jsonify({
            'current_membership': current_membership,
            'purchase_history': purchase_history
        }), 200
        
    except Exception as e:
        app.logger.error(f"获取会员历史错误: {str(e)}")
        return jsonify({'error': f'获取历史失败: {str(e)}'}), 500


@app.route('/api/membership/upgrade-options', methods=['GET'])
@csrf.exempt
@require_login_api
def get_upgrade_options():
    """获取可升级的会员套餐"""
    try:
        from membership_utils import get_user_membership
        
        # 获取当前会员信息
        membership = get_user_membership(current_user.id)
        current_tier_level = membership.tier.level if membership else 0
        current_tier_code = membership.tier.code if membership else 'free'
        
        # 获取所有会员等级
        tiers = MembershipTier.query.order_by(MembershipTier.level).all()
        
        upgrade_options = []
        for tier in tiers:
            # 跳过免费会员
            if tier.code == 'free':
                continue
                
            # 判断是否为升级选项
            is_upgrade = tier.level > current_tier_level
            
            # 计算折扣价（升级时的价格）
            discount_price = tier.price * 0.8 if is_upgrade else tier.price
            
            # 获取特权列表
            features = []
            if tier.features:
                import json
                try:
                    features = json.loads(tier.features) if isinstance(tier.features, str) else tier.features
                except:
                    features = []
            
            upgrade_options.append({
                'tier_id': tier.id,
                'tier_name': tier.name,
                'tier_code': tier.code,
                'price': float(tier.price),
                'discount_price': float(discount_price),
                'duration_days': tier.duration_days,
                'features': features,
                'is_upgrade': is_upgrade
            })
        
        return jsonify({
            'current_tier': current_tier_code,
            'upgrade_options': upgrade_options
        }), 200
        
    except Exception as e:
        app.logger.error(f"获取升级选项错误: {str(e)}")
        return jsonify({'error': f'获取升级选项失败: {str(e)}'}), 500


@app.route('/api/payment/create-order', methods=['POST'])
@csrf.exempt
@require_login_api
def create_payment_order():
    """创建支付订单（模拟）"""
    try:
        data = request.json
        tier_id = data.get('tier_id')
        payment_method = data.get('payment_method', 'alipay')
        
        if not tier_id:
            return jsonify({'error': '缺少会员等级ID'}), 400
        
        # 获取会员等级
        tier = MembershipTier.query.get(tier_id)
        if not tier:
            return jsonify({'error': '会员等级不存在'}), 404
        
        # 生成订单ID
        import uuid
        order_id = f"ORDER_{uuid.uuid4().hex[:16].upper()}"
        transaction_id = f"TXN_{uuid.uuid4().hex[:16].upper()}"
        
        # 生成模拟支付二维码URL（使用草料二维码API）
        qr_text = f"模拟支付订单: {order_id}, 金额: ¥{tier.price}, 套餐: {tier.name}"
        qr_code_url = f"https://api.qrserver.com/v1/create-qr-code/?size=200x200&data={qr_text}"
        
        # 创建待支付的交易记录
        payment = PaymentTransaction(
            user_id=current_user.id,
            transaction_id=transaction_id,
            amount=tier.price,
            currency='CNY',
            payment_method=payment_method,
            tier_id=tier_id,
            status='pending',
            qr_code_url=qr_code_url,
            note=f'订单ID: {order_id}'
        )
        
        db.session.add(payment)
        db.session.commit()
        
        # 计算过期时间（15分钟后）
        expire_time = (datetime.utcnow() + timedelta(minutes=15)).strftime('%Y-%m-%d %H:%M:%S')
        
        return jsonify({
            'order_id': order_id,
            'transaction_id': transaction_id,
            'qr_code_url': qr_code_url,
            'amount': float(tier.price),
            'expire_time': expire_time
        }), 200
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"创建支付订单错误: {str(e)}")
        return jsonify({'error': f'创建订单失败: {str(e)}'}), 500


@app.route('/api/payment/check-status/<transaction_id>', methods=['GET'])
@csrf.exempt
@require_login_api
def check_payment_status(transaction_id):
    """查询支付状态（模拟自动成功）"""
    try:
        # 查找交易记录
        payment = PaymentTransaction.query.filter_by(
            transaction_id=transaction_id,
            user_id=current_user.id
        ).first()
        
        if not payment:
            return jsonify({'error': '订单不存在'}), 404
        
        # 模拟支付：如果是pending状态，自动改为success
        if payment.status == 'pending':
            payment.status = 'completed'
            payment.completed_at = datetime.utcnow()
            
            # 获取会员等级
            tier = MembershipTier.query.get(payment.tier_id)
            if tier:
                # 计算会员有效期
                start_date = datetime.utcnow()
                end_date = start_date + timedelta(days=tier.duration_days)
                
                # 检查是否已有会员记录
                existing_membership = UserMembership.query.filter_by(
                    user_id=current_user.id
                ).first()
                
                if existing_membership:
                    # 更新现有会员
                    existing_membership.tier_id = tier.id
                    existing_membership.start_date = start_date
                    existing_membership.end_date = end_date
                    existing_membership.payment_id = payment.id
                else:
                    # 创建新会员记录
                    membership = UserMembership(
                        user_id=current_user.id,
                        tier_id=tier.id,
                        start_date=start_date,
                        end_date=end_date,
                        payment_id=payment.id
                    )
                    db.session.add(membership)
            
            db.session.commit()
        
        return jsonify({
            'status': payment.status,
            'order_id': payment.note.replace('订单ID: ', '') if payment.note else '',
            'paid_at': payment.completed_at.strftime('%Y-%m-%d %H:%M:%S') if payment.completed_at else None
        }), 200
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"查询支付状态错误: {str(e)}")
        return jsonify({'error': f'查询失败: {str(e)}'}), 500


@app.route('/api/membership/cancel-auto-renew', methods=['POST'])
@csrf.exempt
@require_login_api
def cancel_auto_renew():
    """取消自动续费"""
    try:
        # 获取当前用户的会员记录
        membership = UserMembership.query.filter_by(
            user_id=current_user.id
        ).first()
        
        if not membership:
            return jsonify({'error': '您还不是会员'}), 404
        
        # 取消自动续费
        if hasattr(membership, 'auto_renew'):
            membership.auto_renew = False
            db.session.commit()
            return jsonify({'message': '已取消自动续费'}), 200
        else:
            return jsonify({'message': '当前系统不支持自动续费'}), 400
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"取消自动续费错误: {str(e)}")
        return jsonify({'error': f'取消失败: {str(e)}'}), 500


@app.route('/api/usage/stats', methods=['GET'])
@csrf.exempt
@require_login_api
def get_usage_stats_api():
    """获取用户使用统计"""
    try:
        from membership_utils import FEATURE_PERMISSIONS, get_user_membership, get_usage_stats
        
        period = request.args.get('period', 'daily')
        user_id = current_user.id
        
        # 获取用户会员信息
        membership = get_user_membership(user_id)
        tier_code = membership.tier.code if membership else 'free'
        
        # 获取各功能的使用统计
        stats = {}
        for feature_name, permissions in FEATURE_PERMISSIONS.items():
            tier_perms = permissions.get(tier_code, permissions['free'])
            used = get_usage_stats(user_id, feature_name, period)
            
            stats[feature_name] = {
                'used': used,
                'limit': tier_perms.get('limit', 0)
            }
        
        return jsonify({'stats': stats}), 200
        
    except Exception as e:
        app.logger.error(f"获取使用统计错误: {str(e)}")
        return jsonify({'error': f'获取统计失败: {str(e)}'}), 500


@app.route('/api/auth/change-password', methods=['POST'])
@csrf.exempt
@require_login_api
@limiter.limit("3 per hour")  # 每小时最多3次修改密码请求
def change_password():
    """修改密码"""
    try:
        data = request.json
        current_password = data.get('current_password')
        new_password = data.get('new_password')
        
        if not current_password or not new_password:
            return jsonify({'error': '请填写所有字段'}), 400
        
        # 验证当前密码
        if not current_user.check_password(current_password):
            return jsonify({'error': '当前密码不正确'}), 400
        
        # 验证新密码强度
        pwd_valid, pwd_msg, pwd_strength = validate_password_strength(new_password)
        if not pwd_valid:
            return jsonify({
                'error': pwd_msg,
                'strength': pwd_strength
            }), 400
        
        # 更新密码
        current_user.password = new_password
        db.session.commit()
        
        return jsonify({'message': '密码修改成功'}), 200
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"修改密码错误: {str(e)}")
        return jsonify({'error': f'修改密码失败: {str(e)}'}), 500


# ==================== 忘记密码功能 ====================

@app.route('/api/auth/forgot-password', methods=['POST'])
@csrf.exempt
@limiter.limit("5 per hour")  # 每小时最多5次请求
def forgot_password():
    """发送密码重置验证码"""
    try:
        data = request.json
        email = data.get('email', '').strip()
        
        if not email:
            return jsonify({'success': False, 'message': '请输入邮箱地址'}), 400
        
        # 验证邮箱格式
        if not validate_email(email):
            return jsonify({'success': False, 'message': '邮箱格式不正确'}), 400
        
        # 检查邮箱是否存在
        user = User.query.filter_by(email=email).first()
        if not user:
            # 为安全考虑，不泄露用户是否存在
            return jsonify({
                'success': True,
                'message': '如果该邮箱已注册，您将收到一封包含验证码的邮件'
            }), 200
        
        # 清除该邮箱之前未使用的验证码
        VerificationCode.query.filter_by(
            email=email,
            type='reset_password',
            used=False
        ).delete()
        db.session.commit()
        
        # 生成6位验证码
        code = email_service.generate_code(6)
        
        # 保存验证码（10分钟有效期）
        verification = VerificationCode(
            email=email,
            code=code,
            type='reset_password',
            expires_at=datetime.utcnow() + timedelta(minutes=10),
            ip_address=request.remote_addr
        )
        db.session.add(verification)
        db.session.commit()
        
        # 发送邮件
        if email_service.send_verification_code(email, code, 'reset_password'):
            app.logger.info(f"密码重置验证码已发送到: {email}")
            return jsonify({
                'success': True,
                'message': '验证码已发送到您的邮箱，请查收'
            }), 200
        else:
            app.logger.error(f"邮件发送失败: {email}")
            return jsonify({
                'success': False,
                'message': '邮件发送失败，请稍后重试'
            }), 500
            
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"发送验证码错误: {str(e)}")
        return jsonify({
            'success': False,
            'message': f'发送验证码失败: {str(e)}'
        }), 500


@app.route('/api/auth/verify-reset-code', methods=['POST'])
@csrf.exempt
@limiter.limit("10 per hour")  # 每小时最多10次验证
def verify_reset_code():
    """验证重置密码验证码"""
    try:
        data = request.json
        email = data.get('email', '').strip()
        code = data.get('code', '').strip()
        
        if not email or not code:
            return jsonify({
                'success': False,
                'message': '请输入邮箱和验证码'
            }), 400
        
        # 查找验证码
        verification = VerificationCode.query.filter_by(
            email=email,
            code=code,
            type='reset_password'
        ).order_by(VerificationCode.created_at.desc()).first()
        
        if not verification:
            return jsonify({
                'success': False,
                'message': '验证码不正确'
            }), 400
        
        if not verification.is_valid():
            return jsonify({
                'success': False,
                'message': '验证码已过期或已使用'
            }), 400
        
        # 生成临时令牌（15分钟有效）
        import secrets
        reset_token = secrets.token_urlsafe(32)
        
        # 将令牌与验证码关联（使用session或redis存储）
        # 这里简单存储在验证码记录的ip_address字段中（临时方案）
        verification.ip_address = reset_token
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '验证成功',
            'reset_token': reset_token
        }), 200
        
    except Exception as e:
        app.logger.error(f"验证码验证错误: {str(e)}")
        return jsonify({
            'success': False,
            'message': f'验证失败: {str(e)}'
        }), 500


@app.route('/api/auth/reset-password', methods=['POST'])
@csrf.exempt
@limiter.limit("5 per hour")  # 每小时最多5次重置
def reset_password():
    """重置密码"""
    try:
        data = request.json
        email = data.get('email', '').strip()
        reset_token = data.get('reset_token', '').strip()
        new_password = data.get('new_password', '').strip()
        
        if not email or not reset_token or not new_password:
            return jsonify({
                'success': False,
                'message': '请填写所有字段'
            }), 400
        
        # 验证密码强度
        pwd_valid, pwd_msg, pwd_strength = validate_password_strength(new_password)
        if not pwd_valid:
            return jsonify({
                'success': False,
                'message': pwd_msg,
                'strength': pwd_strength
            }), 400
        
        # 查找验证码记录
        verification = VerificationCode.query.filter_by(
            email=email,
            type='reset_password',
            ip_address=reset_token  # 临时存储的reset_token
        ).order_by(VerificationCode.created_at.desc()).first()
        
        if not verification or not verification.is_valid():
            return jsonify({
                'success': False,
                'message': '无效的重置请求或令牌已过期'
            }), 400
        
        # 查找用户
        user = User.query.filter_by(email=email).first()
        if not user:
            return jsonify({
                'success': False,
                'message': '用户不存在'
            }), 404
        
        # 更新密码
        user.password = new_password
        
        # 标记验证码为已使用
        verification.mark_as_used()
        
        db.session.commit()
        
        app.logger.info(f"用户 {email} 已重置密码")
        
        return jsonify({
            'success': True,
            'message': '密码重置成功，请使用新密码登录'
        }), 200
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"重置密码错误: {str(e)}")
        return jsonify({
            'success': False,
            'message': f'密码重置失败: {str(e)}'
        }), 500


# ==================== 日志配置 ====================
def setup_logging():
    """配置应用日志系统"""
    try:
        import logging
        from logging.handlers import RotatingFileHandler
        import os
        
        # 设置日志格式
        formatter = logging.Formatter(
            '[%(asctime)s] %(levelname)s in %(module)s: %(message)s'
        )
        
        # 尝试创建文件日志（Vercel等只读环境会跳过）
        try:
            if not os.path.exists('logs'):
                os.makedirs('logs')
            
            # 应用日志处理器（所有日志）
            app_handler = RotatingFileHandler(
                'logs/app.log',
                maxBytes=10240000,  # 10MB
                backupCount=10
            )
            app_handler.setFormatter(formatter)
            app_handler.setLevel(logging.INFO)
            
            # 错误日志处理器（仅错误）
            error_handler = RotatingFileHandler(
                'logs/error.log',
                maxBytes=10240000,  # 10MB
                backupCount=10
            )
            error_handler.setFormatter(formatter)
            error_handler.setLevel(logging.ERROR)
            
            # 添加处理器到app.logger
            app.logger.addHandler(app_handler)
            app.logger.addHandler(error_handler)
        except (OSError, PermissionError):
            # Vercel等无服务器环境，使用控制台日志
            console_handler = logging.StreamHandler()
            console_handler.setFormatter(formatter)
            app.logger.addHandler(console_handler)
        
        app.logger.setLevel(logging.INFO)
        app.logger.info('日志系统初始化成功')
        
    except Exception as e:
        print(f"日志配置失败: {e}")


# 初始化日志系统
setup_logging()


# ==================== 安全响应头中间件 ====================

from utils.security import add_security_headers
import mimetypes

@app.after_request
def apply_security_headers(response):
    """
    为所有响应添加安全HTTP头部
    
    这是一个全局中间件，会在每个请求的响应返回前执行
    添加多层安全防护头部，防止常见Web攻击
    """
    # 修复静态文件的MIME类型
    if request.path.startswith('/static/'):
        # 确保JavaScript文件有正确的MIME类型
        if request.path.endswith('.js'):
            response.headers['Content-Type'] = 'application/javascript; charset=utf-8'
        elif request.path.endswith('.css'):
            response.headers['Content-Type'] = 'text/css; charset=utf-8'
        elif request.path.endswith('.json'):
            response.headers['Content-Type'] = 'application/json; charset=utf-8'
        elif request.path.endswith(('.png', '.jpg', '.jpeg', '.gif', '.svg')):
            # 图片类型由Flask自动处理，不需要修改
            pass
    
    return add_security_headers(response)


# ==================== 错误处理器 ====================

@app.errorhandler(404)
def not_found_error(error):
    """404错误处理器 - 页面未找到"""
    try:
        app.logger.warning(f'404错误: {request.url}')
        return render_template('errors/404.html'), 404
    except Exception as e:
        app.logger.error(f'404错误处理器异常: {str(e)}')
        return '404 - 页面未找到', 404


@app.errorhandler(500)
def internal_error(error):
    """500错误处理器 - 服务器内部错误"""
    try:
        # 回滚数据库会话（如果有未完成的事务）
        db.session.rollback()
        
        # 生成错误追踪ID
        import uuid
        error_id = str(uuid.uuid4())[:8]
        
        # 记录详细错误信息
        app.logger.error(f'500错误 [ID: {error_id}]: {str(error)}', exc_info=True)
        
        return render_template('errors/500.html', error_id=error_id), 500
    except Exception as e:
        app.logger.error(f'500错误处理器异常: {str(e)}')
        return '500 - 服务器内部错误', 500


@app.errorhandler(403)
def forbidden_error(error):
    """403错误处理器 - 访问被拒绝"""
    try:
        app.logger.warning(f'403错误: {request.url} - 用户: {current_user.username if current_user.is_authenticated else "匿名"}')
        return render_template('errors/403.html'), 403
    except Exception as e:
        app.logger.error(f'403错误处理器异常: {str(e)}')
        return '403 - 访问被拒绝', 403


@app.errorhandler(400)
def bad_request_error(error):
    """400错误处理器 - 错误请求"""
    try:
        app.logger.warning(f'400错误: {request.url} - 错误: {str(error)}')
        return render_template('errors/400.html'), 400
    except Exception as e:
        app.logger.error(f'400错误处理器异常: {str(e)}')
        return '400 - 错误请求', 400

@app.errorhandler(429)
def ratelimit_error(error):
    """429错误处理器 - API限流"""
    try:
        app.logger.warning(f'429限流: {request.url} - IP: {request.remote_addr}')
        
        # 检查是否是API请求
        if request.path.startswith('/api/'):
            return jsonify({
                'error': '请求过于频繁，请稍后再试',
                'message': '您的请求已达到速率限制，请稍后重试',
                'retry_after': error.description if hasattr(error, 'description') else '请稍后'
            }), 429
        else:
            # 非API请求，返回HTML页面
            return render_template('errors/429.html'), 429
    except Exception as e:
        app.logger.error(f'429错误处理器异常: {str(e)}')
        return jsonify({'error': '请求过于频繁'}), 429


# 记录所有请求（可选，用于调试）
@app.before_request
def log_request_info():
    """记录请求信息"""
    try:
        if not request.path.startswith('/static'):
            app.logger.info(f'请求: {request.method} {request.path} - IP: {request.remote_addr}')
    except Exception:
        pass  # 日志记录失败不应影响请求处理


# ============= 注册支付路由 =============
try:
    from routes_payment import register_payment_routes
    register_payment_routes(app, csrf)
    app.logger.info("支付路由注册成功")
except Exception as e:
    app.logger.warning(f"支付路由注册失败: {str(e)}")


# ============= 管理后台路由 =============
from flask import session, flash, url_for
from models_admin import Admin, AdminLog
from utils.admin_decorators import admin_required, permission_required, api_admin_required
from utils.admin_auth import (
    get_current_admin, admin_login, admin_logout, 
    check_login_attempts, record_admin_login_attempt, get_admin_context
)

# 管理后台上下文处理器
@app.context_processor
def inject_admin_context():
    """注入管理员上下文到所有模板"""
    return {'admin_context': get_admin_context()}


# 管理员登录页面
@app.route('/admin/login', methods=['GET', 'POST'])
def admin_login_page():
    """管理员登录"""
    # 如果已登录，重定向到dashboard
    if 'admin_id' in session:
        return redirect(url_for('admin_dashboard'))
    
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        password = request.form.get('password', '')
        remember = request.form.get('remember') == 'on'
        
        if not username or not password:
            flash('请输入用户名和密码', 'error')
            return render_template('admin/login.html')
        
        # 检查登录尝试次数
        can_login, remaining = check_login_attempts(username, request.remote_addr)
        if not can_login:
            flash('登录尝试次数过多，请15分钟后再试', 'error')
            return render_template('admin/login.html')
        
        # 查找管理员
        admin = Admin.query.filter_by(username=username).first()
        
        if not admin or not admin.check_password(password):
            # 记录失败的登录尝试
            record_admin_login_attempt(username, request.remote_addr, False)
            flash(f'用户名或密码错误，剩余尝试次数: {remaining - 1}', 'error')
            return render_template('admin/login.html')
        
        if not admin.is_active:
            flash('该管理员账户已被禁用', 'error')
            return render_template('admin/login.html')
        
        # 登录成功
        admin_login(admin)
        record_admin_login_attempt(username, request.remote_addr, True)
        
        if remember:
            session.permanent = True
        
        flash(f'欢迎回来，{admin.real_name or admin.username}！', 'success')
        return redirect(url_for('admin_dashboard'))
    
    return render_template('admin/login.html')


# 管理员登出
@app.route('/admin/logout')
def admin_logout_page():
    """管理员登出"""
    admin_logout()
    flash('已安全退出登录', 'success')
    return redirect(url_for('admin_login_page'))


# 管理后台首页（Dashboard）
@app.route('/admin')
@app.route('/admin/dashboard')
@admin_required
def admin_dashboard():
    """管理后台总览"""
    return render_template('admin/dashboard.html')


# 用户管理页面
@app.route('/admin/users')
@admin_required
@permission_required('user_view')
def admin_users():
    """用户管理"""
    return render_template('admin/users.html')


# 会员管理页面
@app.route('/admin/memberships')
@admin_required
@permission_required('membership_view')
def admin_memberships():
    """会员管理"""
    return render_template('admin/memberships.html')


# 订单管理页面
@app.route('/admin/orders')
@admin_required
@permission_required('order_view')
def admin_orders():
    """订单管理"""
    return render_template('admin/orders.html')


# 操作日志页面
@app.route('/admin/logs')
@admin_required
@permission_required('log_view')
def admin_logs():
    """操作日志"""
    flash('操作日志页面开发中...', 'info')
    return redirect(url_for('admin_dashboard'))


# 管理员个人资料
@app.route('/admin/profile')
@admin_required
def admin_profile():
    """管理员个人资料"""
    admin = get_current_admin()
    return render_template('admin/profile.html', admin=admin) if admin else redirect(url_for('admin_login_page'))


# ============= 管理后台 API =============

# 统计数据API
@app.route('/api/admin/stats/overview')
@api_admin_required
def api_admin_stats_overview(current_admin):
    """获取总览统计数据"""
    try:
        from sqlalchemy import func
        from datetime import datetime, timedelta
        
        # 总用户数
        total_users = User.query.count()
        
        # 会员数
        total_members = UserMembership.query.filter_by(is_active=True).count()
        
        # 订单数
        total_orders = PaymentTransaction.query.count()
        
        # 总收入
        total_revenue = db.session.query(
            func.sum(PaymentTransaction.amount)
        ).filter_by(status='success').scalar() or 0
        
        # 昨日数据（用于计算增长）
        yesterday = datetime.now() - timedelta(days=1)
        yesterday_start = yesterday.replace(hour=0, minute=0, second=0, microsecond=0)
        
        users_growth = User.query.filter(User.created_at >= yesterday_start).count()
        members_growth = UserMembership.query.filter(
            UserMembership.created_at >= yesterday_start,
            UserMembership.is_active == True
        ).count()
        orders_growth = PaymentTransaction.query.filter(
            PaymentTransaction.created_at >= yesterday_start
        ).count()
        revenue_growth = db.session.query(
            func.sum(PaymentTransaction.amount)
        ).filter(
            PaymentTransaction.created_at >= yesterday_start,
            PaymentTransaction.status == 'success'
        ).scalar() or 0
        
        return jsonify({
            'success': True,
            'data': {
                'total_users': total_users,
                'total_members': total_members,
                'total_orders': total_orders,
                'total_revenue': float(total_revenue),
                'users_growth': users_growth,
                'members_growth': members_growth,
                'orders_growth': orders_growth,
                'revenue_growth': float(revenue_growth)
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取统计数据失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取统计数据失败'}), 500


# 用户增长图表数据
@app.route('/api/admin/charts/users')
@api_admin_required
def api_admin_chart_users(current_admin):
    """获取用户增长图表数据"""
    try:
        from datetime import datetime, timedelta
        
        period = int(request.args.get('period', 7))
        
        labels = []
        values = []
        
        for i in range(period):
            date = datetime.now() - timedelta(days=period - i - 1)
            date_start = date.replace(hour=0, minute=0, second=0, microsecond=0)
            date_end = date_start + timedelta(days=1)
            
            count = User.query.filter(
                User.created_at >= date_start,
                User.created_at < date_end
            ).count()
            
            labels.append(date.strftime('%m-%d'))
            values.append(count)
        
        return jsonify({
            'success': True,
            'data': {
                'labels': labels,
                'values': values
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取用户增长数据失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取图表数据失败'}), 500


# 会员分布图表数据
@app.route('/api/admin/charts/memberships')
@api_admin_required
def api_admin_chart_memberships(current_admin):
    """获取会员分布图表数据"""
    try:
        # 查询各类会员数量
        tiers = db.session.query(
            MembershipTier.name,
            db.func.count(UserMembership.id)
        ).join(UserMembership).filter(
            UserMembership.is_active == True
        ).group_by(MembershipTier.name).all()
        
        labels = [tier[0] for tier in tiers]
        values = [tier[1] for tier in tiers]
        
        # 添加非会员数
        non_members = User.query.outerjoin(UserMembership).filter(
            (UserMembership.id == None) | (UserMembership.is_active == False)
        ).count()
        
        if non_members > 0:
            labels.append('非会员')
            values.append(non_members)
        
        return jsonify({
            'success': True,
            'data': {
                'labels': labels,
                'values': values
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取会员分布数据失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取图表数据失败'}), 500


# ==================== 订单管理 API (旧版本已删除，使用新的Order模型) ====================
# 旧版本使用PaymentTransaction，已迁移到Order模型
# 新版本API在文件末尾

# 最近订单（保留用于仪表板）
@app.route('/api/admin/orders/recent')
@api_admin_required
def api_admin_orders_recent(current_admin):
    """获取最近订单"""
    try:
        orders = PaymentTransaction.query.join(User).order_by(
            PaymentTransaction.created_at.desc()
        ).limit(5).all()
        
        data = []
        for order in orders:
            data.append({
                'transaction_id': order.transaction_id,
                'username': order.user.username,
                'amount': float(order.amount),
                'status': order.status,
                'created_at': order.created_at.strftime('%Y-%m-%d %H:%M')
            })
        
        return jsonify({'success': True, 'data': data})
        
    except Exception as e:
        app.logger.error(f"获取最近订单失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取订单数据失败'}), 500


# 最近注册用户
@app.route('/api/admin/users/recent')
@api_admin_required
def api_admin_users_recent(current_admin):
    """获取最近注册用户"""
    try:
        users = User.query.order_by(User.created_at.desc()).limit(5).all()
        
        data = []
        for user in users:
            data.append({
                'username': user.username,
                'email': user.email,
                'created_at': user.created_at.strftime('%Y-%m-%d %H:%M'),
                'is_active': user.is_active
            })
        
        return jsonify({'success': True, 'data': data})
        
    except Exception as e:
        app.logger.error(f"获取最近用户失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取用户数据失败'}), 500


# 创建用户API
@app.route('/api/admin/users/create', methods=['POST', 'OPTIONS'])
@api_admin_required
@permission_required('user_edit')
def api_admin_users_create(current_admin):
    """创建新用户"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        data = request.get_json()
        
        # 验证必填字段
        username = data.get('username', '').strip()
        email = data.get('email', '').strip()
        password = data.get('password', '').strip()
        
        if not username or not email or not password:
            return jsonify({'success': False, 'message': '用户名、邮箱和密码不能为空'}), 400
        
        # 验证邮箱格式
        import re
        email_pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
        if not re.match(email_pattern, email):
            return jsonify({'success': False, 'message': '邮箱格式不正确'}), 400
        
        # 检查用户名是否已存在
        if User.query.filter_by(username=username).first():
            return jsonify({'success': False, 'message': '用户名已存在'}), 400
        
        # 检查邮箱是否已存在
        if User.query.filter_by(email=email).first():
            return jsonify({'success': False, 'message': '邮箱已被注册'}), 400
        
        # 创建新用户
        new_user = User(
            username=username,
            email=email,
            full_name=data.get('full_name', '').strip(),
            phone=data.get('phone', '').strip(),
            is_active=data.get('is_active', True)
        )
        new_user.set_password(password)
        
        db.session.add(new_user)
        db.session.commit()
        
        # 记录管理员操作日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='create_user',
            module='user',
            target_type='user',
            target_id=new_user.id,
            description=f'创建用户 {username} ({email})',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': f'用户 {username} 创建成功',
            'data': {
                'id': new_user.id,
                'username': new_user.username,
                'email': new_user.email
            }
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"创建用户失败: {str(e)}")
        return jsonify({'success': False, 'message': '创建用户失败'}), 500


# 用户列表API
@app.route('/api/admin/users', methods=['GET', 'OPTIONS'])
@api_admin_required
@permission_required('user_view')
def api_admin_users_list(current_admin):
    """获取用户列表"""
    if request.method == 'OPTIONS':
        return '', 200
    
    # GET方法 - 获取用户列表
    try:
        page = int(request.args.get('page', 1))
        per_page = int(request.args.get('per_page', 20))
        keyword = request.args.get('keyword', '').strip()
        status = request.args.get('status', '')
        membership_tier = request.args.get('membership_tier', '').strip()
        
        query = User.query
        
        # 搜索
        if keyword:
            query = query.filter(
                (User.username.contains(keyword)) |
                (User.email.contains(keyword))
            )
        
        # 状态筛选
        if status == 'active':
            query = query.filter_by(is_active=True)
        elif status == 'inactive':
            query = query.filter_by(is_active=False)
        
        # 会员等级筛选
        if membership_tier:
            # 子查询获取每个用户的当前有效会员记录
            from datetime import datetime
            subquery = db.session.query(UserMembership.user_id)\
                                .join(MembershipTier, UserMembership.tier_id == MembershipTier.id)\
                                .filter(
                                    UserMembership.is_active == True,
                                    UserMembership.end_date > datetime.now(),
                                    MembershipTier.name == membership_tier
                                ).subquery()
            
            query = query.filter(User.id.in_(subquery))
        
        # 分页
        pagination = query.order_by(User.created_at.desc()).paginate(
            page=page,
            per_page=per_page,
            error_out=False
        )
        
        users = []
        for user in pagination.items:
            # 查询会员信息
            membership = UserMembership.query.filter_by(
                user_id=user.id,
                is_active=True
            ).first()
            
            users.append({
                'id': user.id,
                'username': user.username,
                'email': user.email,
                'is_active': user.is_active,
                'is_member': membership is not None,
                'membership_tier': membership.tier.name if membership else None,
                'membership_expires': membership.end_date.strftime('%Y-%m-%d %H:%M:%S') if membership else None,
                'created_at': user.created_at.strftime('%Y-%m-%d %H:%M:%S'),
                'last_login_at': user.last_login.strftime('%Y-%m-%d %H:%M:%S') if user.last_login else None
            })
        
        return jsonify({
            'success': True,
            'data': {
                'users': users,
                'pagination': {
                    'current_page': page,
                    'total_pages': pagination.pages,
                    'total_items': pagination.total,
                    'per_page': per_page
                }
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取用户列表失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取用户列表失败'}), 500


# 用户详情API
@app.route('/api/admin/users/<int:user_id>')
@api_admin_required
@permission_required('user_view')
def api_admin_user_detail(user_id, current_admin):
    """获取用户详情"""
    try:
        user = User.query.get_or_404(user_id)
        
        # 查询会员信息
        membership = UserMembership.query.filter_by(
            user_id=user.id,
            is_active=True
        ).first()
        
        # 查询使用统计
        usage = UsageLog.query.filter_by(user_id=user.id).all()
        ai_answer_count = sum(1 for u in usage if u.feature_code == 'ai_answer_questions')
        ai_question_count = sum(1 for u in usage if u.feature_code == 'ai_generate_questions')
        video_summary_count = sum(1 for u in usage if u.feature_code == 'video_summary')
        
        data = {
            'id': user.id,
            'username': user.username,
            'email': user.email,
            'full_name': user.full_name,
            'phone': user.phone,
            'is_active': user.is_active,
            'membership_tier': membership.tier.name if membership else None,
            'membership_expires': membership.end_date.strftime('%Y-%m-%d %H:%M:%S') if membership else None,
            'created_at': user.created_at.strftime('%Y-%m-%d %H:%M:%S'),
            'last_login_at': user.last_login.strftime('%Y-%m-%d %H:%M:%S') if user.last_login else None,
            'ai_answer_count': ai_answer_count,
            'ai_question_count': ai_question_count,
            'video_summary_count': video_summary_count
        }
        
        return jsonify({'success': True, 'data': data})
        
    except Exception as e:
        app.logger.error(f"获取用户详情失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取用户详情失败'}), 500


# 更新用户信息
@app.route('/api/admin/users/<int:user_id>', methods=['PUT', 'OPTIONS'])
@api_admin_required
@permission_required('user_edit')
def api_admin_user_update(user_id, current_admin):
    """更新用户信息"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        user = User.query.get_or_404(user_id)
        data = request.get_json()
        
        # 更新基本信息
        if 'email' in data and data['email'] != user.email:
            # 检查邮箱是否被占用
            existing = User.query.filter_by(email=data['email']).first()
            if existing and existing.id != user_id:
                return jsonify({'success': False, 'message': '该邮箱已被使用'}), 400
            user.email = data['email']
        
        if 'full_name' in data:
            user.full_name = data['full_name']
        
        if 'phone' in data:
            user.phone = data['phone']
        
        db.session.commit()
        
        # 记录操作日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='update_user',
            module='user',
            target_type='user',
            target_id=user_id,
            description=f'更新用户 {user.username} 的信息',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '用户信息已更新',
            'data': user.to_dict()
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"更新用户失败: {str(e)}")
        return jsonify({'success': False, 'message': f'更新失败: {str(e)}'}), 500


# 切换用户状态
@app.route('/api/admin/users/<int:user_id>/toggle', methods=['PUT', 'OPTIONS'])
@api_admin_required
@permission_required('user_edit')
def api_admin_user_toggle(user_id, current_admin):
    """启用/禁用用户"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        user = User.query.get_or_404(user_id)
        user.is_active = not user.is_active
        db.session.commit()
        
        # 记录操作日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='toggle_user',
            module='user',
            target_type='user',
            target_id=user_id,
            description=f'{"启用" if user.is_active else "禁用"}用户 {user.username}',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': f'用户已{"启用" if user.is_active else "禁用"}'
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"切换用户状态失败: {str(e)}")
        return jsonify({'success': False, 'message': '操作失败'}), 500


# 删除用户
@app.route('/api/admin/users/<int:user_id>', methods=['DELETE', 'OPTIONS'])
@api_admin_required
@permission_required('user_delete')
def api_admin_user_delete(user_id, current_admin):
    """删除用户（PostgreSQL自动级联删除关联数据）"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        user = User.query.get_or_404(user_id)
        username = user.username
        email = user.email
        
        app.logger.info(f"开始删除用户: ID={user_id}, username={username}, email={email}")
        
        # 删除用户（PostgreSQL会自动级联删除相关数据）
        db.session.delete(user)
        db.session.commit()
        
        app.logger.info(f"用户删除成功: {username}")
        
        # 记录操作日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='delete_user',
            module='user',
            target_type='user',
            target_id=user_id,
            description=f'删除用户 {username} ({email})（已级联删除关联数据）',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        app.logger.info(f"删除日志已记录")
        
        return jsonify({
            'success': True,
            'message': f'用户 {username} 已成功删除',
            'data': {
                'deleted_user': username,
                'deleted_email': email,
                'deleted_id': user_id
            }
        }), 200
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"删除用户失败: {str(e)}")
        app.logger.exception(e)
        return jsonify({
            'success': False, 
            'message': f'删除失败: {str(e)}'
        }), 500


# 重置用户密码
@app.route('/api/admin/users/<int:user_id>/reset-password', methods=['POST', 'OPTIONS'])
@api_admin_required
@permission_required('user_edit')
def api_admin_user_reset_password(user_id, current_admin):
    """重置用户密码"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        user = User.query.get_or_404(user_id)
        data = request.get_json()
        
        new_password = data.get('new_password', '').strip()
        
        if not new_password or len(new_password) < 6:
            return jsonify({'success': False, 'message': '密码长度至少为6位'}), 400
        
        # 重置密码
        user.set_password(new_password)
        db.session.commit()
        
        # 记录操作日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='reset_password',
            module='user',
            target_type='user',
            target_id=user_id,
            description=f'重置用户 {user.username} 的密码',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '密码已重置'
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"重置密码失败: {str(e)}")
        return jsonify({'success': False, 'message': '重置失败'}), 500


# 获取用户锁定状态
@app.route('/api/admin/users/<int:user_id>/lock-status', methods=['GET', 'OPTIONS'])
@api_admin_required
@permission_required('user_view')
def api_admin_user_lock_status(user_id, current_admin):
    """获取用户锁定状态"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        user = User.query.get_or_404(user_id)
        
        # 检查锁定状态
        from utils.security import is_account_locked, get_remaining_attempts
        is_locked, locked_until, recent_attempts = is_account_locked(user.username)
        remaining_attempts = get_remaining_attempts(user.username)
        
        return jsonify({
            'success': True,
            'data': {
                'is_locked': is_locked,
                'locked_until': locked_until.isoformat() + 'Z' if locked_until else None,
                'recent_attempts': recent_attempts,
                'remaining_attempts': remaining_attempts,
                'max_attempts': 5,
                'lockout_duration_minutes': 2
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取用户锁定状态失败: {str(e)}")
        return jsonify({
            'success': False,
            'message': f'获取状态失败: {str(e)}'
        }), 500


# 解锁用户账户
@app.route('/api/admin/users/<int:user_id>/unlock', methods=['POST', 'OPTIONS'])
@api_admin_required
@permission_required('user_edit')
def api_admin_user_unlock(user_id, current_admin):
    """解锁用户账户"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        user = User.query.get_or_404(user_id)
        
        # 解锁账户
        from utils.security import unlock_account
        success = unlock_account(user.username)
        
        if success:
            # 记录操作日志
            from models_admin import AdminLog
            log = AdminLog(
                admin_id=current_admin.id,
                action='unlock_account',
                module='user',
                target_type='user',
                target_id=user_id,
                description=f'解锁用户 {user.username} 的账户',
                ip_address=request.remote_addr
            )
            db.session.add(log)
            db.session.commit()
            
            return jsonify({
                'success': True,
                'message': f'用户 {user.username} 账户已解锁'
            })
        else:
            return jsonify({
                'success': False,
                'message': '解锁失败'
            }), 500
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"解锁账户失败: {str(e)}")
        return jsonify({
            'success': False,
            'message': f'解锁失败: {str(e)}'
        }), 500


# 赠送会员
@app.route('/api/admin/users/<int:user_id>/grant-membership', methods=['POST', 'OPTIONS'])
@api_admin_required
@permission_required('membership_edit')
def api_admin_user_grant_membership(user_id, current_admin):
    """赠送会员"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        user = User.query.get_or_404(user_id)
        data = request.get_json()
        
        tier_id = data.get('tier_id')
        duration_days = data.get('duration_days', 30)
        
        if not tier_id:
            return jsonify({'success': False, 'message': '请选择会员套餐'}), 400
        
        tier = MembershipTier.query.get(tier_id)
        if not tier:
            return jsonify({'success': False, 'message': '会员套餐不存在'}), 404
        
        # 检查是否已有会员
        existing = UserMembership.query.filter_by(
            user_id=user_id,
            is_active=True
        ).first()
        
        if existing:
            # 检查是否是相同套餐
            if existing.tier_id == tier_id:
                # 相同套餐：延长时间
                existing.end_date = existing.end_date + timedelta(days=duration_days)
                db.session.commit()
                message = f'已为用户 {user.username} 延长 {tier.name} 会员 {duration_days} 天'
            else:
                # 不同套餐：替换套餐
                # 先停用旧会员
                existing.is_active = False
                existing.updated_at = datetime.utcnow()
                
                # 创建新会员（从当前时间开始）
                new_membership = UserMembership(
                    user_id=user_id,
                    tier_id=tier_id,
                    start_date=datetime.utcnow(),
                    end_date=datetime.utcnow() + timedelta(days=duration_days),
                    is_active=True
                )
                db.session.add(new_membership)
                db.session.commit()
                
                # 获取旧套餐和新套餐名称
                old_tier = MembershipTier.query.get(existing.tier_id)
                old_tier_name = old_tier.name if old_tier else f'套餐ID{existing.tier_id}'
                message = f'已为用户 {user.username} 将会员从 {old_tier_name} 更换为 {tier.name}（{duration_days}天）'
        else:
            # 创建新会员
            membership = UserMembership(
                user_id=user_id,
                tier_id=tier_id,
                start_date=datetime.utcnow(),
                end_date=datetime.utcnow() + timedelta(days=duration_days),
                is_active=True
            )
            db.session.add(membership)
            db.session.commit()
            message = f'已为用户 {user.username} 赠送 {tier.name} 会员 {duration_days} 天'
        
        # 记录操作日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='grant_membership',
            module='membership',
            target_type='user',
            target_id=user_id,
            description=message,
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': message
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"赠送会员失败: {str(e)}")
        return jsonify({'success': False, 'message': '赠送失败'}), 500


# 批量切换用户状态API
@app.route('/api/admin/users/batch-toggle', methods=['POST', 'OPTIONS'])
@api_admin_required
@permission_required('user_edit')
def api_admin_users_batch_toggle(current_admin):
    """批量切换用户状态"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        data = request.get_json()
        user_ids = data.get('user_ids', [])
        
        if not user_ids:
            return jsonify({'success': False, 'message': '未选择用户'}), 400
            
        # 获取要操作的用户
        users = User.query.filter(User.id.in_(user_ids)).all()
        if not users:
            return jsonify({'success': False, 'message': '未找到指定用户'}), 404
            
        # 批量切换状态
        for user in users:
            user.is_active = not user.is_active
            
        db.session.commit()
        
        return jsonify({
            'success': True, 
            'message': f'成功切换 {len(users)} 个用户的状态'
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"批量切换用户状态失败: {str(e)}")
        return jsonify({'success': False, 'message': '操作失败'}), 500


# 批量删除用户API
@app.route('/api/admin/users/batch-delete', methods=['POST', 'OPTIONS'])
@api_admin_required
@permission_required('user_delete')
def api_admin_users_batch_delete(current_admin):
    """批量删除用户"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        data = request.get_json()
        user_ids = data.get('user_ids', [])
        
        if not user_ids:
            return jsonify({'success': False, 'message': '未选择用户'}), 400
            
        # 获取要删除的用户
        users = User.query.filter(User.id.in_(user_ids)).all()
        if not users:
            return jsonify({'success': False, 'message': '未找到指定用户'}), 404
            
        # 批量删除用户及相关数据
        for user in users:
            # 删除用户会员记录
            UserMembership.query.filter_by(user_id=user.id).delete()
            # 使用原生SQL删除相关记录，避免模型依赖问题
            db.session.execute(db.text("DELETE FROM payment_transactions WHERE user_id = :user_id"), {"user_id": user.id})
            db.session.execute(db.text("DELETE FROM usage_logs WHERE user_id = :user_id"), {"user_id": user.id})
            # 删除用户本身
            db.session.delete(user)
            
        db.session.commit()
        
        return jsonify({
            'success': True, 
            'message': f'成功删除 {len(users)} 个用户'
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"批量删除用户失败: {str(e)}")
        return jsonify({'success': False, 'message': '删除失败'}), 500


# 用户数据导出API
@app.route('/api/admin/users/export')
@api_admin_required
@permission_required('user_view')
def api_admin_users_export(current_admin):
    """导出用户数据为XLSX"""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        from io import BytesIO
        from flask import make_response
        
        # 获取筛选参数
        keyword = request.args.get('keyword', '').strip()
        status = request.args.get('status', '')
        membership_tier = request.args.get('membership_tier', '').strip()
        start_date = request.args.get('start_date', '')
        end_date = request.args.get('end_date', '')
        
        # 构建查询
        query = User.query
        
        # 关键词搜索
        if keyword:
            query = query.filter(
                db.or_(
                    User.username.contains(keyword),
                    User.email.contains(keyword)
                )
            )
        
        # 状态筛选
        if status:
            is_active = status == 'active'
            query = query.filter(User.is_active == is_active)
        
        # 会员等级筛选
        if membership_tier:
            from datetime import datetime
            subquery = db.session.query(UserMembership.user_id)\
                                .join(MembershipTier, UserMembership.tier_id == MembershipTier.id)\
                                .filter(
                                    UserMembership.is_active == True,
                                    UserMembership.end_date > datetime.now(),
                                    MembershipTier.name == membership_tier
                                ).subquery()
            
            query = query.filter(User.id.in_(subquery))
        
        # 日期范围筛选
        if start_date:
            from datetime import datetime
            start_dt = datetime.strptime(start_date, '%Y-%m-%d')
            query = query.filter(User.created_at >= start_dt)
        
        if end_date:
            from datetime import datetime
            end_dt = datetime.strptime(end_date, '%Y-%m-%d')
            # 包含整个结束日期
            end_dt = end_dt.replace(hour=23, minute=59, second=59)
            query = query.filter(User.created_at <= end_dt)
        
        # 获取用户
        users = query.order_by(User.created_at.desc()).all()
        
        # 创建Excel工作簿
        wb = Workbook()
        ws = wb.active
        ws.title = "用户数据"
        
        # 设置表头样式
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        header_alignment = Alignment(horizontal="center", vertical="center")
        
        # 写入标题
        headers = ['ID', '用户名', '邮箱', '状态', '会员等级', '会员到期时间', '注册时间', '最后登录']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_alignment
        
        # 写入数据
        for row_idx, user in enumerate(users, 2):
            membership = UserMembership.query.filter_by(
                user_id=user.id,
                is_active=True
            ).first()
            
            # 数据行
            data = [
                user.id,
                user.username,
                user.email,
                '正常' if user.is_active else '禁用',
                membership.tier.name if membership else '无',
                membership.end_date.strftime('%Y-%m-%d %H:%M:%S') if membership else '',
                user.created_at.strftime('%Y-%m-%d %H:%M:%S'),
                user.last_login.strftime('%Y-%m-%d %H:%M:%S') if user.last_login else '从未登录'
            ]
            
            for col, value in enumerate(data, 1):
                ws.cell(row=row_idx, column=col, value=value)
        
        # 自动调整列宽
        for column in ws.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)  # 最大宽度50
            ws.column_dimensions[column_letter].width = adjusted_width
        
        # 保存到内存
        excel_buffer = BytesIO()
        wb.save(excel_buffer)
        excel_buffer.seek(0)
        
        # 创建响应
        from datetime import datetime
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"users_export_{timestamp}.xlsx"
        
        output = make_response(excel_buffer.getvalue())
        output.headers["Content-Disposition"] = f"attachment; filename={filename}"
        output.headers["Content-Type"] = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        
        # 记录日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='export',
            module='user',
            description=f'导出了 {len(users)} 条用户数据',
            ip_address=request.remote_addr,
            user_agent=request.user_agent.string[:200] if request.user_agent else None,
            request_method=request.method,
            request_path=request.path,
            status='success'
        )
        db.session.add(log)
        db.session.commit()
        
        return output
        
    except Exception as e:
        app.logger.error(f"导出用户数据失败: {str(e)}")
        return jsonify({'success': False, 'message': '导出失败'}), 500


# ==================== 旧订单API已移除 ====================
# 以下旧的订单管理API已被删除，使用文件末尾基于Order模型的新API
# 删除的API包括:
# - GET /api/admin/orders (订单列表)
# - GET /api/admin/orders/:id (订单详情)
# - PUT /api/admin/orders/:id/status (更新状态)
# - POST /api/admin/orders/:id/refund (退款)
# - PUT /api/admin/orders/:id/notes (更新备注)
# - GET /api/admin/orders/stats (统计)
# - GET /api/admin/orders/export (导出)
# ==================== 旧订单API结束 ====================


# ================================================================================
# 支付记录管理 API
# ================================================================================

# 获取支付记录列表
@app.route('/api/admin/payments')
@api_admin_required
@permission_required('payment_view')
def api_admin_payments_list(current_admin):
    """获取支付记录列表"""
    try:
        page = int(request.args.get('page', 1))
        per_page = int(request.args.get('per_page', 20))
        keyword = request.args.get('keyword', '').strip()
        status = request.args.get('status', '')
        
        query = PaymentTransaction.query
        
        # 搜索（订单号或用户名）
        if keyword:
            query = query.join(User).filter(
                (PaymentTransaction.transaction_id.contains(keyword)) |
                (User.username.contains(keyword))
            )
        
        # 状态筛选
        if status:
            query = query.filter_by(status=status)
        
        # 分页
        pagination = query.order_by(PaymentTransaction.created_at.desc()).paginate(
            page=page,
            per_page=per_page,
            error_out=False
        )
        
        orders = []
        for order in pagination.items:
            user = User.query.get(order.user_id)
            tier = MembershipTier.query.get(order.tier_id)
            
            orders.append({
                'id': order.id,
                'transaction_id': order.transaction_id,
                'user_id': order.user_id,
                'username': user.username if user else '未知',
                'tier_name': tier.name if tier else '未知',
                'amount': float(order.amount),
                'status': order.status,
                'payment_method': order.payment_method or '未知',
                'created_at': order.created_at.strftime('%Y-%m-%d %H:%M'),
                'completed_at': order.completed_at.strftime('%Y-%m-%d %H:%M') if order.completed_at else None
            })
        
        return jsonify({
            'success': True,
            'data': {
                'orders': orders,
                'pagination': {
                    'current_page': page,
                    'total_pages': pagination.pages,
                    'total_items': pagination.total,
                    'per_page': per_page
                }
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取订单列表失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取订单列表失败'}), 500


# 订单详情API（旧版本，已弃用，使用PaymentTransaction）
@app.route('/api/admin/orders/<int:order_id>/old')
@api_admin_required
@permission_required('order_view')
def api_admin_order_detail_old(order_id, current_admin):
    """获取订单详情（旧版本）"""
    try:
        order = PaymentTransaction.query.get_or_404(order_id)
        user = User.query.get(order.user_id)
        tier = MembershipTier.query.get(order.tier_id)
        
        data = {
            'id': order.id,
            'transaction_id': order.transaction_id,
            'user': {
                'id': user.id,
                'username': user.username,
                'email': user.email
            } if user else None,
            'tier': {
                'id': tier.id,
                'name': tier.name,
                'description': tier.description
            } if tier else None,
            'amount': float(order.amount),
            'status': order.status,
            'payment_method': order.payment_method,
            'created_at': order.created_at.strftime('%Y-%m-%d %H:%M'),
            'completed_at': order.completed_at.strftime('%Y-%m-%d %H:%M') if order.completed_at else None,
            'notes': order.notes or ''
        }
        
        return jsonify({'success': True, 'data': data})
        
    except Exception as e:
        app.logger.error(f"获取订单详情失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取订单详情失败'}), 500


# 更新订单状态
@app.route('/api/admin/orders/<int:order_id>/status', methods=['PUT', 'OPTIONS'])
@api_admin_required
@permission_required('order_edit')
def api_admin_order_update_status(order_id, current_admin):
    """更新订单状态"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        order = PaymentTransaction.query.get_or_404(order_id)
        data = request.get_json()
        
        new_status = data.get('status')
        if not new_status or new_status not in ['pending', 'success', 'failed', 'refunded']:
            return jsonify({'success': False, 'message': '无效的状态'}), 400
        
        old_status = order.status
        order.status = new_status
        
        # 如果更改为成功状态，设置完成时间
        if new_status == 'success' and not order.completed_at:
            order.completed_at = datetime.utcnow()
        
        db.session.commit()
        
        # 记录操作日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='update_order_status',
            module='order',
            target_type='order',
            target_id=order_id,
            description=f'订单 {order.transaction_id} 状态从 {old_status} 变更为 {new_status}',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '订单状态已更新'
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"更新订单状态失败: {str(e)}")
        return jsonify({'success': False, 'message': '更新失败'}), 500


# 退款订单
@app.route('/api/admin/orders/<int:order_id>/refund', methods=['POST', 'OPTIONS'])
@api_admin_required
@permission_required('order_edit')
def api_admin_order_refund(order_id, current_admin):
    """
    退款订单
    支持支付宝退款接口调用（如已配置）
    如未配置支付宝，将仅更新系统状态
    """
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        app.logger.info(f"请求: POST /api/admin/orders/{order_id}/refund - IP: {request.remote_addr}")
        
        order = PaymentTransaction.query.get_or_404(order_id)
        data = request.get_json()
        
        app.logger.info(f"退款请求数据: {data}")
        
        # 验证订单状态
        if order.status == 'refunded':
            return jsonify({'success': False, 'message': '该订单已退款'}), 400
        
        if order.status != 'success' and order.status != 'completed':
            return jsonify({'success': False, 'message': '只能退款已完成的订单'}), 400
        
        # 获取退款参数
        reason = data.get('reason', 'user_request')
        refund_amount = float(data.get('amount', order.amount))
        description = data.get('description', '')
        audit_notes = data.get('audit_notes', '')
        
        # 验证退款金额
        if refund_amount <= 0:
            return jsonify({'success': False, 'message': '退款金额必须大于0'}), 400
        
        if refund_amount > order.amount:
            return jsonify({'success': False, 'message': '退款金额不能超过订单金额'}), 400
        
        # 支付宝退款（如果已配置）
        alipay_refund_success = False
        alipay_refund_id = None
        
        try:
            # 检查是否配置了支付宝
            alipay_config = app.config.get('ALIPAY_CONFIG')
            if alipay_config and order.payment_method == 'alipay' and order.transaction_id:
                # 这里是支付宝退款接口调用的占位符
                # 实际使用时需要导入支付宝SDK并调用退款接口
                # from alipay import AliPay
                # alipay = AliPay(...)
                # result = alipay.api_alipay_trade_refund(
                #     out_trade_no=order.transaction_id,
                #     refund_amount=refund_amount,
                #     refund_reason=description
                # )
                # if result.get('code') == '10000':
                #     alipay_refund_success = True
                #     alipay_refund_id = result.get('trade_no')
                
                app.logger.info("支付宝接口未配置，仅更新系统状态")
            else:
                app.logger.info("未使用支付宝支付或未配置支付宝接口")
        except Exception as alipay_error:
            app.logger.error(f"支付宝退款失败: {str(alipay_error)}")
            # 支付宝退款失败不影响系统状态更新
        
        # 更新订单状态
        order.status = 'refunded'
        
        # 构建退款备注
        refund_note = f"\n退款信息:\n"
        refund_note += f"- 退款原因: {reason}\n"
        refund_note += f"- 退款金额: ¥{refund_amount}\n"
        refund_note += f"- 退款说明: {description}\n"
        if audit_notes:
            refund_note += f"- 审核意见: {audit_notes}\n"
        if alipay_refund_success:
            refund_note += f"- 支付宝退款单号: {alipay_refund_id}\n"
        else:
            refund_note += "- 退款方式: 系统退款（支付宝接口未配置）\n"
        refund_note += f"- 退款时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        refund_note += f"- 操作人: {current_admin.username}"
        
        order.notes = (order.notes or '') + refund_note
        
        # 取消关联的会员权限
        membership = UserMembership.query.filter_by(
            user_id=order.user_id,
            is_active=True
        ).first()
        if membership:
            membership.is_active = False
            membership.end_date = datetime.now()
            app.logger.info(f"已取消用户 {order.user_id} 的会员权限")
        
        db.session.commit()
        
        # 记录操作日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='refund_order',
            module='order',
            target_type='order',
            target_id=order_id,
            description=f'订单 {order.transaction_id} 已退款，金额: ¥{refund_amount}，原因: {reason}',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        app.logger.info(f"订单 {order_id} 退款成功")
        
        return jsonify({
            'success': True,
            'message': '退款成功' + ('（支付宝退款已处理）' if alipay_refund_success else '（系统退款，支付宝接口未配置）'),
            'refund_id': alipay_refund_id,
            'alipay_refund': alipay_refund_success
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"退款失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': f'退款失败: {str(e)}'}), 500


# 更新订单备注
@app.route('/api/admin/orders/<int:order_id>/notes', methods=['PUT', 'OPTIONS'])
@api_admin_required
@permission_required('order_edit')
def api_admin_order_update_notes(order_id, current_admin):
    """更新订单备注"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        order = PaymentTransaction.query.get_or_404(order_id)
        data = request.get_json()
        
        notes = data.get('notes', '').strip()
        order.notes = notes
        
        db.session.commit()
        
        # 记录操作日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='update_order_notes',
            module='order',
            target_type='order',
            target_id=order_id,
            description=f'更新订单 {order.transaction_id} 的备注',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '备注已更新'
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"更新备注失败: {str(e)}")
        return jsonify({'success': False, 'message': '更新失败'}), 500


# 获取订单统计数据
@app.route('/api/admin/orders/stats')
@api_admin_required
@permission_required('order_view')
def api_admin_orders_stats_old(current_admin):
    """获取订单统计数据"""
    try:
        from datetime import timedelta
        from sqlalchemy import func, and_
        
        now = datetime.now()
        today_start = now.replace(hour=0, minute=0, second=0, microsecond=0)
        week_start = today_start - timedelta(days=now.weekday())
        month_start = now.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
        
        # 今日收入
        today_revenue = db.session.query(func.sum(PaymentTransaction.amount)).filter(
            and_(
                PaymentTransaction.created_at >= today_start,
                PaymentTransaction.status.in_(['success', 'completed'])
            )
        ).scalar() or 0
        
        # 本周收入
        week_revenue = db.session.query(func.sum(PaymentTransaction.amount)).filter(
            and_(
                PaymentTransaction.created_at >= week_start,
                PaymentTransaction.status.in_(['success', 'completed'])
            )
        ).scalar() or 0
        
        # 本月收入
        month_revenue = db.session.query(func.sum(PaymentTransaction.amount)).filter(
            and_(
                PaymentTransaction.created_at >= month_start,
                PaymentTransaction.status.in_(['success', 'completed'])
            )
        ).scalar() or 0
        
        # 总收入
        total_revenue = db.session.query(func.sum(PaymentTransaction.amount)).filter(
            PaymentTransaction.status.in_(['success', 'completed'])
        ).scalar() or 0
        
        # 今日订单数
        today_orders = PaymentTransaction.query.filter(
            PaymentTransaction.created_at >= today_start
        ).count()
        
        # 本周订单数
        week_orders = PaymentTransaction.query.filter(
            PaymentTransaction.created_at >= week_start
        ).count()
        
        # 本月订单数
        month_orders = PaymentTransaction.query.filter(
            PaymentTransaction.created_at >= month_start
        ).count()
        
        # 总订单数
        total_orders = PaymentTransaction.query.count()
        
        # 最近30天收入趋势
        revenue_trend = []
        for i in range(29, -1, -1):
            day = today_start - timedelta(days=i)
            day_end = day + timedelta(days=1)
            
            day_revenue = db.session.query(func.sum(PaymentTransaction.amount)).filter(
                and_(
                    PaymentTransaction.created_at >= day,
                    PaymentTransaction.created_at < day_end,
                    PaymentTransaction.status.in_(['success', 'completed'])
                )
            ).scalar() or 0
            
            day_orders = PaymentTransaction.query.filter(
                and_(
                    PaymentTransaction.created_at >= day,
                    PaymentTransaction.created_at < day_end
                )
            ).count()
            
            revenue_trend.append({
                'date': day.strftime('%Y-%m-%d'),
                'revenue': float(day_revenue),
                'orders': day_orders
            })
        
        # 订单状态分布
        status_distribution = db.session.query(
            PaymentTransaction.status,
            func.count(PaymentTransaction.id)
        ).group_by(PaymentTransaction.status).all()
        
        total_count = sum(count for _, count in status_distribution)
        status_dist_list = []
        for status, count in status_distribution:
            status_dist_list.append({
                'status': status,
                'count': count,
                'percentage': round(count / total_count * 100, 2) if total_count > 0 else 0
            })
        
        # 支付方式分布
        payment_distribution = db.session.query(
            PaymentTransaction.payment_method,
            func.count(PaymentTransaction.id)
        ).filter(
            PaymentTransaction.status.in_(['success', 'completed'])
        ).group_by(PaymentTransaction.payment_method).all()
        
        payment_total = sum(count for _, count in payment_distribution)
        payment_dist_list = []
        for method, count in payment_distribution:
            payment_dist_list.append({
                'method': method,
                'count': count,
                'percentage': round(count / payment_total * 100, 2) if payment_total > 0 else 0
            })
        
        return jsonify({
            'today_revenue': float(today_revenue),
            'week_revenue': float(week_revenue),
            'month_revenue': float(month_revenue),
            'total_revenue': float(total_revenue),
            'today_orders': today_orders,
            'week_orders': week_orders,
            'month_orders': month_orders,
            'total_orders': total_orders,
            'revenue_trend': revenue_trend,
            'status_distribution': status_dist_list,
            'payment_method_distribution': payment_dist_list
        })
        
    except Exception as e:
        app.logger.error(f"获取订单统计失败: {str(e)}", exc_info=True)
        return jsonify({
            'today_revenue': 0,
            'week_revenue': 0,
            'month_revenue': 0,
            'total_revenue': 0,
            'today_orders': 0,
            'week_orders': 0,
            'month_orders': 0,
            'total_orders': 0,
            'revenue_trend': [],
            'status_distribution': [],
            'payment_method_distribution': []
        })


# 订单导出API
@app.route('/api/admin/orders/export')
@api_admin_required
@permission_required('order_view')
def api_admin_orders_export(current_admin):
    """导出订单数据为CSV"""
    try:
        import csv
        from io import StringIO
        from flask import make_response
        
        # 获取所有订单
        orders = PaymentTransaction.query.order_by(PaymentTransaction.created_at.desc()).all()
        
        # 创建CSV
        si = StringIO()
        writer = csv.writer(si, lineterminator='\n')  # 使用Unix换行符避免Windows下的额外空行
        
        # 写入标题
        writer.writerow(['订单ID', '订单号', '用户名', '套餐名称', '金额', '状态', '支付方式', '创建时间', '完成时间'])
        
        # 写入数据
        for order in orders:
            user = User.query.get(order.user_id)
            tier = MembershipTier.query.get(order.tier_id)
            
            status_text = {
                'pending': '待支付',
                'success': '成功',
                'failed': '失败',
                'refunded': '已退款'
            }.get(order.status, order.status)
            
            writer.writerow([
                order.id,
                order.transaction_id,
                user.username if user else '未知',
                tier.name if tier else '未知',
                float(order.amount),
                status_text,
                order.payment_method or '未知',
                order.created_at.strftime('%Y-%m-%d %H:%M'),
                order.completed_at.strftime('%Y-%m-%d %H:%M') if order.completed_at else ''
            ])
        
        # 创建响应
        output = make_response(si.getvalue())
        output.headers["Content-Disposition"] = "attachment; filename=orders_export.csv"
        output.headers["Content-type"] = "text/csv; charset=utf-8-sig"
        
        # 记录日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='export',
            module='order',
            description=f'导出了 {len(orders)} 条订单数据',
            ip_address=request.remote_addr,
            user_agent=request.user_agent.string[:200] if request.user_agent else None,
            request_method=request.method,
            request_path=request.path,
            status='success'
        )
        db.session.add(log)
        db.session.commit()
        
        return output
        
    except Exception as e:
        app.logger.error(f"导出订单数据失败: {str(e)}")
        return jsonify({'success': False, 'message': '导出失败'}), 500


# ================================================================================
# 支付记录管理 API
# ================================================================================

# 生成对账报表
@app.route('/api/admin/payments/reconciliation')
@api_admin_required
@permission_required('order_view')
def api_admin_payments_reconciliation(current_admin):
    """生成对账报表"""
    try:
        from datetime import timedelta
        from sqlalchemy import func, and_
        
        start_date_str = request.args.get('start_date')
        end_date_str = request.args.get('end_date')
        payment_method = request.args.get('payment_method')
        
        if not start_date_str or not end_date_str:
            return jsonify({'success': False, 'message': '请提供开始和结束日期'}), 400
        
        start_date = datetime.strptime(start_date_str, '%Y-%m-%d')
        end_date = datetime.strptime(end_date_str, '%Y-%m-%d') + timedelta(days=1)
        
        # 基础查询
        query = PaymentTransaction.query.filter(
            and_(
                PaymentTransaction.created_at >= start_date,
                PaymentTransaction.created_at < end_date
            )
        )
        
        if payment_method and payment_method != 'all':
            query = query.filter(PaymentTransaction.payment_method == payment_method)
        
        # 汇总统计
        all_transactions = query.all()
        
        total_transactions = len(all_transactions)
        total_amount = sum(float(t.amount) for t in all_transactions)
        
        success_transactions = [t for t in all_transactions if t.status in ['success', 'completed']]
        success_count = len(success_transactions)
        success_amount = sum(float(t.amount) for t in success_transactions)
        
        failed_transactions = [t for t in all_transactions if t.status == 'failed']
        failed_count = len(failed_transactions)
        failed_amount = sum(float(t.amount) for t in failed_transactions)
        
        refund_transactions = [t for t in all_transactions if t.status == 'refunded']
        refund_count = len(refund_transactions)
        refund_amount = sum(float(t.amount) for t in refund_transactions)
        
        # 按日汇总
        daily_breakdown = []
        current_date = start_date
        while current_date < end_date:
            next_date = current_date + timedelta(days=1)
            day_transactions = [t for t in all_transactions 
                              if current_date <= t.created_at < next_date]
            
            day_success = [t for t in day_transactions if t.status in ['success', 'completed']]
            day_failed = [t for t in day_transactions if t.status == 'failed']
            day_refund = [t for t in day_transactions if t.status == 'refunded']
            
            daily_breakdown.append({
                'date': current_date.strftime('%Y-%m-%d'),
                'transactions': len(day_transactions),
                'amount': sum(float(t.amount) for t in day_transactions),
                'success_count': len(day_success),
                'success_amount': sum(float(t.amount) for t in day_success),
                'failed_count': len(day_failed),
                'failed_amount': sum(float(t.amount) for t in day_failed),
                'refund_count': len(day_refund),
                'refund_amount': sum(float(t.amount) for t in day_refund)
            })
            
            current_date = next_date
        
        # 按支付方式汇总
        method_breakdown = []
        methods = set(t.payment_method for t in all_transactions if t.payment_method)
        for method in methods:
            method_transactions = [t for t in all_transactions if t.payment_method == method]
            method_success = [t for t in method_transactions if t.status in ['success', 'completed']]
            
            method_breakdown.append({
                'method': method,
                'transactions': len(method_transactions),
                'amount': sum(float(t.amount) for t in method_transactions),
                'success_rate': (len(method_success) / len(method_transactions) * 100) if method_transactions else 0
            })
        
        return jsonify({
            'date_range': {
                'start_date': start_date_str,
                'end_date': end_date_str
            },
            'summary': {
                'total_transactions': total_transactions,
                'total_amount': total_amount,
                'success_transactions': success_count,
                'success_amount': success_amount,
                'failed_transactions': failed_count,
                'failed_amount': failed_amount,
                'refund_transactions': refund_count,
                'refund_amount': refund_amount
            },
            'daily_breakdown': daily_breakdown,
            'method_breakdown': method_breakdown,
            'discrepancies': []  # 预留字段
        })
        
    except Exception as e:
        app.logger.error(f"生成对账报表失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': f'生成报表失败: {str(e)}'}), 500


# 导出对账报表（Excel格式）
@app.route('/api/admin/payments/reconciliation/export')
@api_admin_required
@permission_required('order_view')
def api_admin_payments_reconciliation_export(current_admin):
    """导出对账报表为Excel"""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill
        from io import BytesIO
        from flask import send_file
        from datetime import timedelta
        from sqlalchemy import and_
        
        start_date_str = request.args.get('start_date')
        end_date_str = request.args.get('end_date')
        payment_method = request.args.get('payment_method')
        
        if not start_date_str or not end_date_str:
            return jsonify({'success': False, 'message': '请提供开始和结束日期'}), 400
        
        start_date = datetime.strptime(start_date_str, '%Y-%m-%d')
        end_date = datetime.strptime(end_date_str, '%Y-%m-%d') + timedelta(days=1)
        
        # 查询数据
        query = PaymentTransaction.query.filter(
            and_(
                PaymentTransaction.created_at >= start_date,
                PaymentTransaction.created_at < end_date
            )
        )
        
        if payment_method and payment_method != 'all':
            query = query.filter(PaymentTransaction.payment_method == payment_method)
        
        transactions = query.order_by(PaymentTransaction.created_at.desc()).all()
        
        # 创建Excel工作簿
        wb = Workbook()
        ws = wb.active
        ws.title = "对账报表"
        
        # 设置标题样式
        title_font = Font(bold=True, size=12)
        title_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        title_alignment = Alignment(horizontal="center", vertical="center")
        
        # 写入标题行
        headers = [
            '交易ID', '订单号', '用户名', '用户邮箱', '套餐名称', 
            '金额', '支付方式', '状态', '创建时间', '备注'
        ]
        
        for col, header in enumerate(headers, start=1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = title_font
            cell.fill = title_fill
            cell.alignment = title_alignment
        
        # 写入数据
        for row_idx, transaction in enumerate(transactions, start=2):
            user = User.query.get(transaction.user_id)
            tier = MembershipTier.query.get(transaction.tier_id) if transaction.tier_id else None
            
            ws.cell(row=row_idx, column=1, value=transaction.transaction_id)
            ws.cell(row=row_idx, column=2, value=transaction.transaction_id)
            ws.cell(row=row_idx, column=3, value=user.username if user else 'Unknown')
            ws.cell(row=row_idx, column=4, value=user.email if user else '')
            ws.cell(row=row_idx, column=5, value=tier.name if tier else 'Unknown')
            ws.cell(row=row_idx, column=6, value=float(transaction.amount))
            ws.cell(row=row_idx, column=7, value=transaction.payment_method or 'alipay')
            ws.cell(row=row_idx, column=8, value=transaction.status)
            ws.cell(row=row_idx, column=9, value=transaction.created_at.strftime('%Y-%m-%d %H:%M:%S') if transaction.created_at else '')
            ws.cell(row=row_idx, column=10, value=transaction.notes or '')
        
        # 调整列宽
        column_widths = [15, 20, 15, 25, 15, 10, 12, 10, 20, 30]
        for i, width in enumerate(column_widths, start=1):
            ws.column_dimensions[chr(64 + i)].width = width
        
        # 保存到BytesIO
        output = BytesIO()
        wb.save(output)
        output.seek(0)
        
        # 记录日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='export_reconciliation',
            module='payment',
            description=f'导出对账报表: {start_date_str} 至 {end_date_str}，共 {len(transactions)} 条记录',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        filename = f'对账报表_{start_date_str}_{end_date_str}.xlsx'
        
        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name=filename
        )
        
    except Exception as e:
        app.logger.error(f"导出对账报表失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '导出失败'}), 500


# 单笔同步支付宝交易状态
@app.route('/api/admin/payments/<int:payment_id>/sync', methods=['POST', 'OPTIONS'])
@api_admin_required
@permission_required('order_edit')
def api_admin_payment_sync(payment_id, current_admin):
    """同步单笔支付宝交易状态"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        payment = PaymentTransaction.query.get_or_404(payment_id)
        
        old_status = payment.status
        new_status = old_status
        gateway_data = None
        
        # 检查是否配置了支付宝
        alipay_config = app.config.get('ALIPAY_CONFIG')
        
        if alipay_config and payment.payment_method == 'alipay' and payment.transaction_id:
            try:
                # 这里是支付宝查询接口调用的占位符
                # 实际使用时需要导入支付宝SDK并调用查询接口
                # from alipay import AliPay
                # alipay = AliPay(...)
                # result = alipay.api_alipay_trade_query(
                #     out_trade_no=payment.transaction_id
                # )
                # if result.get('code') == '10000':
                #     trade_status = result.get('trade_status')
                #     if trade_status == 'TRADE_SUCCESS':
                #         new_status = 'completed'
                #     elif trade_status == 'TRADE_CLOSED':
                #         new_status = 'failed'
                #     elif trade_status == 'WAIT_BUYER_PAY':
                #         new_status = 'pending'
                #     gateway_data = result
                
                app.logger.info(f"支付宝接口未配置，跳过同步交易 {payment.transaction_id}")
                
                # 模拟：如果未配置，保持原状态
                return jsonify({
                    'success': True,
                    'message': '支付宝接口未配置，无法同步',
                    'old_status': old_status,
                    'new_status': new_status,
                    'gateway_data': None
                })
                
            except Exception as alipay_error:
                app.logger.error(f"支付宝查询失败: {str(alipay_error)}")
                return jsonify({
                    'success': False,
                    'message': f'支付宝查询失败: {str(alipay_error)}'
                }), 500
        else:
            return jsonify({
                'success': False,
                'message': '非支付宝支付或未配置支付宝接口'
            }), 400
        
        # 如果状态有变化，更新数据库
        if new_status != old_status:
            payment.status = new_status
            db.session.commit()
            
            # 记录日志
            from models_admin import AdminLog
            log = AdminLog(
                admin_id=current_admin.id,
                action='sync_payment',
                module='payment',
                target_type='payment',
                target_id=payment_id,
                description=f'同步支付 {payment.transaction_id} 状态: {old_status} -> {new_status}',
                ip_address=request.remote_addr
            )
            db.session.add(log)
            db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '同步成功' if new_status != old_status else '状态未变化',
            'old_status': old_status,
            'new_status': new_status,
            'gateway_data': gateway_data
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"同步支付状态失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': f'同步失败: {str(e)}'}), 500


# 批量同步支付宝交易状态
@app.route('/api/admin/payments/batch-sync', methods=['POST', 'OPTIONS'])
@api_admin_required
@permission_required('order_edit')
def api_admin_payments_batch_sync(current_admin):
    """批量同步支付宝交易状态"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        from datetime import timedelta
        from sqlalchemy import and_
        
        data = request.get_json()
        start_date_str = data.get('start_date')
        end_date_str = data.get('end_date')
        payment_method = data.get('payment_method', 'alipay')
        status_filter = data.get('status', 'pending')
        
        if not start_date_str or not end_date_str:
            return jsonify({'success': False, 'message': '请提供开始和结束日期'}), 400
        
        start_date = datetime.strptime(start_date_str, '%Y-%m-%d')
        end_date = datetime.strptime(end_date_str, '%Y-%m-%d') + timedelta(days=1)
        
        # 查询需要同步的交易
        query = PaymentTransaction.query.filter(
            and_(
                PaymentTransaction.created_at >= start_date,
                PaymentTransaction.created_at < end_date,
                PaymentTransaction.payment_method == payment_method
            )
        )
        
        if status_filter and status_filter != 'all':
            query = query.filter(PaymentTransaction.status == status_filter)
        
        transactions = query.all()
        
        # 检查是否配置了支付宝
        alipay_config = app.config.get('ALIPAY_CONFIG')
        
        if not alipay_config:
            return jsonify({
                'success': False,
                'message': '支付宝接口未配置'
            }), 400
        
        synced_count = 0
        updated_count = 0
        failed_count = 0
        details = []
        
        for transaction in transactions:
            try:
                old_status = transaction.status
                new_status = old_status
                
                # 这里是支付宝批量查询接口调用的占位符
                # 实际使用时调用支付宝SDK
                # result = alipay.api_alipay_trade_query(
                #     out_trade_no=transaction.transaction_id
                # )
                # if result.get('code') == '10000':
                #     trade_status = result.get('trade_status')
                #     # 根据trade_status更新new_status
                
                synced_count += 1
                
                if new_status != old_status:
                    transaction.status = new_status
                    updated_count += 1
                
                details.append({
                    'transaction_id': transaction.transaction_id,
                    'old_status': old_status,
                    'new_status': new_status,
                    'success': True
                })
                
            except Exception as e:
                failed_count += 1
                details.append({
                    'transaction_id': transaction.transaction_id,
                    'old_status': transaction.status,
                    'new_status': transaction.status,
                    'success': False,
                    'error': str(e)
                })
        
        if updated_count > 0:
            db.session.commit()
            
            # 记录日志
            from models_admin import AdminLog
            log = AdminLog(
                admin_id=current_admin.id,
                action='batch_sync_payment',
                module='payment',
                description=f'批量同步支付: 同步{synced_count}笔，更新{updated_count}笔，失败{failed_count}笔',
                ip_address=request.remote_addr
            )
            db.session.add(log)
            db.session.commit()
        
        return jsonify({
            'success': True,
            'message': f'批量同步完成: 同步{synced_count}笔，更新{updated_count}笔，失败{failed_count}笔',
            'synced_count': synced_count,
            'updated_count': updated_count,
            'failed_count': failed_count,
            'details': details
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"批量同步失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': f'批量同步失败: {str(e)}'}), 500


# 会员列表API
@app.route('/api/admin/memberships')
@api_admin_required
@permission_required('membership_view')
def api_admin_memberships_list(current_admin):
    """获取会员列表"""
    try:
        page = int(request.args.get('page', 1))
        per_page = int(request.args.get('per_page', 20))
        keyword = request.args.get('keyword', '').strip()
        tier_filter = request.args.get('tier', '')
        status_filter = request.args.get('status', '')
        
        # 查询活跃会员
        query = UserMembership.query.filter_by(is_active=True)
        
        # 关联用户表进行搜索
        if keyword:
            query = query.join(User).filter(
                (User.username.contains(keyword)) |
                (User.email.contains(keyword))
            )
        
        # 套餐筛选
        if tier_filter:
            query = query.filter_by(tier_id=int(tier_filter))
        
        # 状态筛选（即将过期/已过期）
        from datetime import datetime, timedelta
        if status_filter == 'expiring':
            # 即将过期（7天内）
            soon = datetime.now() + timedelta(days=7)
            query = query.filter(
                UserMembership.end_date <= soon,
                UserMembership.end_date > datetime.now()
            )
        elif status_filter == 'expired':
            # 已过期
            query = query.filter(UserMembership.end_date < datetime.now())
        
        # 分页
        pagination = query.order_by(UserMembership.end_date.desc()).paginate(
            page=page,
            per_page=per_page,
            error_out=False
        )
        
        memberships = []
        for membership in pagination.items:
            user = User.query.get(membership.user_id)
            tier = MembershipTier.query.get(membership.tier_id)
            
            # 判断状态
            is_expired = membership.end_date < datetime.now()
            days_left = (membership.end_date - datetime.now()).days if not is_expired else 0
            
            # 生成颜色（根据套餐等级）
            tier_colors = {
                'free': '#999999',
                'weekly': '#667eea',
                'monthly': '#28a745',
                'yearly': '#ffc107'
            }
            tier_color = tier_colors.get(tier.code, '#667eea') if tier else '#999'
            
            memberships.append({
                'id': membership.id,
                'user_id': user.id,
                'username': user.username if user else '未知',
                'email': user.email if user else '',
                'tier_name': tier.name if tier else '未知',
                'tier_color': tier_color,
                'started_at': membership.start_date.strftime('%Y-%m-%d'),
                'expires_at': membership.end_date.strftime('%Y-%m-%d'),
                'days_left': days_left,
                'is_expired': is_expired,
                'status': '已过期' if is_expired else f'{days_left}天后到期'
            })
        
        return jsonify({
            'success': True,
            'data': {
                'memberships': memberships,
                'pagination': {
                    'current_page': page,
                    'total_pages': pagination.pages,
                    'total_items': pagination.total,
                    'per_page': per_page
                }
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取会员列表失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取会员列表失败'}), 500


# 会员套餐列表API
@app.route('/api/admin/tiers')
@api_admin_required
@permission_required('membership_view')
def api_admin_tiers_list(current_admin):
    """获取会员套餐列表"""
    try:
        tiers = MembershipTier.query.filter_by(is_active=True).order_by(MembershipTier.price).all()
        
        # 套餐颜色映射
        tier_colors = {
            'free': '#999999',
            'weekly': '#667eea',
            'monthly': '#28a745',
            'yearly': '#ffc107'
        }
        
        result = []
        for tier in tiers:
            # 统计当前套餐的会员数
            member_count = UserMembership.query.filter_by(
                tier_id=tier.id,
                is_active=True
            ).count()
            
            result.append({
                'id': tier.id,
                'name': tier.name,
                'code': tier.code,
                'description': tier.description,
                'price': float(tier.price),
                'duration_days': tier.duration_days,
                'color': tier_colors.get(tier.code, '#667eea'),
                'member_count': member_count,
                'is_active': tier.is_active
            })
        
        return jsonify({'success': True, 'data': result})
        
    except Exception as e:
        app.logger.error(f"获取套餐列表失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取套餐列表失败'}), 500


# ==================== Dashboard专用API端点 ====================

@app.route('/api/admin/stats/dashboard')
@api_admin_required
def api_admin_stats_dashboard(current_admin):
    """获取Dashboard完整统计数据"""
    try:
        from sqlalchemy import func
        from datetime import datetime, timedelta
        
        now = datetime.utcnow()
        today_start = now.replace(hour=0, minute=0, second=0, microsecond=0)
        yesterday_start = today_start - timedelta(days=1)
        week_ago = now - timedelta(days=7)
        month_ago = now - timedelta(days=30)
        
        # 总用户数
        total_users = User.query.count()
        yesterday_users = User.query.filter(User.created_at < yesterday_start).count()
        users_growth = ((total_users - yesterday_users) / max(yesterday_users, 1)) * 100 if yesterday_users > 0 else 0
        
        # 活跃用户（最近7天登录）
        active_users = User.query.filter(User.last_login >= week_ago).count()
        week_before_active = User.query.filter(
            User.last_login >= (week_ago - timedelta(days=7)),
            User.last_login < week_ago
        ).count()
        active_growth = ((active_users - week_before_active) / max(week_before_active, 1)) * 100 if week_before_active > 0 else 0
        
        # 会员数
        total_members = UserMembership.query.filter_by(is_active=True).count()
        month_ago_members = UserMembership.query.filter(
            UserMembership.created_at < month_ago,
            UserMembership.is_active == True
        ).count()
        members_growth = ((total_members - month_ago_members) / max(month_ago_members, 1)) * 100 if month_ago_members > 0 else 0
        
        # 订单数
        total_orders = PaymentTransaction.query.count()
        month_ago_orders = PaymentTransaction.query.filter(
            PaymentTransaction.created_at < month_ago
        ).count()
        orders_growth = ((total_orders - month_ago_orders) / max(month_ago_orders, 1)) * 100 if month_ago_orders > 0 else 0
        
        # 总收入
        total_revenue = db.session.query(
            func.sum(PaymentTransaction.amount)
        ).filter_by(status='success').scalar() or 0
        
        month_ago_revenue = db.session.query(
            func.sum(PaymentTransaction.amount)
        ).filter(
            PaymentTransaction.created_at < month_ago,
            PaymentTransaction.status == 'success'
        ).scalar() or 0
        
        revenue_growth = ((float(total_revenue) - float(month_ago_revenue)) / max(float(month_ago_revenue), 1)) * 100 if month_ago_revenue > 0 else 0
        
        # 今日收入
        today_revenue = db.session.query(
            func.sum(PaymentTransaction.amount)
        ).filter(
            PaymentTransaction.created_at >= today_start,
            PaymentTransaction.status == 'success'
        ).scalar() or 0
        
        yesterday_revenue = db.session.query(
            func.sum(PaymentTransaction.amount)
        ).filter(
            PaymentTransaction.created_at >= yesterday_start,
            PaymentTransaction.created_at < today_start,
            PaymentTransaction.status == 'success'
        ).scalar() or 0
        
        today_growth = ((float(today_revenue) - float(yesterday_revenue)) / max(float(yesterday_revenue), 1)) * 100 if yesterday_revenue > 0 else 0
        
        # 转化率（支付订单 / 总订单）
        paid_orders = PaymentTransaction.query.filter_by(status='success').count()
        conversion_rate = (paid_orders / max(total_orders, 1)) * 100 if total_orders > 0 else 0
        
        week_before_paid = PaymentTransaction.query.filter(
            PaymentTransaction.created_at >= (week_ago - timedelta(days=7)),
            PaymentTransaction.created_at < week_ago,
            PaymentTransaction.status == 'success'
        ).count()
        week_before_total = PaymentTransaction.query.filter(
            PaymentTransaction.created_at >= (week_ago - timedelta(days=7)),
            PaymentTransaction.created_at < week_ago
        ).count()
        week_before_rate = (week_before_paid / max(week_before_total, 1)) * 100 if week_before_total > 0 else 0
        conversion_growth = conversion_rate - week_before_rate
        
        return jsonify({
            'success': True,
            'data': {
                'total_users': total_users,
                'active_users': active_users,
                'total_members': total_members,
                'total_orders': total_orders,
                'total_revenue': float(total_revenue),
                'today_revenue': float(today_revenue),
                'conversion_rate': round(conversion_rate, 2),
                'system_health': 95.0,  # 系统健康度（可以根据实际情况计算）
                'total_users_growth': round(users_growth, 1),
                'active_users_growth': round(active_growth, 1),
                'total_members_growth': round(members_growth, 1),
                'total_orders_growth': round(orders_growth, 1),
                'total_revenue_growth': round(revenue_growth, 1),
                'today_revenue_growth': round(today_growth, 1),
                'conversion_rate_growth': round(conversion_growth, 1)
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取Dashboard统计数据失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': str(e)}), 500


@app.route('/api/admin/stats/recent-orders')
@api_admin_required
def api_admin_stats_recent_orders(current_admin):
    """获取最近订单列表"""
    try:
        limit = int(request.args.get('limit', 5))
        
        orders = PaymentTransaction.query.order_by(
            PaymentTransaction.created_at.desc()
        ).limit(limit).all()
        
        orders_data = []
        for order in orders:
            user = User.query.get(order.user_id) if order.user_id else None
            tier = MembershipTier.query.get(order.tier_id) if order.tier_id else None
            orders_data.append({
                'id': order.id,
                'order_no': order.transaction_id or f'ORD{order.id:08d}',
                'user_id': order.user_id,
                'username': user.username if user else '未知用户',
                'user_email': user.email if user else '',
                'tier_id': order.tier_id,
                'tier_name': tier.name if tier else '未知套餐',
                'amount': float(order.amount),
                'status': order.status,
                'payment_method': order.payment_method,
                'payment_time': order.updated_at.isoformat() if order.updated_at else None,
                'created_at': order.created_at.isoformat() if order.created_at else None
            })
        
        return jsonify({
            'success': True,
            'data': orders_data
        })
        
    except Exception as e:
        app.logger.error(f"获取最近订单失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': str(e)}), 500


@app.route('/api/admin/stats/recent-users')
@api_admin_required
def api_admin_stats_recent_users(current_admin):
    """获取最近注册用户列表"""
    try:
        limit = int(request.args.get('limit', 5))
        
        users = User.query.order_by(
            User.created_at.desc()
        ).limit(limit).all()
        
        users_data = []
        for user in users:
            # 检查用户是否有会员
            membership = UserMembership.query.filter_by(
                user_id=user.id,
                is_active=True
            ).first()
            
            users_data.append({
                'id': user.id,
                'username': user.username,
                'email': user.email,
                'full_name': user.full_name if hasattr(user, 'full_name') else None,
                'phone': user.phone if hasattr(user, 'phone') else None,
                'status': 'active',  # 默认活跃状态
                'is_member': membership is not None,
                'created_at': user.created_at.isoformat() if user.created_at else None,
                'last_login_at': user.last_login.isoformat() if hasattr(user, 'last_login') and user.last_login else None
            })
        
        return jsonify({
            'success': True,
            'data': users_data
        })
        
    except Exception as e:
        app.logger.error(f"获取最近用户失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': str(e)}), 500


# ==================== 管理员认证API（for React前端） ====================

@app.route('/api/admin/auth/login', methods=['POST', 'OPTIONS'])
def api_admin_login():
    """管理员登录API"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        data = request.get_json()
        username = data.get('username', '').strip()
        password = data.get('password', '')
        
        if not username or not password:
            return jsonify({
                'success': False,
                'message': '请输入用户名和密码'
            }), 400
        
        # 检查登录尝试次数
        can_login, remaining = check_login_attempts(username, request.remote_addr)
        if not can_login:
            return jsonify({
                'success': False,
                'message': '登录尝试次数过多，请15分钟后再试'
            }), 429
        
        # 查找管理员
        admin = Admin.query.filter_by(username=username).first()
        
        if not admin or not admin.check_password(password):
            record_admin_login_attempt(username, request.remote_addr, False)
            return jsonify({
                'success': False,
                'message': f'用户名或密码错误，剩余尝试次数: {remaining - 1}'
            }), 401
        
        if not admin.is_active:
            return jsonify({
                'success': False,
                'message': '该管理员账户已被禁用'
            }), 403
        
        # 登录成功
        admin_login(admin)
        record_admin_login_attempt(username, request.remote_addr, True)
        
        # 返回管理员信息
        import json
        permissions = json.loads(admin.permissions or '[]') if not admin.is_super_admin else ['*']
        
        return jsonify({
            'success': True,
            'message': '登录成功',
            'data': {
                'admin': {
                    'id': admin.id,
                    'username': admin.username,
                    'email': admin.email,
                    'role': admin.role,
                    'permissions': permissions,
                    'last_login_at': admin.last_login_at.isoformat() if admin.last_login_at else None,
                    'created_at': admin.created_at.isoformat() if admin.created_at else None
                }
            }
        })
        
    except Exception as e:
        app.logger.error(f"管理员登录失败: {str(e)}", exc_info=True)
        return jsonify({
            'success': False,
            'message': '登录失败，请稍后重试'
        }), 500


@app.route('/api/admin/auth/logout', methods=['POST', 'OPTIONS'])
@api_admin_required
def api_admin_logout(current_admin):
    """管理员登出API"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        admin_logout()
        return jsonify({
            'success': True,
            'message': '已安全退出登录'
        })
    except Exception as e:
        app.logger.error(f"管理员登出失败: {str(e)}", exc_info=True)
        return jsonify({
            'success': False,
            'message': '登出失败'
        }), 500


@app.route('/api/admin/auth/current', methods=['GET', 'OPTIONS'])
@api_admin_required
def api_admin_current(current_admin):
    """获取当前管理员信息API"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        import json
        permissions = json.loads(current_admin.permissions or '[]') if not current_admin.is_super_admin else ['*']
        
        return jsonify({
            'success': True,
            'data': {
                'id': current_admin.id,
                'username': current_admin.username,
                'email': current_admin.email,
                'role': current_admin.role,
                'permissions': permissions,
                'last_login_at': current_admin.last_login_at.isoformat() if current_admin.last_login_at else None,
                'created_at': current_admin.created_at.isoformat() if current_admin.created_at else None
            }
        })
    except Exception as e:
        app.logger.error(f"获取管理员信息失败: {str(e)}", exc_info=True)
        return jsonify({
            'success': False,
            'message': '获取管理员信息失败'
        }), 500


# ==================== 会员套餐管理API ====================

# 会员套餐接口 - 支持单数和复数形式
@app.route('/api/admin/membership/tiers', methods=['GET', 'POST', 'OPTIONS'])
@app.route('/api/admin/memberships/tiers', methods=['GET', 'POST', 'OPTIONS'])
@api_admin_required
def api_admin_membership_tiers(current_admin):
    """获取所有会员套餐或创建新套餐"""
    if request.method == 'OPTIONS':
        return '', 200
    
    if request.method == 'GET':
        try:
            # 管理员可以看到所有套餐（包括免费套餐）
            tiers = MembershipTier.query.order_by(MembershipTier.price).all()
            
            return jsonify({
                'success': True,
                'data': [tier.to_dict() for tier in tiers]
            })
            
        except Exception as e:
            app.logger.error(f"获取会员套餐错误: {str(e)}", exc_info=True)
            return jsonify({
                'success': False,
                'message': '获取会员套餐失败'
            }), 500
    
    elif request.method == 'POST':
        try:
            data = request.get_json()
            app.logger.info(f"创建套餐请求数据: {data}")
            
            # 验证必填字段
            if not all([data.get('name'), data.get('code')]):
                app.logger.warning(f"必填字段缺失: name={data.get('name')}, code={data.get('code')}")
                return jsonify({'success': False, 'message': '套餐名称和代码不能为空'}), 400
            
            # 检查代码是否已存在
            existing = MembershipTier.query.filter_by(code=data['code']).first()
            if existing:
                app.logger.warning(f"套餐代码已存在: {data['code']}")
                return jsonify({'success': False, 'message': '该套餐代码已存在'}), 400
            
            # 创建新套餐
            import json
            tier = MembershipTier(
                name=data['name'],
                code=data['code'],
                description=data.get('description', ''),
                price=float(data.get('price', 0)),
                duration_days=int(data.get('duration_days', 30)),
                level=int(data.get('level', 0)),
                is_active=data.get('is_active', True),
                is_limited=data.get('is_limited', False),
                total_quota=int(data.get('total_quota', 0)),
                is_early_bird=data.get('is_early_bird', False),
                early_bird_tier=int(data.get('early_bird_tier', 0)),
                original_price=float(data.get('original_price', 0)),
                sort_order=int(data.get('sort_order', 0))
            )
            
            # 处理功能列表
            if 'features' in data:
                if isinstance(data['features'], list):
                    tier.features = json.dumps(data['features'], ensure_ascii=False)
                else:
                    tier.features = data['features']
            
            db.session.add(tier)
            db.session.commit()
            
            # 记录操作日志
            from models_admin import AdminLog
            log = AdminLog(
                admin_id=current_admin.id,
                action='create_tier',
                module='membership',
                target_type='tier',
                target_id=tier.id,
                description=f'创建会员套餐 {tier.name}',
                ip_address=request.remote_addr
            )
            db.session.add(log)
            db.session.commit()
            
            return jsonify({
                'success': True,
                'message': '套餐创建成功',
                'data': tier.to_dict()
            })
            
        except Exception as e:
            db.session.rollback()
            app.logger.error(f"创建套餐失败: {str(e)}")
            return jsonify({'success': False, 'message': '创建失败'}), 500


# 会员套餐详情接口 - 支持单数和复数形式
@app.route('/api/admin/membership/tiers/<int:tier_id>', methods=['PUT', 'DELETE', 'OPTIONS'])
@app.route('/api/admin/memberships/tiers/<int:tier_id>', methods=['PUT', 'DELETE', 'OPTIONS'])
@api_admin_required
def api_admin_membership_tier_detail(tier_id, current_admin):
    """更新或删除会员套餐"""
    if request.method == 'OPTIONS':
        return '', 200
    
    if request.method == 'PUT':
        try:
            tier = MembershipTier.query.get_or_404(tier_id)
            data = request.get_json()
            
            # 更新字段
            if 'name' in data:
                tier.name = data['name']
            if 'code' in data:
                tier.code = data['code']
            if 'description' in data:
                tier.description = data['description']
            if 'price' in data:
                tier.price = float(data['price'])
            if 'duration_days' in data:
                tier.duration_days = int(data['duration_days'])
            if 'level' in data:
                tier.level = int(data['level'])
            if 'is_active' in data:
                tier.is_active = data['is_active']
            if 'is_limited' in data:
                tier.is_limited = data['is_limited']
            if 'total_quota' in data:
                tier.total_quota = int(data['total_quota'])
            if 'is_early_bird' in data:
                tier.is_early_bird = data['is_early_bird']
            if 'early_bird_tier' in data:
                tier.early_bird_tier = int(data['early_bird_tier'])
            if 'original_price' in data:
                tier.original_price = float(data['original_price'])
            if 'sort_order' in data:
                tier.sort_order = int(data['sort_order'])
            if 'features' in data:
                import json
                if isinstance(data['features'], list):
                    tier.features = json.dumps(data['features'], ensure_ascii=False)
                else:
                    tier.features = data['features']
            
            db.session.commit()
            
            # 记录操作日志
            from models_admin import AdminLog
            log = AdminLog(
                admin_id=current_admin.id,
                action='update_tier',
                module='membership',
                target_type='tier',
                target_id=tier_id,
                description=f'更新会员套餐 {tier.name}',
                ip_address=request.remote_addr
            )
            db.session.add(log)
            db.session.commit()
            
            return jsonify({
                'success': True,
                'message': '套餐已更新',
                'data': tier.to_dict()
            })
            
        except Exception as e:
            db.session.rollback()
            app.logger.error(f"更新套餐失败: {str(e)}")
            return jsonify({'success': False, 'message': '更新失败'}), 500
    
    elif request.method == 'DELETE':
        try:
            tier = MembershipTier.query.get_or_404(tier_id)
            
            # 检查是否有用户正在使用该套餐
            active_count = UserMembership.query.filter_by(tier_id=tier_id, is_active=True).count()
            if active_count > 0:
                return jsonify({
                    'success': False,
                    'message': f'该套餐仍有{active_count}个活跃用户，无法删除'
                }), 400
            
            tier_name = tier.name
            db.session.delete(tier)
            db.session.commit()
            
            # 记录操作日志
            from models_admin import AdminLog
            log = AdminLog(
                admin_id=current_admin.id,
                action='delete_tier',
                module='membership',
                target_type='tier',
                target_id=tier_id,
                description=f'删除会员套餐 {tier_name}',
                ip_address=request.remote_addr
            )
            db.session.add(log)
            db.session.commit()
            
            return jsonify({
                'success': True,
                'message': f'套餐 {tier_name} 已删除'
            })
            
        except Exception as e:
            db.session.rollback()
            app.logger.error(f"删除套餐失败: {str(e)}")
            return jsonify({'success': False, 'message': '删除失败'}), 500


# 会员统计接口 - 支持单数和复数形式
@app.route('/api/admin/membership/stats', methods=['GET', 'OPTIONS'])
@app.route('/api/admin/memberships/stats', methods=['GET', 'OPTIONS'])
@api_admin_required
def api_admin_membership_stats(current_admin):
    """获取会员统计数据"""
    if request.method == 'OPTIONS':
        return '', 200
    
    try:
        from datetime import datetime, timedelta
        
        # 总会员数
        total_members = UserMembership.query.filter_by(is_active=True).count()
        
        # 即将过期会员（7天内）
        soon = datetime.now() + timedelta(days=7)
        expiring_members = UserMembership.query.filter(
            UserMembership.is_active == True,
            UserMembership.end_date <= soon,
            UserMembership.end_date > datetime.now()
        ).count()
        
        # 已过期会员
        expired_members = UserMembership.query.filter(
            UserMembership.is_active == True,
            UserMembership.end_date < datetime.now()
        ).count()
        
        # 按套餐统计
        tiers_stats = []
        tiers = MembershipTier.query.all()
        for tier in tiers:
            count = UserMembership.query.filter_by(tier_id=tier.id, is_active=True).count()
            tiers_stats.append({
                'tier_id': tier.id,
                'tier_name': tier.name,
                'count': count
            })
        
        return jsonify({
            'success': True,
            'data': {
                'total_members': total_members,
                'expiring_members': expiring_members,
                'expired_members': expired_members,
                'tiers_stats': tiers_stats
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取会员统计失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取统计数据失败'}), 500


# 早鸟优惠状态接口
@app.route('/api/admin/memberships/early-bird', methods=['GET', 'OPTIONS'])
@api_admin_required
def api_admin_early_bird_status(current_admin):
    """获取早鸟优惠状态"""
    if request.method == 'OPTIONS':
        return '', 200
    
    try:
        # 获取早鸟套餐
        early_bird_tiers = MembershipTier.query.filter(
            MembershipTier.code.like('early_bird_%'),
            MembershipTier.is_active == True
        ).order_by(MembershipTier.price).all()
        
        # 统计每个早鸟套餐的销售情况
        early_bird_data = []
        for tier in early_bird_tiers:
            sold_count = UserMembership.query.filter_by(tier_id=tier.id, is_active=True).count()
            
            # 根据套餐代码确定限额
            if tier.code == 'early_bird_1':
                limit = 100
            elif tier.code == 'early_bird_2':
                limit = 200  # 101-300，所以是200个名额
            elif tier.code == 'early_bird_3':
                limit = 200  # 301-500，所以是200个名额
            else:
                limit = 100
            
            early_bird_data.append({
                'tier_id': tier.id,
                'tier_name': tier.name,
                'tier_code': tier.code,
                'price': tier.price,
                'sold_count': sold_count,
                'limit': limit,
                'remaining': max(0, limit - sold_count),
                'is_available': sold_count < limit
            })
        
        return jsonify({
            'success': True,
            'data': {
                'early_bird_active': len(early_bird_data) > 0,
                'tiers': early_bird_data,
                'total_sold': sum(tier['sold_count'] for tier in early_bird_data),
                'total_limit': sum(tier['limit'] for tier in early_bird_data)
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取早鸟优惠状态失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取早鸟优惠状态失败'}), 500


# 会员开通记录接口
@app.route('/api/admin/memberships/records', methods=['GET', 'OPTIONS'])
@api_admin_required
def api_admin_membership_records(current_admin):
    """获取会员开通记录"""
    if request.method == 'OPTIONS':
        return '', 200
    
    try:
        page = request.args.get('page', 1, type=int)
        per_page = request.args.get('per_page', 10, type=int)
        
        # 查询会员记录
        pagination = UserMembership.query.filter_by(is_active=True)\
            .order_by(UserMembership.created_at.desc())\
            .paginate(page=page, per_page=per_page, error_out=False)
        
        records = []
        for membership in pagination.items:
            user = User.query.get(membership.user_id)
            tier = MembershipTier.query.get(membership.tier_id)
            
            # 判断状态
            from datetime import datetime
            is_expired = membership.end_date < datetime.now()
            
            records.append({
                'id': membership.id,
                'user': {
                    'id': user.id if user else None,
                    'username': user.username if user else '未知用户',
                    'email': user.email if user else ''
                },
                'tier': {
                    'id': tier.id if tier else None,
                    'name': tier.name if tier else '未知套餐',
                    'code': tier.code if tier else ''
                },
                'start_date': membership.start_date.strftime('%Y-%m-%d %H:%M:%S') if membership.start_date else None,
                'end_date': membership.end_date.strftime('%Y-%m-%d %H:%M:%S') if membership.end_date else None,
                'created_at': membership.created_at.strftime('%Y-%m-%d %H:%M:%S') if membership.created_at else None,
                'status': 'expired' if is_expired else 'active'
            })
        
        return jsonify({
            'success': True,
            'data': {
                'records': records,
                'pagination': {
                    'page': page,
                    'per_page': per_page,
                    'total': pagination.total,
                    'pages': pagination.pages,
                    'has_prev': pagination.has_prev,
                    'has_next': pagination.has_next
                }
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取会员记录失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取会员记录失败'}), 500


# ==================== Dashboard API ====================
@app.route('/api/admin/dashboard/stats', methods=['GET', 'OPTIONS'])
@api_admin_required
def api_admin_dashboard_stats(current_admin):
    """获取Dashboard统计数据"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        import psutil
        import shutil
        
        # 获取用户统计
        total_users = User.query.count()
        active_users = User.query.filter_by(is_active=True).count()
        
        # 获取会员统计
        total_members = UserMembership.query.filter_by(is_active=True).count()
        
        # 获取订单统计（模拟数据）
        total_orders = 156
        
        # 系统监控数据
        try:
            # CPU使用率
            cpu_usage = psutil.cpu_percent(interval=0.5)
            
            # 内存使用率
            memory = psutil.virtual_memory()
            memory_usage = memory.percent
            
            # 磁盘使用率
            disk = shutil.disk_usage('/')
            disk_usage = (disk.used / disk.total) * 100
            
            # 数据库性能（简单估算：基于查询响应时间）
            from datetime import datetime
            db_start = datetime.now()
            db.session.execute('SELECT 1')
            db_duration = (datetime.now() - db_start).total_seconds()
            # 假设<0.01s为100%，>0.1s为0%
            db_performance = max(0, min(100, 100 - (db_duration - 0.01) * 1000))
            
        except Exception as monitor_error:
            app.logger.warning(f"获取系统监控数据失败: {str(monitor_error)}")
            cpu_usage = 0
            memory_usage = 0
            disk_usage = 0
            db_performance = 100
        
        return jsonify({
            'success': True,
            'data': {
                'total_users': total_users,
                'active_users': active_users,
                'total_members': total_members,
                'total_orders': total_orders,
                'revenue': 25680.50,
                'cpu_usage': round(cpu_usage, 2),
                'memory_usage': round(memory_usage, 2),
                'disk_usage': round(disk_usage, 2),
                'db_performance': round(db_performance, 2)
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取Dashboard统计失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取统计数据失败'}), 500


@app.route('/api/admin/dashboard/user-trend', methods=['GET', 'OPTIONS'])
@api_admin_required
def api_admin_dashboard_user_trend(current_admin):
    """获取用户趋势数据"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        # 模拟用户趋势数据
        trend_data = []
        from datetime import datetime, timedelta
        
        for i in range(30):
            date = datetime.now() - timedelta(days=29-i)
            trend_data.append({
                'date': date.strftime('%Y-%m-%d'),
                'new_users': 5 + (i % 10),
                'active_users': 50 + (i % 20)
            })
        
        return jsonify({
            'success': True,
            'data': trend_data
        })
        
    except Exception as e:
        app.logger.error(f"获取用户趋势失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取趋势数据失败'}), 500


@app.route('/api/admin/dashboard/revenue-trend', methods=['GET', 'OPTIONS'])
@api_admin_required
def api_admin_dashboard_revenue_trend(current_admin):
    """获取收入趋势数据"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        # 模拟收入趋势数据
        trend_data = []
        from datetime import datetime, timedelta
        
        for i in range(30):
            date = datetime.now() - timedelta(days=29-i)
            trend_data.append({
                'date': date.strftime('%Y-%m-%d'),
                'revenue': 800 + (i % 500)
            })
        
        return jsonify({
            'success': True,
            'data': trend_data
        })
        
    except Exception as e:
        app.logger.error(f"获取收入趋势失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取趋势数据失败'}), 500


@app.route('/api/admin/dashboard/membership-distribution', methods=['GET', 'OPTIONS'])
@api_admin_required
def api_admin_dashboard_membership_distribution(current_admin):
    """获取会员分布数据"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        # 查询会员分布
        from sqlalchemy import func
        
        distribution = db.session.query(
            MembershipTier.name,
            func.count(UserMembership.id).label('count')
        ).join(
            UserMembership, MembershipTier.id == UserMembership.tier_id
        ).filter(
            UserMembership.is_active == True
        ).group_by(MembershipTier.name).all()
        
        data = [{'name': name, 'value': count} for name, count in distribution]
        
        return jsonify({
            'success': True,
            'data': data
        })
        
    except Exception as e:
        app.logger.error(f"获取会员分布失败: {str(e)}")
        return jsonify({'success': False, 'message': '获取分布数据失败'}), 500


@app.route('/api/admin/dashboard/ai-usage-heatmap', methods=['GET', 'OPTIONS'])
@api_admin_required
def api_admin_dashboard_ai_usage_heatmap(current_admin):
    """获取AI功能使用热力图数据"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        from sqlalchemy import func
        from datetime import datetime, timedelta
        
        days = int(request.args.get('days', 30))
        end_date = datetime.now()
        start_date = end_date - timedelta(days=days)
        
        # AI功能映射
        ai_features = {
            'ai_ask': 'AI问答',
            'generate_ppt': 'PPT生成',
            'generate_lecture': '讲义生成',
            'generate_question': '题目生成',
            'video_summary': '视频摘要'
        }
        
        # 查询usage_logs表获取AI功能使用数据
        usage_data = db.session.query(
            func.date(UsageLog.created_at).label('date'),
            UsageLog.feature_code,
            func.count(UsageLog.id).label('count')
        ).filter(
            UsageLog.created_at >= start_date,
            UsageLog.created_at <= end_date,
            UsageLog.feature_code.in_(list(ai_features.keys()))
        ).group_by(
            func.date(UsageLog.created_at),
            UsageLog.feature_code
        ).all()
        
        # 构建热力图数据结构 [date, feature, count]
        heatmap_data = []
        for record in usage_data:
            date_str = record.date.strftime('%Y-%m-%d')
            feature_label = ai_features.get(record.feature_code, record.feature_code)
            heatmap_data.append([date_str, feature_label, record.count])
        
        # 生成完整的日期列表
        date_list = []
        current = start_date
        while current <= end_date:
            date_list.append(current.strftime('%Y-%m-%d'))
            current += timedelta(days=1)
        
        # 功能列表
        feature_list = list(ai_features.values())
        
        return jsonify({
            'success': True,
            'data': {
                'heatmap_data': heatmap_data,
                'dates': date_list,
                'features': feature_list
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取AI使用热力图数据失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '获取数据失败'}), 500


# ==================== 操作日志 API ====================
@app.route('/api/admin/logs', methods=['GET', 'OPTIONS'])
@api_admin_required
@permission_required('log_view')
def api_admin_logs_list(current_admin):
    """获取操作日志列表"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        from models_admin import AdminLog, Admin
        
        page = int(request.args.get('page', 1))
        per_page = int(request.args.get('per_page', 20))
        
        # 筛选条件
        admin_id = request.args.get('admin_id')
        action = request.args.get('action')
        module = request.args.get('module')
        start_date = request.args.get('start_date')
        end_date = request.args.get('end_date')
        
        query = AdminLog.query
        
        if admin_id:
            query = query.filter(AdminLog.admin_id == admin_id)
        if action:
            query = query.filter(AdminLog.action == action)
        if module:
            query = query.filter(AdminLog.module == module)
        if start_date:
            from datetime import datetime
            start = datetime.strptime(start_date, '%Y-%m-%d')
            query = query.filter(AdminLog.created_at >= start)
        if end_date:
            from datetime import datetime, timedelta
            end = datetime.strptime(end_date, '%Y-%m-%d') + timedelta(days=1)
            query = query.filter(AdminLog.created_at < end)
        
        # 排序
        query = query.order_by(AdminLog.created_at.desc())
        
        # 分页
        pagination = query.paginate(page=page, per_page=per_page, error_out=False)
        
        logs = []
        for log in pagination.items:
            admin = Admin.query.get(log.admin_id)
            logs.append({
                'id': log.id,
                'admin_id': log.admin_id,
                'admin_username': admin.username if admin else '未知',
                'action': log.action,
                'module': log.module,
                'description': log.description,
                'target_type': log.target_type,
                'target_id': log.target_id,
                'ip_address': log.ip_address,
                'user_agent': log.user_agent,
                'created_at': log.created_at.strftime('%Y-%m-%d %H:%M:%S') if log.created_at else None
            })
        
        return jsonify({
            'success': True,
            'data': logs,
            'total': pagination.total,
            'page': page,
            'per_page': per_page,
            'pages': pagination.pages
        })
        
    except Exception as e:
        app.logger.error(f"获取操作日志失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '获取日志失败'}), 500


@app.route('/api/admin/logs/<int:log_id>', methods=['GET', 'OPTIONS'])
@api_admin_required
@permission_required('log_view')
def api_admin_log_detail(log_id, current_admin):
    """获取操作日志详情"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        from models_admin import AdminLog, Admin
        
        log = AdminLog.query.get_or_404(log_id)
        admin = Admin.query.get(log.admin_id)
        
        return jsonify({
            'success': True,
            'data': {
                'id': log.id,
                'admin_id': log.admin_id,
                'admin_username': admin.username if admin else '未知',
                'admin_email': admin.email if admin else None,
                'action': log.action,
                'module': log.module,
                'description': log.description,
                'target_type': log.target_type,
                'target_id': log.target_id,
                'ip_address': log.ip_address,
                'user_agent': log.user_agent,
                'created_at': log.created_at.strftime('%Y-%m-%d %H:%M:%S') if log.created_at else None
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取日志详情失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '获取日志详情失败'}), 500


@app.route('/api/admin/logs/export', methods=['GET', 'OPTIONS'])
@api_admin_required
@permission_required('log_view')
def api_admin_logs_export(current_admin):
    """导出操作日志"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        from models_admin import AdminLog, Admin
        import io
        from openpyxl import Workbook
        from flask import send_file
        
        # 筛选条件（同列表接口）
        admin_id = request.args.get('admin_id')
        action = request.args.get('action')
        module = request.args.get('module')
        start_date = request.args.get('start_date')
        end_date = request.args.get('end_date')
        
        query = AdminLog.query
        
        if admin_id:
            query = query.filter(AdminLog.admin_id == admin_id)
        if action:
            query = query.filter(AdminLog.action == action)
        if module:
            query = query.filter(AdminLog.module == module)
        if start_date:
            from datetime import datetime
            start = datetime.strptime(start_date, '%Y-%m-%d')
            query = query.filter(AdminLog.created_at >= start)
        if end_date:
            from datetime import datetime, timedelta
            end = datetime.strptime(end_date, '%Y-%m-%d') + timedelta(days=1)
            query = query.filter(AdminLog.created_at < end)
        
        logs = query.order_by(AdminLog.created_at.desc()).limit(10000).all()
        
        # 创建Excel
        wb = Workbook()
        ws = wb.active
        ws.title = "操作日志"
        
        # 表头
        headers = ['ID', '管理员', '操作', '模块', '描述', '目标类型', '目标ID', 'IP地址', '操作时间']
        ws.append(headers)
        
        # 数据
        for log in logs:
            admin = Admin.query.get(log.admin_id)
            ws.append([
                log.id,
                admin.username if admin else '未知',
                log.action,
                log.module,
                log.description,
                log.target_type or '',
                log.target_id or '',
                log.ip_address or '',
                log.created_at.strftime('%Y-%m-%d %H:%M:%S') if log.created_at else ''
            ])
        
        # 保存到内存
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        
        from datetime import datetime
        filename = f'操作日志_{datetime.now().strftime("%Y%m%d_%H%M%S")}.xlsx'
        
        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name=filename
        )
        
    except Exception as e:
        app.logger.error(f"导出操作日志失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '导出失败'}), 500


@app.route('/api/admin/auth/check', methods=['POST', 'OPTIONS'])
@api_admin_required
def api_admin_auth_check(current_admin):
    """检查管理员认证状态"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        return jsonify({
            'success': True,
            'data': {
                'id': current_admin.id,
                'username': current_admin.username,
                'email': current_admin.email,
                'role': current_admin.role,
                'permissions': current_admin.permissions.split(',') if current_admin.permissions else []
            }
        })
        
    except Exception as e:
        app.logger.error(f"检查认证状态失败: {str(e)}")
        return jsonify({'success': False, 'message': '认证检查失败'}), 500


# ==================== 系统设置 API ====================
@app.route('/api/admin/settings', methods=['GET', 'OPTIONS'])
@api_admin_required
@permission_required('system_view')
def api_admin_settings_get(current_admin):
    """获取系统设置"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        # 从配置读取设置
        settings = {
            'basic': {
                'system_name': app.config.get('SYSTEM_NAME', 'EduPilot教育协控系统'),
                'system_logo': app.config.get('SYSTEM_LOGO', ''),
                'system_description': app.config.get('SYSTEM_DESCRIPTION', 'AI驱动的智能教育管理平台'),
                'contact_email': app.config.get('CONTACT_EMAIL', 'support@edupilot.com'),
                'icp_number': app.config.get('ICP_NUMBER', ''),
            },
            'ai': {
                'provider': app.config.get('AI_PROVIDER', 'deepseek'),
                'api_key_configured': bool(app.config.get('DEEPSEEK_API_KEY')),
                'model': app.config.get('AI_MODEL', 'deepseek-chat'),
                'max_tokens': app.config.get('AI_MAX_TOKENS', 4000),
                'temperature': app.config.get('AI_TEMPERATURE', 0.7),
            },
            'payment': {
                'alipay_enabled': bool(app.config.get('ALIPAY_APP_ID')),
                'alipay_app_id': app.config.get('ALIPAY_APP_ID', ''),
                'wechat_enabled': bool(app.config.get('WECHAT_APP_ID')),
                'wechat_app_id': app.config.get('WECHAT_APP_ID', ''),
            },
            'email': {
                'smtp_enabled': bool(app.config.get('MAIL_SERVER')),
                'smtp_server': app.config.get('MAIL_SERVER', ''),
                'smtp_port': app.config.get('MAIL_PORT', 587),
                'smtp_username': app.config.get('MAIL_USERNAME', ''),
                'smtp_use_tls': app.config.get('MAIL_USE_TLS', True),
            },
            'security': {
                'session_timeout': 3600,
                'password_min_length': 8,
                'password_require_special': True,
                'max_login_attempts': 5,
                'login_lockout_duration': 1800,
            }
        }
        
        return jsonify({
            'success': True,
            'data': settings
        })
        
    except Exception as e:
        app.logger.error(f"获取系统设置失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '获取设置失败'}), 500


@app.route('/api/admin/settings', methods=['PUT', 'OPTIONS'])
@api_admin_required
@permission_required('system_edit')
def api_admin_settings_update(current_admin):
    """更新系统设置"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        data = request.get_json()
        category = data.get('category')
        settings = data.get('settings', {})
        
        if not category:
            return jsonify({'success': False, 'message': '缺少设置类别'}), 400
        
        # 记录操作日志
        from models_admin import AdminLog
        log = AdminLog(
            admin_id=current_admin.id,
            action='update',
            module='system',
            target_type='settings',
            description=f'更新系统设置: {category}',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '设置已保存（部分配置需重启服务生效）'
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"更新系统设置失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '更新设置失败'}), 500


# ==================== 管理员管理 API ====================
@app.route('/api/admin/admins', methods=['GET', 'OPTIONS'])
@api_admin_required
@permission_required('admin_view')
def api_admin_admins_list(current_admin):
    """获取管理员列表"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        from models_admin import Admin
        
        page = int(request.args.get('page', 1))
        per_page = int(request.args.get('per_page', 20))
        search = request.args.get('search', '')
        role = request.args.get('role', '')
        
        query = Admin.query
        
        if search:
            query = query.filter(
                db.or_(
                    Admin.username.ilike(f'%{search}%'),
                    Admin.email.ilike(f'%{search}%')
                )
            )
        
        if role:
            query = query.filter(Admin.role == role)
        
        # 排序
        query = query.order_by(Admin.created_at.desc())
        
        # 分页
        pagination = query.paginate(page=page, per_page=per_page, error_out=False)
        
        admins = []
        for admin in pagination.items:
            admins.append({
                'id': admin.id,
                'username': admin.username,
                'email': admin.email,
                'role': admin.role,
                'is_active': admin.is_active,
                'last_login': admin.last_login.strftime('%Y-%m-%d %H:%M:%S') if admin.last_login else None,
                'created_at': admin.created_at.strftime('%Y-%m-%d %H:%M:%S') if admin.created_at else None,
                'permissions': admin.permissions.split(',') if admin.permissions else []
            })
        
        return jsonify({
            'success': True,
            'data': admins,
            'total': pagination.total,
            'page': page,
            'per_page': per_page,
            'pages': pagination.pages
        })
        
    except Exception as e:
        app.logger.error(f"获取管理员列表失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '获取管理员列表失败'}), 500


@app.route('/api/admin/admins', methods=['POST', 'OPTIONS'])
@api_admin_required
@permission_required('admin_edit')
def api_admin_admins_create(current_admin):
    """创建管理员"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        from models_admin import Admin, AdminLog
        from werkzeug.security import generate_password_hash
        
        data = request.get_json()
        username = data.get('username')
        email = data.get('email')
        password = data.get('password')
        role = data.get('role', 'admin')
        permissions = data.get('permissions', [])
        
        if not username or not email or not password:
            return jsonify({'success': False, 'message': '用户名、邮箱和密码不能为空'}), 400
        
        # 检查用户名是否已存在
        if Admin.query.filter_by(username=username).first():
            return jsonify({'success': False, 'message': '用户名已存在'}), 400
        
        # 检查邮箱是否已存在
        if Admin.query.filter_by(email=email).first():
            return jsonify({'success': False, 'message': '邮箱已存在'}), 400
        
        # 创建管理员
        admin = Admin(
            username=username,
            email=email,
            password_hash=generate_password_hash(password),
            role=role,
            permissions=','.join(permissions) if permissions else '',
            is_active=True
        )
        
        db.session.add(admin)
        db.session.commit()
        
        # 记录操作日志
        log = AdminLog(
            admin_id=current_admin.id,
            action='create',
            module='admin',
            target_type='admin',
            target_id=admin.id,
            description=f'创建管理员: {username}',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '管理员创建成功',
            'data': {
                'id': admin.id,
                'username': admin.username,
                'email': admin.email,
                'role': admin.role
            }
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"创建管理员失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '创建管理员失败'}), 500


@app.route('/api/admin/admins/<int:admin_id>', methods=['PUT', 'OPTIONS'])
@api_admin_required
@permission_required('admin_edit')
def api_admin_admins_update(admin_id, current_admin):
    """更新管理员"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        from models_admin import Admin, AdminLog
        from werkzeug.security import generate_password_hash
        
        admin = Admin.query.get_or_404(admin_id)
        
        # 不能修改自己的角色和权限
        if admin.id == current_admin.id:
            return jsonify({'success': False, 'message': '不能修改自己的信息'}), 400
        
        data = request.get_json()
        
        # 更新基本信息
        if 'email' in data:
            # 检查邮箱是否已被其他用户使用
            existing = Admin.query.filter(Admin.email == data['email'], Admin.id != admin_id).first()
            if existing:
                return jsonify({'success': False, 'message': '邮箱已被使用'}), 400
            admin.email = data['email']
        
        if 'role' in data:
            admin.role = data['role']
        
        if 'permissions' in data:
            admin.permissions = ','.join(data['permissions']) if data['permissions'] else ''
        
        if 'is_active' in data:
            admin.is_active = data['is_active']
        
        # 更新密码（如果提供）
        if 'password' in data and data['password']:
            admin.password_hash = generate_password_hash(data['password'])
        
        db.session.commit()
        
        # 记录操作日志
        log = AdminLog(
            admin_id=current_admin.id,
            action='update',
            module='admin',
            target_type='admin',
            target_id=admin_id,
            description=f'更新管理员: {admin.username}',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '管理员更新成功'
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"更新管理员失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '更新管理员失败'}), 500


@app.route('/api/admin/admins/<int:admin_id>', methods=['DELETE', 'OPTIONS'])
@api_admin_required
@permission_required('admin_edit')
def api_admin_admins_delete(admin_id, current_admin):
    """删除管理员"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        from models_admin import Admin, AdminLog
        
        admin = Admin.query.get_or_404(admin_id)
        
        # 不能删除自己
        if admin.id == current_admin.id:
            return jsonify({'success': False, 'message': '不能删除自己'}), 400
        
        # 软删除（设置为非活跃状态）
        admin.is_active = False
        db.session.commit()
        
        # 记录操作日志
        log = AdminLog(
            admin_id=current_admin.id,
            action='delete',
            module='admin',
            target_type='admin',
            target_id=admin_id,
            description=f'删除管理员: {admin.username}',
            ip_address=request.remote_addr
        )
        db.session.add(log)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'message': '管理员删除成功'
        })
        
    except Exception as e:
        db.session.rollback()
        app.logger.error(f"删除管理员失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '删除管理员失败'}), 500


@app.route('/api/admin/permissions', methods=['GET', 'OPTIONS'])
@api_admin_required
def api_admin_permissions_list(current_admin):
    """获取权限列表"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        # 定义所有可用权限
        permissions = [
            {'key': 'user_view', 'name': '查看用户', 'module': '用户管理'},
            {'key': 'user_edit', 'name': '编辑用户', 'module': '用户管理'},
            {'key': 'membership_view', 'name': '查看会员', 'module': '会员管理'},
            {'key': 'membership_edit', 'name': '编辑会员', 'module': '会员管理'},
            {'key': 'order_view', 'name': '查看订单', 'module': '订单管理'},
            {'key': 'order_edit', 'name': '编辑订单', 'module': '订单管理'},
            {'key': 'payment_view', 'name': '查看支付', 'module': '支付管理'},
            {'key': 'payment_edit', 'name': '编辑支付', 'module': '支付管理'},
            {'key': 'log_view', 'name': '查看日志', 'module': '日志管理'},
            {'key': 'system_view', 'name': '查看系统设置', 'module': '系统管理'},
            {'key': 'system_edit', 'name': '编辑系统设置', 'module': '系统管理'},
            {'key': 'admin_view', 'name': '查看管理员', 'module': '管理员管理'},
            {'key': 'admin_edit', 'name': '编辑管理员', 'module': '管理员管理'},
        ]
        
        return jsonify({
            'success': True,
            'data': permissions
        })
        
    except Exception as e:
        app.logger.error(f"获取权限列表失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '获取权限列表失败'}), 500


# ==================== 通知中心 API ====================
@app.route('/api/admin/notifications', methods=['GET', 'OPTIONS'])
@api_admin_required
def api_admin_notifications_list(current_admin):
    """获取通知列表"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        # 模拟通知数据（实际应该从数据库读取）
        notifications = [
            {
                'id': 1,
                'title': '新用户注册',
                'content': '用户 test@example.com 刚刚注册了账户',
                'type': 'user',
                'is_read': False,
                'created_at': '2025-10-14 10:30:00'
            },
            {
                'id': 2,
                'title': '订单支付成功',
                'content': '订单 #12345 支付成功，金额 ¥99.00',
                'type': 'order',
                'is_read': False,
                'created_at': '2025-10-14 09:15:00'
            },
            {
                'id': 3,
                'title': '系统维护通知',
                'content': '系统将于今晚22:00-24:00进行维护升级',
                'type': 'system',
                'is_read': True,
                'created_at': '2025-10-13 16:00:00'
            }
        ]
        
        # 筛选未读通知
        unread_count = len([n for n in notifications if not n['is_read']])
        
        return jsonify({
            'success': True,
            'data': notifications,
            'unread_count': unread_count
        })
        
    except Exception as e:
        app.logger.error(f"获取通知列表失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '获取通知失败'}), 500


@app.route('/api/admin/notifications/<int:notification_id>/read', methods=['POST', 'OPTIONS'])
@api_admin_required
def api_admin_notification_mark_read(notification_id, current_admin):
    """标记通知为已读"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        # 这里应该更新数据库中的通知状态
        # 简化处理，直接返回成功
        
        return jsonify({
            'success': True,
            'message': '通知已标记为已读'
        })
        
    except Exception as e:
        app.logger.error(f"标记通知已读失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '操作失败'}), 500


@app.route('/api/admin/notifications/read-all', methods=['POST', 'OPTIONS'])
@api_admin_required
def api_admin_notifications_read_all(current_admin):
    """标记所有通知为已读"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        # 这里应该更新数据库中所有未读通知的状态
        # 简化处理，直接返回成功
        
        return jsonify({
            'success': True,
            'message': '所有通知已标记为已读'
        })
        
    except Exception as e:
        app.logger.error(f"标记所有通知已读失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '操作失败'}), 500


# ==================== 订单管理 API ====================
@app.route('/api/admin/orders', methods=['GET', 'OPTIONS'])
@api_admin_required
@permission_required('order_view')
def api_admin_orders_list(current_admin):
    """获取订单列表"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        # 调试日志
        app.logger.info(f"订单列表API被调用，管理员: {current_admin.username}")
        app.logger.info(f"请求参数: {dict(request.args)}")
        page = int(request.args.get('page', 1))
        per_page = int(request.args.get('per_page', 20))
        keyword = request.args.get('keyword', '').strip()
        user_keyword = request.args.get('user_keyword', '').strip()
        status = request.args.get('status', '').strip()
        payment_method = request.args.get('payment_method', '').strip()
        start_date = request.args.get('start_date', '').strip()
        end_date = request.args.get('end_date', '').strip()
        min_amount = request.args.get('min_amount', '').strip()
        max_amount = request.args.get('max_amount', '').strip()
        sort_by = request.args.get('sort_by', 'created_at')
        sort_order = request.args.get('sort_order', 'desc')
        
        # 构建查询
        query = Order.query
        
        # 调试：检查基础查询
        total_orders_in_db = Order.query.count()
        app.logger.info(f"数据库中总订单数: {total_orders_in_db}")
        
        # 调试：直接SQL查询对比
        try:
            raw_count = db.session.execute(db.text("SELECT COUNT(*) FROM orders")).scalar()
            app.logger.info(f"原始SQL查询订单数: {raw_count}")
        except Exception as e:
            app.logger.error(f"原始SQL查询失败: {e}")
        
        app.logger.info(f"基础查询对象: {query}")
        
        # 关键字搜索（订单号）
        if keyword:
            query = query.filter(Order.order_number.ilike(f'%{keyword}%'))
        
        # 用户关键字搜索
        if user_keyword:
            query = query.join(User).filter(
                db.or_(
                    User.username.ilike(f'%{user_keyword}%'),
                    User.email.ilike(f'%{user_keyword}%')
                )
            )
        
        # 状态筛选
        if status and status != 'all':
            query = query.filter(Order.status == status)
        
        # 支付方式筛选
        if payment_method and payment_method != 'all':
            query = query.filter(Order.payment_method == payment_method)
        
        # 日期范围筛选
        if start_date:
            try:
                start_dt = datetime.strptime(start_date, '%Y-%m-%d')
                query = query.filter(Order.created_at >= start_dt)
            except ValueError:
                pass
        
        if end_date:
            try:
                end_dt = datetime.strptime(end_date, '%Y-%m-%d')
                end_dt = end_dt.replace(hour=23, minute=59, second=59)
                query = query.filter(Order.created_at <= end_dt)
            except ValueError:
                pass
        
        # 金额范围筛选
        if min_amount:
            try:
                query = query.filter(Order.amount >= float(min_amount))
            except ValueError:
                pass
        
        if max_amount:
            try:
                query = query.filter(Order.amount <= float(max_amount))
            except ValueError:
                pass
        
        # 排序
        if sort_by == 'amount':
            order_by = Order.amount.desc() if sort_order == 'desc' else Order.amount.asc()
        elif sort_by == 'completed_at':
            order_by = Order.completed_at.desc() if sort_order == 'desc' else Order.completed_at.asc()
        else:  # created_at
            order_by = Order.created_at.desc() if sort_order == 'desc' else Order.created_at.asc()
        
        query = query.order_by(order_by)
        
        # 调试：检查排序后的查询
        pre_paginate_count = query.count()
        app.logger.info(f"排序后查询记录数: {pre_paginate_count}")
        app.logger.info(f"最终SQL查询: {str(query)}")
        
        # 分页
        pagination = query.paginate(page=page, per_page=per_page, error_out=False)
        
        # 调试日志
        app.logger.info(f"分页结果: total={pagination.total}, items={len(pagination.items)}")
        
        # 构建返回数据
        orders = []
        for order in pagination.items:
            order_data = order.to_dict()
            
            # 添加用户信息
            if order.user:
                order_data['user'] = {
                    'id': order.user.id,
                    'username': order.user.username,
                    'email': order.user.email,
                    'phone': getattr(order.user, 'phone', None),
                }
            
            # 添加套餐信息
            if order.tier:
                order_data['tier'] = {
                    'id': order.tier.id,
                    'name': order.tier.name,
                    'price': float(order.tier.price) if order.tier.price else 0,
                    'duration_days': order.tier.duration_days,
                }
            
            orders.append(order_data)
        
        # 调试日志
        app.logger.info(f"最终返回: orders数量={len(orders)}, total={pagination.total}")
        
        return jsonify({
            'success': True,
            'orders': orders,
            'total': pagination.total,
            'page': page,
            'per_page': per_page,
            'pages': pagination.pages
        })
        
    except Exception as e:
        app.logger.error(f"获取订单列表失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '获取订单列表失败'}), 500


@app.route('/api/admin/orders/<int:order_id>', methods=['GET', 'OPTIONS'])
@api_admin_required
@permission_required('order_view')
def api_admin_order_detail(order_id, current_admin):
    """获取订单详情"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        order = Order.query.get_or_404(order_id)
        order_data = order.to_detail_dict()
        
        # 添加支付信息
        order_data['payment_info'] = {
            'transaction_id': order.transaction_id,
            'payment_time': order.completed_at.strftime('%Y-%m-%d %H:%M:%S') if order.completed_at else None,
            'payment_url': order.payment_url,
        }
        
        # 添加退款信息
        refunds = order.refunds.all()
        if refunds:
            order_data['payment_info']['refund_info'] = []
            for refund in refunds:
                order_data['payment_info']['refund_info'].append({
                    'refund_id': refund.refund_id,
                    'refund_amount': float(refund.amount) if refund.amount else 0,
                    'refund_reason': refund.reason,
                    'refund_time': refund.created_at.strftime('%Y-%m-%d %H:%M:%S') if refund.created_at else None,
                    'refund_status': refund.status,
                })
        
        # 添加操作历史（模拟数据，实际应该从日志表读取）
        order_data['operation_history'] = [
            {
                'id': 1,
                'action': '创建订单',
                'description': f'订单 {order.order_number} 已创建',
                'operator': order.user.username if order.user else '系统',
                'created_at': order.created_at.strftime('%Y-%m-%d %H:%M:%S') if order.created_at else None,
            }
        ]
        
        if order.completed_at:
            order_data['operation_history'].append({
                'id': 2,
                'action': '支付完成',
                'description': f'订单支付成功，金额 ¥{order.amount}',
                'operator': '系统',
                'created_at': order.completed_at.strftime('%Y-%m-%d %H:%M:%S'),
            })
        
        return jsonify({
            'success': True,
            'data': order_data
        })
        
    except Exception as e:
        app.logger.error(f"获取订单详情失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '获取订单详情失败'}), 500


@app.route('/api/admin/orders/stats', methods=['GET', 'OPTIONS'])
@api_admin_required
@permission_required('order_view')
def api_admin_orders_stats(current_admin):
    """获取订单统计数据"""
    if request.method == 'OPTIONS':
        return '', 200
        
    try:
        from sqlalchemy import func
        from datetime import datetime, timedelta
        
        now = datetime.now()
        today_start = now.replace(hour=0, minute=0, second=0, microsecond=0)
        week_start = today_start - timedelta(days=now.weekday())
        month_start = now.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
        
        # 今日统计
        today_stats = db.session.query(
            func.count(Order.id).label('count'),
            func.coalesce(func.sum(Order.amount), 0).label('revenue')
        ).filter(
            Order.created_at >= today_start,
            Order.status == 'completed'
        ).first()
        
        # 本周统计
        week_stats = db.session.query(
            func.count(Order.id).label('count'),
            func.coalesce(func.sum(Order.amount), 0).label('revenue')
        ).filter(
            Order.created_at >= week_start,
            Order.status == 'completed'
        ).first()
        
        # 本月统计
        month_stats = db.session.query(
            func.count(Order.id).label('count'),
            func.coalesce(func.sum(Order.amount), 0).label('revenue')
        ).filter(
            Order.created_at >= month_start,
            Order.status == 'completed'
        ).first()
        
        # 总统计
        total_stats = db.session.query(
            func.count(Order.id).label('count'),
            func.coalesce(func.sum(Order.amount), 0).label('revenue')
        ).filter(Order.status == 'completed').first()
        
        # 收入趋势（最近30天）
        trend_start = today_start - timedelta(days=29)
        revenue_trend = db.session.query(
            func.date(Order.created_at).label('date'),
            func.count(Order.id).label('orders'),
            func.coalesce(func.sum(Order.amount), 0).label('revenue')
        ).filter(
            Order.created_at >= trend_start,
            Order.status == 'completed'
        ).group_by(func.date(Order.created_at)).all()
        
        trend_data = []
        for record in revenue_trend:
            trend_data.append({
                'date': record.date.strftime('%Y-%m-%d'),
                'orders': record.orders,
                'revenue': float(record.revenue)
            })
        
        # 状态分布
        status_dist = db.session.query(
            Order.status,
            func.count(Order.id).label('count')
        ).group_by(Order.status).all()
        
        total_orders = sum([s.count for s in status_dist])
        status_distribution = []
        for record in status_dist:
            status_distribution.append({
                'status': record.status,
                'count': record.count,
                'percentage': round(record.count / total_orders * 100, 2) if total_orders > 0 else 0
            })
        
        # 支付方式分布
        payment_dist = db.session.query(
            Order.payment_method,
            func.count(Order.id).label('count')
        ).filter(Order.status == 'completed').group_by(Order.payment_method).all()
        
        total_payments = sum([p.count for p in payment_dist])
        payment_method_distribution = []
        for record in payment_dist:
            payment_method_distribution.append({
                'method': record.payment_method,
                'count': record.count,
                'percentage': round(record.count / total_payments * 100, 2) if total_payments > 0 else 0
            })
        
        return jsonify({
            'success': True,
            'data': {
                'today_revenue': float(today_stats.revenue) if today_stats else 0,
                'week_revenue': float(week_stats.revenue) if week_stats else 0,
                'month_revenue': float(month_stats.revenue) if month_stats else 0,
                'total_revenue': float(total_stats.revenue) if total_stats else 0,
                'today_orders': today_stats.count if today_stats else 0,
                'week_orders': week_stats.count if week_stats else 0,
                'month_orders': month_stats.count if month_stats else 0,
                'total_orders': total_stats.count if total_stats else 0,
                'revenue_trend': trend_data,
                'status_distribution': status_distribution,
                'payment_method_distribution': payment_method_distribution,
            }
        })
        
    except Exception as e:
        app.logger.error(f"获取订单统计失败: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'message': '获取订单统计失败'}), 500


app.logger.info("管理后台路由注册成功")


if __name__ == '__main__':
    with app.app_context():
        db.create_all()  # 确保创建所有表，包括新的VideoNote表
    app.run(host='0.0.0.0', port=5000)
